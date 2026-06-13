"""Inject a native-diagram component into a target deck as native editable shapes.

The component's flattened shapes are appended into a slide's ``spTree`` verbatim.
Two pieces of bookkeeping make them collision-free in their new home:

1. **Media** — each referenced bitmap is re-added through python-pptx's
   ``get_or_add_image_part`` (which registers the content type + relationship),
   and the shape's ``r:embed`` is rewritten to the freshly issued ``rId``.
2. **Shape ids** — every ``<p:cNvPr id="...">`` is renumbered above the target
   slide's current maximum so ids stay unique within the slide.

Optional ``pos=(x, y, w, h)`` (EMU) re-frames the whole diagram by wrapping the
lifted shapes in a group whose ``a:off`` / ``a:ext`` scales the original canvas
bounds to the requested rectangle; omit it to drop the diagram at its original
slide coordinates.
"""
from __future__ import annotations

import copy
import gzip
import itertools
import json
from pathlib import Path

from lxml import etree

from .component import A, P, NS, local


def _max_cnvpr_id(sp_tree) -> int:
    ids = [
        int(c.get("id"))
        for c in sp_tree.iter("{%s}cNvPr" % P)
        if (c.get("id") or "").isdigit()
    ]
    return max(ids) if ids else 1


def _wrap_in_group(shapes: list, group_id: int, canvas, pos) -> etree._Element:
    """Wrap *shapes* in a grpSp that maps the canvas bounds onto *pos* (EMU)."""
    cw, ch = canvas
    x, y, w, h = pos
    grp = etree.SubElement(etree.Element("{%s}_tmp" % P), "{%s}grpSp" % P)
    nv = etree.SubElement(grp, "{%s}nvGrpSpPr" % P)
    etree.SubElement(nv, "{%s}cNvPr" % P, id=str(group_id), name=f"Diagram {group_id}")
    etree.SubElement(nv, "{%s}cNvGrpSpPr" % P)
    etree.SubElement(nv, "{%s}nvPr" % P)
    gpr = etree.SubElement(grp, "{%s}grpSpPr" % P)
    xfrm = etree.SubElement(gpr, "{%s}xfrm" % A)
    etree.SubElement(xfrm, "{%s}off" % A, x=str(int(x)), y=str(int(y)))
    etree.SubElement(xfrm, "{%s}ext" % A, cx=str(int(w)), cy=str(int(h)))
    # Child coordinate space stays the original canvas, so the lifted shapes keep
    # their absolute coords; the off/ext vs chOff/chExt ratio does the scaling.
    etree.SubElement(xfrm, "{%s}chOff" % A, x="0", y="0")
    etree.SubElement(xfrm, "{%s}chExt" % A, cx=str(int(cw)), cy=str(int(ch)))
    for s in shapes:
        grp.append(s)
    return grp


def inject_diagram(
    component_dir: str | Path,
    out_pptx: str | Path,
    *,
    target_pptx: str | Path | None = None,
    slide_index: int | None = None,
    new_slide: bool = True,
    pos: tuple[float, float, float, float] | None = None,
) -> dict:
    from pptx import Presentation

    comp_dir = Path(component_dir)
    meta = json.loads((comp_dir / "meta.json").read_text(encoding="utf-8"))
    gz = comp_dir / "shapes.xml.gz"
    if gz.exists():
        root = etree.fromstring(gzip.decompress(gz.read_bytes()))
    else:  # plain .xml fallback (hand-authored components)
        root = etree.fromstring((comp_dir / "shapes.xml").read_bytes())

    if target_pptx:
        prs = Presentation(str(target_pptx))
    else:
        prs = Presentation()
        prs.slide_width, prs.slide_height = meta["canvas_emu"]

    if new_slide or len(prs.slides) == 0:
        layouts = prs.slide_layouts
        blank = layouts[6] if len(layouts) > 6 else layouts[-1]  # "Blank" in the default template
        slide = prs.slides.add_slide(blank)
    else:
        idx = slide_index if slide_index is not None else 0
        slide = prs.slides[idx]

    sp_tree = slide.shapes._spTree

    # Media: re-add each bitmap, build original-rId -> new-rId map.
    rid_map: dict[str, str] = {}
    for rid, rel_path in meta.get("media", {}).items():
        _, new_rid = slide.part.get_or_add_image_part(str(comp_dir / rel_path))
        rid_map[rid] = new_rid

    shapes = [copy.deepcopy(child) for child in root]

    # Remap embed rIds before id renumbering.
    if rid_map:
        for el in shapes:
            for sub in el.iter():
                for k, v in list(sub.attrib.items()):
                    if (k.endswith("}embed") or k.endswith("}link")) and v in rid_map:
                        sub.set(k, rid_map[v])

    counter = itertools.count(_max_cnvpr_id(sp_tree) + 1)
    payload = shapes
    if pos:
        payload = [_wrap_in_group(shapes, next(counter), meta["canvas_emu"], pos)]

    # Renumber every cNvPr id to stay unique within the slide.
    for el in payload:
        for cnv in el.iter("{%s}cNvPr" % P):
            cnv.set("id", str(next(counter)))
        sp_tree.append(el)

    prs.save(str(out_pptx))
    return {
        "out": str(out_pptx),
        "injected_shapes": len(shapes),
        "media": len(rid_map),
        "repositioned": bool(pos),
    }
