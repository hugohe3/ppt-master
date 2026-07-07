#!/usr/bin/env python3
"""
PPT Master - Beautify Text Fidelity Checker

Verifies that a beautified deck preserves source slide text verbatim, page by
page. Intended for the beautify-pptx workflow where one source slide maps to
one output slide and no source text may be dropped or rewritten.

Usage:
    python3 scripts/verify_beautify_fidelity.py <project_path>
    python3 scripts/verify_beautify_fidelity.py <project_path> --target-pptx exports/out.pptx
    python3 scripts/verify_beautify_fidelity.py <project_path> --target-svg-dir svg_output

Examples:
    python3 scripts/verify_beautify_fidelity.py projects/recruitment_deck
    python3 scripts/verify_beautify_fidelity.py projects/recruitment_deck --write-report

Dependencies:
    python-pptx (only when reading PPTX files)
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from console_encoding import configure_utf8_stdio

configure_utf8_stdio()

SVG_NS = "http://www.w3.org/2000/svg"
_WS_RE = re.compile(r"\s+")
_CONTROL_RE = re.compile(r"[\u0000-\u001f\u007f]")
_MOJIBAKE_HINT_RE = re.compile(r"[�]|(?:[\u00c0-\u00ff]{2,})|(?:[鎴涓鍙闀鏄绛寮]{2,})")


def _norm(text: str) -> str:
    """Normalize text for fidelity matching without changing content order."""
    text = text.replace("\x0b", "\n")
    text = _CONTROL_RE.sub("", text)
    text = _WS_RE.sub("", text)
    return text.strip()


def _display(text: str, limit: int = 120) -> str:
    text = _WS_RE.sub(" ", text.strip())
    if len(text) <= limit:
        return text
    return text[: limit - 1] + "…"


def _text_runs_from_shape_text(text: str) -> list[str]:
    return [line.strip() for line in text.replace("\x0b", "\n").splitlines() if line.strip()]


def _load_source_from_clean_json(project_path: Path) -> list[list[str]] | None:
    path = project_path / "analysis" / "source_text_clean.json"
    if not path.exists():
        return None
    data = json.loads(path.read_text(encoding="utf-8-sig"))
    slides = data.get("slides", [])
    return [
        [
            line
            for block in slide.get("texts", [])
            for line in _text_runs_from_shape_text(str(block))
        ]
        for slide in slides
    ]


def _first_source_pptx(project_path: Path) -> Path | None:
    sources = project_path / "sources"
    if not sources.exists():
        return None
    pptx_files = sorted(sources.glob("*.pptx"), key=lambda p: p.stat().st_mtime, reverse=True)
    return pptx_files[0] if pptx_files else None


def _load_pptx_text(path: Path) -> list[list[str]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to read PPTX files") from exc

    prs = Presentation(str(path))
    slides: list[list[str]] = []
    for slide in prs.slides:
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                lines.extend(_text_runs_from_shape_text(shape.text))
        slides.append(lines)
    return slides


def _load_source_text(project_path: Path, source_pptx: Path | None) -> tuple[list[list[str]], str]:
    clean = _load_source_from_clean_json(project_path)
    if clean is not None:
        return clean, "analysis/source_text_clean.json"
    pptx_path = source_pptx or _first_source_pptx(project_path)
    if pptx_path is None:
        raise FileNotFoundError(
            "No source_text_clean.json or source PPTX found. "
            "Run PPTX intake or pass --source-pptx."
        )
    return _load_pptx_text(pptx_path), str(pptx_path)


def _load_svg_text(svg_dir: Path) -> list[list[str]]:
    svg_files = sorted(svg_dir.glob("*.svg"))
    pages: list[list[str]] = []
    for svg_file in svg_files:
        try:
            root = ET.parse(svg_file).getroot()
        except ET.ParseError as exc:
            pages.append([f"__SVG_PARSE_ERROR__:{exc}"])
            continue
        lines: list[str] = []
        for node in root.iter():
            tag = node.tag.rsplit("}", 1)[-1]
            if tag not in {"text", "tspan"}:
                continue
            pieces = []
            if node.text:
                pieces.append(node.text)
            for child in list(node):
                if child.text:
                    pieces.append(child.text)
                if child.tail:
                    pieces.append(child.tail)
            text = "".join(pieces).strip()
            if text:
                lines.append(text)
        pages.append(lines)
    return pages


def _missing_items(source_lines: list[str], target_lines: list[str]) -> list[str]:
    target_joined = _norm("".join(target_lines))
    missing = []
    for line in source_lines:
        normalized = _norm(line)
        if normalized and normalized not in target_joined:
            missing.append(line)
    return missing


def _duplicate_source_lines(source_lines: list[str]) -> list[str]:
    seen = set()
    duplicates = []
    for line in source_lines:
        n = _norm(line)
        if not n:
            continue
        if n in seen and line not in duplicates:
            duplicates.append(line)
        seen.add(n)
    return duplicates


def _mojibake_hits(lines: list[str]) -> list[str]:
    return [line for line in lines if _MOJIBAKE_HINT_RE.search(line)]


def build_report(
    source_pages: list[list[str]],
    target_pages: list[list[str]],
    *,
    source_label: str,
    target_label: str,
) -> dict[str, Any]:
    page_reports = []
    ok = True

    if len(source_pages) != len(target_pages):
        ok = False

    max_pages = max(len(source_pages), len(target_pages))
    for idx in range(max_pages):
        source_lines = source_pages[idx] if idx < len(source_pages) else []
        target_lines = target_pages[idx] if idx < len(target_pages) else []
        missing = _missing_items(source_lines, target_lines)
        mojibake = _mojibake_hits(target_lines)
        parse_errors = [line for line in target_lines if line.startswith("__SVG_PARSE_ERROR__:")]
        duplicates = _duplicate_source_lines(source_lines)
        page_ok = not missing and not mojibake and not parse_errors and bool(source_lines or target_lines)
        if not page_ok:
            ok = False
        page_reports.append(
            {
                "page": idx + 1,
                "ok": page_ok,
                "source_text_count": len(source_lines),
                "target_text_count": len(target_lines),
                "missing_count": len(missing),
                "missing": [_display(item) for item in missing],
                "target_mojibake_count": len(mojibake),
                "target_mojibake": [_display(item) for item in mojibake[:10]],
                "parse_errors": parse_errors,
                "duplicate_source_text": [_display(item) for item in duplicates],
            }
        )

    return {
        "ok": ok,
        "source": source_label,
        "target": target_label,
        "source_page_count": len(source_pages),
        "target_page_count": len(target_pages),
        "page_count_match": len(source_pages) == len(target_pages),
        "pages": page_reports,
    }


def _print_summary(report: dict[str, Any]) -> None:
    status = "OK" if report["ok"] else "FAIL"
    print(f"[{status}] Beautify text fidelity")
    print(f"  Source: {report['source']} ({report['source_page_count']} pages)")
    print(f"  Target: {report['target']} ({report['target_page_count']} pages)")
    if not report["page_count_match"]:
        print("  [ERROR] page count mismatch")
    for page in report["pages"]:
        if page["ok"]:
            continue
        print(f"  [P{page['page']:02d}] missing={page['missing_count']} mojibake={page['target_mojibake_count']}")
        for item in page["missing"][:8]:
            print(f"    - missing: {item}")
        for item in page["target_mojibake"][:4]:
            print(f"    - mojibake: {item}")
        for item in page["parse_errors"]:
            print(f"    - parse: {item}")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Verify beautify-pptx 1:1 text fidelity.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("project_path", help="PPT Master project directory")
    parser.add_argument("--source-pptx", default=None, help="Source PPTX override")
    parser.add_argument("--target-pptx", default=None, help="Generated PPTX to verify")
    parser.add_argument(
        "--target-svg-dir",
        default="svg_output",
        help="Generated SVG directory relative to project path (default: svg_output)",
    )
    parser.add_argument(
        "--write-report",
        action="store_true",
        help="Write analysis/text_fidelity_report.json",
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    project_path = Path(args.project_path)
    if not project_path.is_dir():
        print(f"[ERROR] Project not found: {project_path}", file=sys.stderr)
        return 2

    source_pptx = Path(args.source_pptx) if args.source_pptx else None
    try:
        source_pages, source_label = _load_source_text(project_path, source_pptx)
        if args.target_pptx:
            target_path = Path(args.target_pptx)
            if not target_path.is_absolute():
                target_path = project_path / target_path
            target_pages = _load_pptx_text(target_path)
            target_label = str(target_path)
        else:
            svg_dir = Path(args.target_svg_dir)
            if not svg_dir.is_absolute():
                svg_dir = project_path / svg_dir
            target_pages = _load_svg_text(svg_dir)
            target_label = str(svg_dir)
    except (FileNotFoundError, RuntimeError, json.JSONDecodeError, OSError) as exc:
        print(f"[ERROR] {exc}", file=sys.stderr)
        return 2

    report = build_report(
        source_pages,
        target_pages,
        source_label=source_label,
        target_label=target_label,
    )
    _print_summary(report)

    if args.write_report:
        analysis_dir = project_path / "analysis"
        analysis_dir.mkdir(parents=True, exist_ok=True)
        report_path = analysis_dir / "text_fidelity_report.json"
        report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"  Report: {report_path}")

    return 0 if report["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
