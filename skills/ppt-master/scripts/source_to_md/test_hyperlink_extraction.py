#!/usr/bin/env python3
"""Regression tests for ppt_to_md.py — hyperlink extraction.

Run from the source_to_md/ directory:
    cd skills/ppt-master/scripts/source_to_md/
    ../../../../.venv/bin/python test_hyperlink_extraction.py

All tests PASS — hyperlink extraction is verified for external URLs,
internal slide jumps, shape-level click actions, empty-text fallback,
consecutive same-URL merging, speaker notes, URL parentheses encoding,
and list items.
"""

import os
import shutil
import sys
import tempfile

# Ensure ppt_to_md.py can be imported regardless of working directory
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_to_md import convert_presentation_to_markdown, normalize_text

# ---------------------------------------------------------------------------
# Regression baseline
# ---------------------------------------------------------------------------
# Captured by running convert_presentation_to_markdown on a 1-slide PPTX with
# title "Hello" and body "World / Foo / Bar" saved as "test_baseline.pptx".
BASELINE = (
    "# test_baseline\n\n"
    "- Source: `test_baseline.pptx`\n"
    "- Total slides: 1\n\n"
    "## Slide 1\n\n"
    "Hello\n\n"
    "- World\n- Foo\n- Bar\n"
)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _run_with_cleanup(build_fn, pptx_name="test.pptx"):
    """Build a PPTX in a temp dir, convert, return markdown, clean up."""
    workdir = tempfile.mkdtemp()
    try:
        pptx_path = os.path.join(workdir, pptx_name)
        prs = Presentation()
        build_fn(prs)
        prs.save(pptx_path)
        return convert_presentation_to_markdown(pptx_path, None)
    finally:
        shutil.rmtree(workdir)


def _add_run_with_internal_jump(paragraph, text, source_slide, target_slide):
    """Add a run whose hyperlink jumps to a slide inside the same deck."""
    run = paragraph.add_run()
    run.text = text
    # Direct XML manipulation to set up an internal slide jump
    # (python-pptx Hyperlink.address is for external URLs only)
    hlink_click = etree.SubElement(run._r.get_or_add_rPr(), qn("a:hlinkClick"))
    hlink_click.set("action", "ppaction://hlinksldjump")
    r_id = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    hlink_click.set(qn("r:id"), r_id)
    return run


def _add_run_with_external_url(paragraph, text, url):
    """Add a run with an external hyperlink."""
    run = paragraph.add_run()
    run.text = text
    run.hyperlink.address = url
    return run


def _ensure_notes_text_frame(slide):
    """Return the first text frame from the notes slide, creating one if needed."""
    notes_slide = slide.notes_slide
    for shape in notes_slide.shapes:
        if shape.has_text_frame:
            return shape.text_frame
    # If no text-bearing shape exists, create one
    tx_box = notes_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(2))
    return tx_box.text_frame


# ---------------------------------------------------------------------------
# Test cases
# ---------------------------------------------------------------------------


def test_regression_no_hyperlinks():
    """Regression: no hyperlinks → output matches the hardcoded baseline exactly."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tx_box.text_frame.text = "Hello"

        tx_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
        tf2 = tx_box2.text_frame
        tf2.text = "World"
        tf2.add_paragraph().text = "Foo"
        tf2.add_paragraph().text = "Bar"

    result = _run_with_cleanup(_build, pptx_name="test_baseline.pptx")
    assert result == BASELINE, (
        f"Regression output mismatch.\n"
        f"Expected:\n{repr(BASELINE)}\n"
        f"Got:\n{repr(result)}"
    )


def test_run_external_url():
    """Run-level external URL → markdown link syntax in output."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.clear()
        _add_run_with_external_url(tf.paragraphs[0], "Click", "https://example.com")

    result = _run_with_cleanup(_build)
    assert "https://example.com" in result, (
        f"Expected 'https://example.com' in output.\nGot:\n{result}"
    )
    assert "[Click](https://example.com)" in result, (
        f"Expected '[Click](https://example.com)' in output.\nGot:\n{result}"
    )


def test_run_internal_jump():
    """Run-level internal slide jump → '#slide-N' anchor in output."""
    def _build(prs):
        # Pre-create all 3 slides (slide 0, 1, 2)
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])

        # Label each so they appear in the output
        for s, label in zip((slide1, slide2, slide3), ("First", "Second", "Third")):
            tx_box = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            tx_box.text_frame.text = label

        # On slide1, add a run that jumps to slide3
        tx_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.clear()
        _add_run_with_internal_jump(tf.paragraphs[0], "Go to slide 3", slide1, slide3)

    result = _run_with_cleanup(_build)
    assert "slide-3" in result, (
        f"Expected '#slide-3' anchor in output.\nGot:\n{result}"
    )
    assert "[Go to slide 3](#slide-3)" in result, (
        f"Expected '[Go to slide 3](#slide-3)' in output.\nGot:\n{result}"
    )


def test_shape_external_url():
    """Shape-level external URL → markdown link wrapping entire shape text."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.text = "This is a shape-level hyperlink"
        tx_box.click_action.hyperlink.address = "https://ext.com"

    result = _run_with_cleanup(_build)
    assert "https://ext.com" in result, (
        f"Expected 'https://ext.com' in output.\nGot:\n{result}"
    )
    assert "[This is a shape-level hyperlink](https://ext.com)" in result, (
        f"Expected shape-level link in output.\nGot:\n{result}"
    )


def test_shape_internal_jump():
    """Shape-level internal jump → '#slide-N' anchor in output."""
    def _build(prs):
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        prs.slides.add_slide(prs.slide_layouts[6])  # slide2
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])

        for s, label in zip((slide1, slide3), ("Start", "Target")):
            tx_box = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            tx_box.text_frame.text = label

        tx_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        tx_box.text_frame.text = "Click whole shape to jump"
        tx_box.click_action.target_slide = slide3

    result = _run_with_cleanup(_build)
    assert "slide-3" in result, (
        f"Expected '#slide-3' in output.\nGot:\n{result}"
    )
    assert "[Click whole shape to jump](#slide-3)" in result, (
        f"Expected shape-level jump link in output.\nGot:\n{result}"
    )


def test_empty_text_hyperlink():
    """Run with empty text but hyperlink → URL used as display text fallback."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = ""
        run.hyperlink.address = "https://example.com"

    result = _run_with_cleanup(_build)
    assert "example.com" in result, (
        f"Expected URL in output even with empty run text.\nGot:\n{result}"
    )


def test_consecutive_same_url_merge():
    """Three consecutive runs with the same URL → single merged markdown link."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        for char in ("A", "B", "C"):
            _add_run_with_external_url(p, char, "https://same-url.com")

    result = _run_with_cleanup(_build)
    assert "same-url.com" in result, (
        f"Expected URL in output for consecutive same-URL runs.\nGot:\n{result}"
    )
    assert "[ABC](https://same-url.com)" in result, (
        f"Expected merged single link '[ABC](https://same-url.com)'.\nGot:\n{result}"
    )


def test_notes_hyperlink():
    """Speaker notes with a hyperlinked run → link preserved in notes section."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tx_box.text_frame.text = "Slide content"

        # Add hyperlink to notes
        notes_tf = _ensure_notes_text_frame(slide)
        notes_tf.clear()
        p = notes_tf.paragraphs[0]
        _add_run_with_external_url(p, "notes with link", "https://notes.example.com")

    result = _run_with_cleanup(_build)
    assert "### Speaker Notes" in result, (
        f"Expected '### Speaker Notes' section in output.\nGot:\n{result}"
    )
    assert "https://notes.example.com" in result, (
        f"Expected URL in speaker notes output.\nGot:\n{result}"
    )


def test_url_parentheses_encoding():
    """URL containing '(' and ')' → percent-encoded %28/%29 in output."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.clear()
        _add_run_with_external_url(
            tf.paragraphs[0],
            "Wikipedia",
            "https://en.wikipedia.org/wiki/PowerPoint_(software)",
        )

    result = _run_with_cleanup(_build)
    assert "wikipedia.org" in result, (
        f"Expected Wikipedia URL in output.\nGot:\n{result}"
    )
    # URL parentheses must be percent-encoded
    assert "%28" in result, (
        f"Expected percent-encoded '(' (%28) in URL.\nGot:\n{result}"
    )
    assert "%29" in result, (
        f"Expected percent-encoded ')' (%29) in URL.\nGot:\n{result}"
    )
    assert "(software)" not in result, (
        f"Raw parentheses should NOT appear after encoding.\nGot:\n{result}"
    )


def test_list_item_hyperlink():
    """List item (level > 0) with hyperlinked run → '- [text](url)' format."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
        tf = tx_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.level = 1
        _add_run_with_external_url(p, "list item", "https://list.example.com")

    result = _run_with_cleanup(_build)
    assert "https://list.example.com" in result, (
        f"Expected URL in list item output.\nGot:\n{result}"
    )
    assert "  - [list item]" in result or "- [list item]" in result, (
        f"Expected list item with link in output.\nGot:\n{result}"
    )


def test_mixed_run_hyperlinks():
    """Two different external URLs in consecutive runs → two separate links."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        _add_run_with_external_url(p, "first", "https://first.example.com")
        _add_run_with_external_url(p, "second", "https://second.example.com")

    result = _run_with_cleanup(_build)
    assert "[first](https://first.example.com)" in result, (
        f"Expected first link.\nGot:\n{result}"
    )
    assert "[second](https://second.example.com)" in result, (
        f"Expected second link.\nGot:\n{result}"
    )


def test_interleaved_text_and_link():
    """Plain text between two hyperlinked runs → text between two links."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        _add_run_with_external_url(p, "before", "https://a.example.com")
        p.add_run().text = " middle "
        _add_run_with_external_url(p, "after", "https://b.example.com")

    result = _run_with_cleanup(_build)
    assert "[before](https://a.example.com)" in result
    assert " middle " not in result, (
        f"Leading/trailing spaces should be stripped by normalize_text.\nGot:\n{result}"
    )
    assert "middle" in result, (
        f"Expected 'middle' text between links.\nGot:\n{result}"
    )
    assert "[after](https://b.example.com)" in result


def test_run_mailto_url():
    """Run with mailto: hyperlink → preserved in markdown output."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.clear()
        _add_run_with_external_url(tf.paragraphs[0], "email us", "mailto:foo@example.com")

    result = _run_with_cleanup(_build)
    assert "mailto:foo@example.com" in result, (
        f"Expected mailto: URL in output.\nGot:\n{result}"
    )
    assert "[email us](mailto:foo@example.com)" in result, (
        f"Expected markdown link '[email us](mailto:...)' .\nGot:\n{result}"
    )


def test_run_tel_url():
    """Run with tel: hyperlink → preserved in markdown output."""
    def _build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tx_box.text_frame
        tf.clear()
        _add_run_with_external_url(tf.paragraphs[0], "call us", "tel:+15550123")

    result = _run_with_cleanup(_build)
    assert "tel:+15550123" in result, (
        f"Expected tel: URL in output.\nGot:\n{result}"
    )
    assert "[call us](tel:+15550123)" in result, (
        f"Expected markdown link '[call us](tel:...)' .\nGot:\n{result}"
    )

# ---------------------------------------------------------------------------
# Test registry & runner
# ---------------------------------------------------------------------------

# (test_func, name)
TEST_CASES = [
    (test_regression_no_hyperlinks, "test_regression_no_hyperlinks"),
    (test_run_external_url, "test_run_external_url"),
    (test_run_internal_jump, "test_run_internal_jump"),
    (test_shape_external_url, "test_shape_external_url"),
    (test_shape_internal_jump, "test_shape_internal_jump"),
    (test_empty_text_hyperlink, "test_empty_text_hyperlink"),
    (test_consecutive_same_url_merge, "test_consecutive_same_url_merge"),
    (test_notes_hyperlink, "test_notes_hyperlink"),
    (test_url_parentheses_encoding, "test_url_parentheses_encoding"),
    (test_list_item_hyperlink, "test_list_item_hyperlink"),
    (test_mixed_run_hyperlinks, "test_mixed_run_hyperlinks"),
    (test_interleaved_text_and_link, "test_interleaved_text_and_link"),
    (test_run_mailto_url, "test_run_mailto_url"),
    (test_run_tel_url, "test_run_tel_url"),
    ]


def main():
    pass_count = 0
    fail_count = 0

    for test_func, name in TEST_CASES:
        try:
            test_func()
            print(f"PASS: {name}")
            pass_count += 1
        except AssertionError as e:
            print(f"FAIL: {name} - {e}")
            fail_count += 1
        except Exception as e:
            print(f"ERROR: {name} - {type(e).__name__}: {e}")
            fail_count += 1

    print()
    print(f"Summary: {pass_count} pass / {fail_count} fail")
    if fail_count:
        print(f"ERROR: {fail_count} test(s) failed!")

    return 0 if fail_count == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
