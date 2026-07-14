#!/usr/bin/env python3
"""Regression tests for SVG checker compatibility severity."""

import base64
import copy
import contextlib
import io
import json
import math
import re
import sys
import tempfile
import unittest
import warnings
import zipfile
from pathlib import Path
from unittest.mock import patch
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from animation_config import main as animation_config_main
from svg_quality_checker import SVGQualityChecker
from pptx_shapes import (
    CONNECTOR_PRESET_TYPES,
    OOXML_COORDINATE_MAX,
    get_preset_registry,
    svg_preset_preview_fingerprint,
    svg_text_fingerprint,
)
from pptx_to_svg.preset_authoring import (
    materialize_compact_authored_preset_tree,
    render_preset_shape_fragment,
    validate_authored_preset_tree,
)
from pptx_to_svg.emu_units import Xfrm
from pptx_to_svg.effect_to_svg import (
    convert_effects,
    unsupported_target_effect_metadata,
)
from pptx_to_svg.preset_registry_to_svg import render_preset_geometry
from pptx_to_svg.preset_svg_markup import serialize_preset_layers
from pptx_to_svg.converter import ConvertOptions, convert_pptx_to_svg
from pptx_to_svg.ooxml_loader import parse_ooxml_boolean
from pptx_to_svg.txbody_to_svg import convert_txbody, convert_vertical_txbody
from svg_to_pptx.animation_config import (
    build_group_listing,
    build_scaffold,
    validate_animation_config,
)
from svg_to_pptx.canvas_contract import (
    CanvasContractError,
    parse_project_viewbox,
)
from svg_to_pptx.drawingml.converter import (
    SvgNativeConversionError,
    convert_svg_to_slide_shapes,
)
from svg_to_pptx.pptx_package.builder import create_pptx_with_native_svg
from svg_to_pptx.pptx_package.dimensions import (
    CANVAS_FORMATS as PACKAGE_CANVAS_FORMATS,
    resolve_svg_canvas,
)
from svg_to_pptx.drawingml.elements import parse_project_nested_svg_crop
from svg_to_pptx.drawingml.utils import (
    EMU_PER_PX,
    estimate_text_width,
    parse_inline_style,
    project_filter_errors,
    split_project_text_clusters,
)
from svg_to_pptx.native_objects import (
    NativeMarkerAttributeError,
    native_fallback_kind,
    native_import_source,
    native_replacement_kind,
    native_replacement_status,
    stamp_native_fallback_baseline,
)
from svg_to_pptx.pptx_package.template_structure import (
    NativeLayoutSpec,
    NativePlaceholderSpec,
    TemplateElementSpec,
    TemplateSlideSpec,
    TemplateStructureError,
    _validate_placeholder_carrier,
    match_native_placeholders,
    parse_template_slide,
)
from svg_to_pptx.use_expander import UseExpansionError, expand_local_use_references
from template_preview_pptx import _canvas_viewbox as preview_canvas_viewbox
from template_import.manifest import (
    _effective_inherited_image_assets,
    build_manifest,
    extract_placeholder_text_style,
)
from template_import.native_structure import build_native_structure
from pptx_effects import (
    project_effect_status_errors,
    txbody_has_run_effects,
    unsupported_effect_metadata,
)


class SVGQualityCheckerCompatibilityTests(unittest.TestCase):
    """Keep supported aliases advisory and unsupported input blocking."""

    _TINY_PNG_URI = (
        'data:image/png;base64,'
        'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk'
        '+A8AAQUBAScY42YAAAAASUVORK5CYII='
    )

    def _check(self, content: str) -> dict:
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'fixture.svg'
            svg_path.write_text(content, encoding='utf-8')
            return SVGQualityChecker().check_file(str(svg_path))

    def _check_with_spec_lock(self, content: str, lock_text: str) -> dict:
        with tempfile.TemporaryDirectory() as tmp_dir:
            project = Path(tmp_dir)
            output_dir = project / 'svg_output'
            output_dir.mkdir()
            (project / 'spec_lock.md').write_text(lock_text, encoding='utf-8')
            svg_path = output_dir / 'fixture.svg'
            svg_path.write_text(content, encoding='utf-8')
            return SVGQualityChecker().check_file(str(svg_path))

    def _check_svg_output(self, content: str) -> dict:
        with tempfile.TemporaryDirectory() as tmp_dir:
            output_dir = Path(tmp_dir) / 'svg_output'
            output_dir.mkdir()
            svg_path = output_dir / 'fixture.svg'
            svg_path.write_text(content, encoding='utf-8')
            return SVGQualityChecker().check_file(str(svg_path))

    @staticmethod
    def _write_show_master_fixture(
        path: Path,
        *,
        slide_values: list[str | None],
        layout_value: str | None = None,
    ) -> None:
        """Write a real PPTX with distinct Master/Layout/Slide contributions."""
        presentation = Presentation()
        layout = presentation.slide_layouts[0]
        for index, _value in enumerate(slide_values, start=1):
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = f'SLIDE LOCAL {index}'
        presentation.save(path)

        pml = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        def inherited_shape(name: str, color: str, shape_id: int, x: int) -> ET.Element:
            return ET.fromstring(f'''<p:sp xmlns:p="{pml}" xmlns:a="{dml}">
  <p:nvSpPr><p:cNvPr id="{shape_id}" name="{name}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="914400"/><a:ext cx="1828800" cy="914400"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
</p:sp>''')

        with zipfile.ZipFile(path, 'r') as source:
            members = [(info, source.read(info.filename)) for info in source.infolist()]
        member_data = {info.filename: data for info, data in members}
        replacements: dict[str, bytes] = {}

        master_part = 'ppt/slideMasters/slideMaster1.xml'
        master_root = ET.fromstring(member_data[master_part])
        master_tree = master_root.find(f'{{{pml}}}cSld/{{{pml}}}spTree')
        assert master_tree is not None
        master_tree.append(inherited_shape('MASTER RED', 'FF0000', 901, 914400))
        master_c_sld = master_root.find(f'{{{pml}}}cSld')
        assert master_c_sld is not None
        old_background = master_c_sld.find(f'{{{pml}}}bg')
        background = ET.fromstring(f'''<p:bg xmlns:p="{pml}" xmlns:a="{dml}">
  <p:bgPr><a:solidFill><a:srgbClr val="00FF00"/></a:solidFill><a:effectLst/></p:bgPr>
</p:bg>''')
        if old_background is None:
            master_c_sld.insert(0, background)
        else:
            position = list(master_c_sld).index(old_background)
            master_c_sld.remove(old_background)
            master_c_sld.insert(position, background)
        replacements[master_part] = ET.tostring(
            master_root, encoding='utf-8', xml_declaration=True
        )

        layout_part = 'ppt/slideLayouts/slideLayout1.xml'
        layout_root = ET.fromstring(member_data[layout_part])
        if layout_value is None:
            layout_root.attrib.pop('showMasterSp', None)
        else:
            layout_root.set('showMasterSp', layout_value)
        layout_tree = layout_root.find(f'{{{pml}}}cSld/{{{pml}}}spTree')
        assert layout_tree is not None
        layout_tree.append(inherited_shape('LAYOUT BLUE', '0000FF', 902, 3657600))
        replacements[layout_part] = ET.tostring(
            layout_root, encoding='utf-8', xml_declaration=True
        )

        for index, value in enumerate(slide_values, start=1):
            slide_part = f'ppt/slides/slide{index}.xml'
            slide_root = ET.fromstring(member_data[slide_part])
            if value is None:
                slide_root.attrib.pop('showMasterSp', None)
            else:
                slide_root.set('showMasterSp', value)
            replacements[slide_part] = ET.tostring(
                slide_root, encoding='utf-8', xml_declaration=True
            )

        rewritten = path.with_name(f'{path.stem}-rewritten.pptx')
        with zipfile.ZipFile(rewritten, 'w') as target:
            for info, data in members:
                target.writestr(info, replacements.get(info.filename, data))
        rewritten.replace(path)

    @staticmethod
    def _write_effect_fixture(path: Path, effects: list[str]) -> None:
        """Write one real slide per supplied shape-level effect container."""
        presentation = Presentation()
        blank = presentation.slide_layouts[6]
        for index, _effect in enumerate(effects, start=1):
            slide = presentation.slides.add_slide(blank)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.CUBE,
                Emu(914400),
                Emu(914400),
                Emu(2743200),
                Emu(1828800),
            )
            shape.name = f'EFFECT TARGET {index}'
        presentation.save(path)

        pml = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        with zipfile.ZipFile(path, 'r') as source:
            members = [(info, source.read(info.filename)) for info in source.infolist()]
        replacements: dict[str, bytes] = {}
        member_data = {info.filename: data for info, data in members}
        for index, effect_xml in enumerate(effects, start=1):
            slide_part = f'ppt/slides/slide{index}.xml'
            slide_root = ET.fromstring(member_data[slide_part])
            target = next(
                shape
                for shape in slide_root.findall(f'.//{{{pml}}}sp')
                if (
                    shape.find(f'{{{pml}}}nvSpPr/{{{pml}}}cNvPr')
                    is not None
                    and shape.find(
                        f'{{{pml}}}nvSpPr/{{{pml}}}cNvPr'
                    ).get('name') == f'EFFECT TARGET {index}'
                )
            )
            sp_pr = target.find(f'{{{pml}}}spPr')
            assert sp_pr is not None
            for child in list(sp_pr):
                if child.tag in {
                    f'{{{dml}}}effectLst',
                    f'{{{dml}}}effectDag',
                }:
                    sp_pr.remove(child)
            sp_pr.append(ET.fromstring(effect_xml))
            replacements[slide_part] = ET.tostring(
                slide_root,
                encoding='utf-8',
                xml_declaration=True,
            )

        rewritten = path.with_name(f'{path.stem}-rewritten.pptx')
        with zipfile.ZipFile(rewritten, 'w') as target:
            for info, data in members:
                target.writestr(info, replacements.get(info.filename, data))
        rewritten.replace(path)

    @staticmethod
    def _write_vertical_run_effect_fixture(
        path: Path,
        effects: list[str | None],
        shape_effects: list[str | None] | None = None,
    ) -> None:
        """Write vertical text slides with optional run and shape effects."""
        if shape_effects is not None and len(shape_effects) != len(effects):
            raise ValueError('shape_effects must match the run-effect count')
        presentation = Presentation()
        blank = presentation.slide_layouts[6]
        for index, _effect in enumerate(effects, start=1):
            slide = presentation.slides.add_slide(blank)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(914400),
                Emu(914400),
                Emu(1371600),
                Emu(2743200),
            )
            shape.name = f'VERTICAL EFFECT TEXT {index}'
            shape.text = '甲乙'
        presentation.save(path)

        pml = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        with zipfile.ZipFile(path, 'r') as source:
            members = [(info, source.read(info.filename)) for info in source.infolist()]
        member_data = {info.filename: data for info, data in members}
        replacements: dict[str, bytes] = {}

        for index, effect_xml in enumerate(effects, start=1):
            slide_part = f'ppt/slides/slide{index}.xml'
            slide_root = ET.fromstring(member_data[slide_part])
            target = next(
                shape
                for shape in slide_root.findall(f'.//{{{pml}}}sp')
                if (
                    shape.find(f'{{{pml}}}nvSpPr/{{{pml}}}cNvPr')
                    is not None
                    and shape.find(
                        f'{{{pml}}}nvSpPr/{{{pml}}}cNvPr'
                    ).get('name') == f'VERTICAL EFFECT TEXT {index}'
                )
            )
            tx_body = target.find(f'{{{pml}}}txBody')
            assert tx_body is not None
            body_pr = tx_body.find(f'{{{dml}}}bodyPr')
            assert body_pr is not None
            body_pr.set('vert', 'eaVert')
            run = tx_body.find(f'.//{{{dml}}}r')
            assert run is not None
            run_properties = run.find(f'{{{dml}}}rPr')
            if run_properties is None:
                run_properties = ET.Element(f'{{{dml}}}rPr')
                run.insert(0, run_properties)
            if effect_xml is not None:
                run_properties.append(ET.fromstring(effect_xml))
            shape_effect_xml = (
                shape_effects[index - 1]
                if shape_effects is not None
                else None
            )
            if shape_effect_xml is not None:
                shape_properties = target.find(f'{{{pml}}}spPr')
                assert shape_properties is not None
                shape_properties.append(ET.fromstring(shape_effect_xml))
            replacements[slide_part] = ET.tostring(
                slide_root,
                encoding='utf-8',
                xml_declaration=True,
            )

        rewritten = path.with_name(f'{path.stem}-rewritten.pptx')
        with zipfile.ZipFile(rewritten, 'w') as target:
            for info, data in members:
                target.writestr(info, replacements.get(info.filename, data))
        rewritten.replace(path)

    @staticmethod
    def _write_relationship_run_effect_fixture(
        path: Path,
        cases: list[tuple[str | None, bool]],
    ) -> None:
        """Write text slides with optional hyperlinks and run effects."""
        presentation = Presentation()
        blank = presentation.slide_layouts[6]
        for index, (_effect, has_relationship) in enumerate(cases, start=1):
            slide = presentation.slides.add_slide(blank)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(914400),
                Emu(914400),
                Emu(2743200),
                Emu(914400),
            )
            shape.name = f'RELATIONSHIP EFFECT TEXT {index}'
            paragraph = shape.text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = f'LINK {index}'
            if has_relationship:
                run.hyperlink.address = 'https://example.com/'
        presentation.save(path)

        pml = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        with zipfile.ZipFile(path, 'r') as source:
            members = [
                (info, source.read(info.filename))
                for info in source.infolist()
            ]
        member_data = {info.filename: data for info, data in members}
        replacements: dict[str, bytes] = {}

        for index, (effect_xml, _has_relationship) in enumerate(
            cases,
            start=1,
        ):
            slide_part = f'ppt/slides/slide{index}.xml'
            slide_root = ET.fromstring(member_data[slide_part])
            target = next(
                shape
                for shape in slide_root.findall(f'.//{{{pml}}}sp')
                if (
                    shape.find(f'{{{pml}}}nvSpPr/{{{pml}}}cNvPr')
                    is not None
                    and shape.find(
                        f'{{{pml}}}nvSpPr/{{{pml}}}cNvPr'
                    ).get('name') == f'RELATIONSHIP EFFECT TEXT {index}'
                )
            )
            run = target.find(f'.//{{{dml}}}r')
            assert run is not None
            run_properties = run.find(f'{{{dml}}}rPr')
            if run_properties is None:
                run_properties = ET.Element(f'{{{dml}}}rPr')
                run.insert(0, run_properties)
            if effect_xml is not None:
                # CT_TextCharacterProperties orders the effect group before
                # hlinkClick / hlinkMouseOver / extLst.
                run_properties.insert(0, ET.fromstring(effect_xml))
            replacements[slide_part] = ET.tostring(
                slide_root,
                encoding='utf-8',
                xml_declaration=True,
            )

        rewritten = path.with_name(f'{path.stem}-rewritten.pptx')
        with zipfile.ZipFile(rewritten, 'w') as target:
            for info, data in members:
                target.writestr(info, replacements.get(info.filename, data))
        rewritten.replace(path)

    @staticmethod
    def _write_inherited_run_effect_fixture(
        path: Path,
        *,
        style_source: str = 'layout',
    ) -> None:
        """Write placeholder text inheriting one Layout or Master run effect."""
        if style_source not in {'layout', 'master'}:
            raise ValueError("style_source must be 'layout' or 'master'")
        presentation = Presentation()
        layout = presentation.slide_layouts[1]
        layout_part = str(layout.part.partname).lstrip('/')
        master_part = str(layout.slide_master.part.partname).lstrip('/')
        slide_parts: list[str] = []
        for index in range(1, 4):
            slide = presentation.slides.add_slide(layout)
            slide_parts.append(str(slide.part.partname).lstrip('/'))
            body = slide.placeholders[1]
            body.text = f'INHERITED EFFECT {index}'
            if index == 3:
                body.text_frame.paragraphs[0].runs[0].hyperlink.address = (
                    'https://example.com/'
                )
        presentation.save(path)

        pml = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        def body_placeholder(
            root: ET.Element,
        ) -> tuple[ET.Element, ET.Element]:
            placeholder_path = (
                f'{{{pml}}}nvSpPr/{{{pml}}}nvPr/{{{pml}}}ph'
            )
            for shape in root.findall(f'.//{{{pml}}}sp'):
                placeholder = shape.find(placeholder_path)
                if placeholder is not None and placeholder.get('idx') == '1':
                    return shape, placeholder
            raise AssertionError('content placeholder idx=1 not found')

        with zipfile.ZipFile(path, 'r') as source:
            members = [
                (info, source.read(info.filename))
                for info in source.infolist()
            ]
        member_data = {info.filename: data for info, data in members}
        replacements: dict[str, bytes] = {}

        layout_root = ET.fromstring(member_data[layout_part])
        layout_placeholder, layout_placeholder_properties = body_placeholder(
            layout_root
        )
        layout_placeholder_properties.attrib.pop('type', None)
        list_style = layout_placeholder.find(
            f'{{{pml}}}txBody/{{{dml}}}lstStyle'
        )
        assert list_style is not None

        if style_source == 'layout':
            style_root = list_style
        else:
            master_root = ET.fromstring(member_data[master_part])
            style_root = master_root.find(
                f'{{{pml}}}txStyles/{{{pml}}}bodyStyle'
            )
            assert style_root is not None

        level_properties = style_root.find(f'{{{dml}}}lvl1pPr')
        if level_properties is None:
            level_properties = ET.SubElement(
                style_root,
                f'{{{dml}}}lvl1pPr',
            )
        default_properties = level_properties.find(f'{{{dml}}}defRPr')
        if default_properties is None:
            default_properties = ET.SubElement(
                level_properties,
                f'{{{dml}}}defRPr',
            )
        default_properties.append(ET.fromstring(f'''<a:effectLst xmlns:a="{dml}">
  <a:glow rad="38100"><a:srgbClr val="2563EB"/></a:glow>
</a:effectLst>'''))
        if style_source == 'master':
            replacements[master_part] = ET.tostring(
                master_root,
                encoding='utf-8',
                xml_declaration=True,
            )
        replacements[layout_part] = ET.tostring(
            layout_root,
            encoding='utf-8',
            xml_declaration=True,
        )

        vertical_root = ET.fromstring(member_data[slide_parts[1]])
        vertical_placeholder, _ = body_placeholder(vertical_root)
        body_properties = vertical_placeholder.find(
            f'{{{pml}}}txBody/{{{dml}}}bodyPr'
        )
        assert body_properties is not None
        body_properties.set('vert', 'eaVert')
        replacements[slide_parts[1]] = ET.tostring(
            vertical_root,
            encoding='utf-8',
            xml_declaration=True,
        )

        rewritten = path.with_name(f'{path.stem}-rewritten.pptx')
        with zipfile.ZipFile(rewritten, 'w') as target:
            for info, data in members:
                target.writestr(info, replacements.get(info.filename, data))
        rewritten.replace(path)

    @staticmethod
    def _write_table_run_effect_fixture(
        path: Path,
        cases: list[tuple[str, bool]],
    ) -> None:
        """Write one real table per run-effect placement and transform case."""
        presentation = Presentation()
        blank = presentation.slide_layouts[6]
        for index, (_placement, _rotated) in enumerate(cases, start=1):
            slide = presentation.slides.add_slide(blank)
            frame = slide.shapes.add_table(
                1,
                1,
                Emu(914400),
                Emu(914400),
                Emu(2743200),
                Emu(914400),
            )
            frame.name = f'TABLE EFFECT TEXT {index}'
            frame.table.cell(0, 0).text = f'CELL {index}'
        presentation.save(path)

        pml = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        effect_list = ET.fromstring(f'''<a:effectLst xmlns:a="{dml}">
  <a:glow rad="38100"><a:srgbClr val="2563EB"/></a:glow>
</a:effectLst>''')
        with zipfile.ZipFile(path, 'r') as source:
            members = [
                (info, source.read(info.filename))
                for info in source.infolist()
            ]
        member_data = {info.filename: data for info, data in members}
        replacements: dict[str, bytes] = {}

        for index, (placement, rotated) in enumerate(cases, start=1):
            slide_part = f'ppt/slides/slide{index}.xml'
            slide_root = ET.fromstring(member_data[slide_part])
            frame = next(
                candidate
                for candidate in slide_root.findall(f'.//{{{pml}}}graphicFrame')
                if (
                    candidate.find(
                        f'{{{pml}}}nvGraphicFramePr/{{{pml}}}cNvPr'
                    ) is not None
                    and candidate.find(
                        f'{{{pml}}}nvGraphicFramePr/{{{pml}}}cNvPr'
                    ).get('name') == f'TABLE EFFECT TEXT {index}'
                )
            )
            tx_body = frame.find(f'.//{{{dml}}}tc/{{{dml}}}txBody')
            assert tx_body is not None
            paragraph = tx_body.find(f'{{{dml}}}p')
            assert paragraph is not None
            run = paragraph.find(f'{{{dml}}}r')
            assert run is not None

            if placement.startswith('rPr:'):
                run_properties = run.find(f'{{{dml}}}rPr')
                if run_properties is None:
                    run_properties = ET.Element(f'{{{dml}}}rPr')
                    run.insert(0, run_properties)
                container = placement.split(':', 1)[1]
                if container == 'effectLst':
                    run_properties.insert(0, copy.deepcopy(effect_list))
                elif container.startswith('empty-'):
                    run_properties.insert(
                        0,
                        ET.Element(f'{{{dml}}}{container.removeprefix("empty-")}'),
                    )
                else:
                    raise AssertionError(
                        f'unknown run effect container: {container}'
                    )
            elif placement == 'defRPr':
                paragraph_properties = paragraph.find(f'{{{dml}}}pPr')
                if paragraph_properties is None:
                    paragraph_properties = ET.Element(f'{{{dml}}}pPr')
                    paragraph.insert(0, paragraph_properties)
                default_properties = ET.SubElement(
                    paragraph_properties,
                    f'{{{dml}}}defRPr',
                )
                default_properties.append(copy.deepcopy(effect_list))
            elif placement == 'endParaRPr':
                end_properties = paragraph.find(f'{{{dml}}}endParaRPr')
                if end_properties is None:
                    end_properties = ET.SubElement(
                        paragraph,
                        f'{{{dml}}}endParaRPr',
                    )
                end_properties.insert(0, copy.deepcopy(effect_list))
            elif placement != 'none':
                raise AssertionError(f'unknown table effect placement: {placement}')

            if rotated:
                xfrm = frame.find(f'{{{pml}}}xfrm')
                assert xfrm is not None
                xfrm.set('rot', '60000')
            replacements[slide_part] = ET.tostring(
                slide_root,
                encoding='utf-8',
                xml_declaration=True,
            )

        rewritten = path.with_name(f'{path.stem}-rewritten.pptx')
        with zipfile.ZipFile(rewritten, 'w') as target:
            for info, data in members:
                target.writestr(info, replacements.get(info.filename, data))
        rewritten.replace(path)

    @classmethod
    def _write_picture_group_effect_fixture(cls, path: Path) -> None:
        """Write one real picture and one group with source effects."""
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        image_bytes = base64.b64decode(cls._TINY_PNG_URI.split(',', 1)[1])
        picture = slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            Emu(914400),
            Emu(914400),
            Emu(914400),
            Emu(914400),
        )
        picture.name = 'EFFECT PICTURE'
        group = slide.shapes.add_group_shape()
        group.name = 'EFFECT GROUP'
        group.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(2743200),
            Emu(914400),
            Emu(1371600),
            Emu(914400),
        )
        presentation.save(path)

        pml = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        with zipfile.ZipFile(path, 'r') as source:
            members = [(info, source.read(info.filename)) for info in source.infolist()]
        slide_part = 'ppt/slides/slide1.xml'
        member_data = {info.filename: data for info, data in members}
        slide_root = ET.fromstring(member_data[slide_part])

        picture_xml = next(
            elem
            for elem in slide_root.findall(f'.//{{{pml}}}pic')
            if elem.find(f'{{{pml}}}nvPicPr/{{{pml}}}cNvPr').get('name')
            == 'EFFECT PICTURE'
        )
        picture_sp_pr = picture_xml.find(f'{{{pml}}}spPr')
        assert picture_sp_pr is not None
        blip_fill = picture_xml.find(f'{{{pml}}}blipFill')
        assert blip_fill is not None
        blip = blip_fill.find(f'{{{dml}}}blip')
        assert blip is not None
        blip_fill.insert(
            list(blip_fill).index(blip) + 1,
            ET.fromstring(
                f'<a:srcRect xmlns:a="{dml}" '
                'l="10000" t="10000" r="10000" b="10000"/>'
            ),
        )
        picture_sp_pr.append(ET.fromstring(f'''<a:effectLst xmlns:a="{dml}">
  <a:outerShdw blurRad="38100" dist="38100" dir="5400000">
    <a:srgbClr val="000000"/>
  </a:outerShdw>
</a:effectLst>'''))

        group_xml = next(
            elem
            for elem in slide_root.findall(f'.//{{{pml}}}grpSp')
            if elem.find(f'{{{pml}}}nvGrpSpPr/{{{pml}}}cNvPr').get('name')
            == 'EFFECT GROUP'
        )
        group_sp_pr = group_xml.find(f'{{{pml}}}grpSpPr')
        assert group_sp_pr is not None
        group_sp_pr.append(ET.fromstring(
            f'<a:effectLst xmlns:a="{dml}"><a:reflection/></a:effectLst>'
        ))

        replacements = {
            slide_part: ET.tostring(
                slide_root,
                encoding='utf-8',
                xml_declaration=True,
            ),
        }
        rewritten = path.with_name(f'{path.stem}-rewritten.pptx')
        with zipfile.ZipFile(rewritten, 'w') as target:
            for info, data in members:
                target.writestr(info, replacements.get(info.filename, data))
        rewritten.replace(path)

    @staticmethod
    def _text_bound_warnings(result: dict) -> list[str]:
        return [
            warning
            for warning in result['warnings']
            if 'estimated horizontal bounds exceed the viewBox' in warning
        ]

    @staticmethod
    def _native_txbody_payload(*, run_effect: bool = False) -> str:
        effect_xml = (
            '<a:effectLst><a:outerShdw blurRad="38100" dist="38100" '
            'dir="5400000"><a:srgbClr val="000000"/>'
            '</a:outerShdw></a:effectLst>'
            if run_effect else ''
        )
        native_txbody = (
            '<p:txBody '
            'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:bodyPr/><a:lstStyle/><a:p><a:r>'
            '<a:rPr lang="en-US" sz="1500" spc="-7500">'
            f'{effect_xml}</a:rPr>'
            '<a:t>AB</a:t></a:r></a:p></p:txBody>'
        )
        return base64.b64encode(native_txbody.encode('utf-8')).decode('ascii')

    @staticmethod
    def _txbody_preview(paragraphs: str) -> str:
        tx_body = ET.fromstring(f'''
<p:txBody xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:bodyPr anchor="t" wrap="none" lIns="0" rIns="0" tIns="0" bIns="0"/>
  <a:lstStyle/>{paragraphs}
</p:txBody>''')
        return convert_txbody(
            tx_body,
            Xfrm(x=20, y=20, w=500, h=300),
            None,
        ).svg

    @staticmethod
    def _visible_text_baselines(svg_fragment: str) -> dict[str, list[float]]:
        root = ET.fromstring(
            '<svg xmlns="http://www.w3.org/2000/svg">'
            f'{svg_fragment}</svg>'
        )
        baselines: dict[str, list[float]] = {}
        namespace = '{http://www.w3.org/2000/svg}'
        for text in root.iter(f'{namespace}text'):
            cursor_y = float(text.get('y') or '0')
            for span in text.findall(f'{namespace}tspan'):
                if span.get('y') is not None:
                    cursor_y = float(span.get('y') or '0')
                if span.get('dy') is not None:
                    cursor_y += float(span.get('dy') or '0')
                value = ''.join(span.itertext())
                if value:
                    baselines.setdefault(value, []).append(cursor_y)
        return baselines

    def _assert_checker_and_exporter_reject(
        self,
        content: str,
        expected_checker_text: str,
        expected_exporter_text: str,
    ) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'fixture.svg'
            svg_path.write_text(content, encoding='utf-8')
            result = SVGQualityChecker().check_file(str(svg_path))
            self.assertFalse(result['passed'])
            self.assertIn(expected_checker_text, '\n'.join(result['errors']))
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                expected_exporter_text,
            ):
                convert_svg_to_slide_shapes(svg_path)

    def test_root_viewbox_uses_one_closed_checker_and_exporter_contract(self):
        canonical = (
            '<svg xmlns="http://www.w3.org/2000/svg" '
            'viewBox="0 0 1280 720">'
            '<rect x="0" y="0" width="1280" height="720" fill="#FFFFFF"/>'
            '</svg>'
        )
        canonical_result = self._check(canonical)
        self.assertFalse(any(
            'non-canonical root viewBox' in warning
            for warning in canonical_result['warnings']
        ))

        compatible = canonical.replace(
            'viewBox="0 0 1280 720"',
            'viewBox="0.0,0,+1.28e3,720.0"',
        )
        compatible_result = self._check(compatible)
        self.assertTrue(compatible_result['passed'])
        self.assertTrue(any(
            'Compatible non-canonical root viewBox' in warning
            for warning in compatible_result['warnings']
        ))

        with tempfile.TemporaryDirectory() as tmp_dir:
            canonical_path = Path(tmp_dir) / 'canonical.svg'
            compatible_path = Path(tmp_dir) / 'compatible.svg'
            canonical_path.write_text(canonical, encoding='utf-8')
            compatible_path.write_text(compatible, encoding='utf-8')
            self.assertEqual(
                convert_svg_to_slide_shapes(canonical_path)[0],
                convert_svg_to_slide_shapes(compatible_path)[0],
            )

        fractional = canonical.replace('1280 720', '1280.25 720.01')
        fractional_result = self._check(fractional)
        self.assertTrue(fractional_result['passed'])
        self.assertTrue(any(
            'fractional dimensions are reserved' in warning
            for warning in fractional_result['warnings']
        ))

        invalid_cases = {
            'missing': (
                canonical.replace(' viewBox="0 0 1280 720"', ''),
                'root viewBox is required',
            ),
            'non-finite': (
                canonical.replace('0 0 1280 720', '0 0 NaN 720'),
                'must contain exactly four SVG numbers',
            ),
            'nonzero origin': (
                canonical.replace('0 0 1280 720', '10 0 1280 720'),
                'origin must be "0 0"',
            ),
            'zero width': (
                canonical.replace('0 0 1280 720', '0 0 0 720'),
                'width and height must be positive',
            ),
            'below PowerPoint range': (
                canonical.replace('0 0 1280 720', '0 0 0.4 0.4'),
                "PowerPoint's supported slide range",
            ),
            'oversized exponent': (
                canonical.replace('0 0 1280 720', '0 0 1e999999 720'),
                "PowerPoint's supported slide range",
            ),
            'non-SVG root': (
                '<g xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 1280 720"/>',
                'root element must be <svg>',
            ),
        }
        for name, (content, message) in invalid_cases.items():
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    content,
                    message,
                    re.escape(message),
                )

    def test_root_viewbox_matches_spec_lock_and_all_project_pages(self):
        content = (
            '<svg xmlns="http://www.w3.org/2000/svg" '
            'viewBox="0 0 1280 720"/>'
        )
        mismatched = self._check_with_spec_lock(
            content,
            '## canvas\n- viewBox: 0 0 1024 768\n',
        )
        self.assertTrue(any(
            "spec_lock canvas requires '0 0 1024 768'" in error
            for error in mismatched['errors']
        ))

        invalid_lock = self._check_with_spec_lock(
            content,
            '## canvas\n- viewBox: 0 0 NaN 720\n',
        )
        self.assertTrue(any(
            'spec_lock canvas viewBox must contain exactly four SVG numbers'
            in error
            for error in invalid_lock['errors']
        ))

        partial_lock = self._check_with_spec_lock(
            content,
            '## colors\n- primary: #000000\n',
        )
        self.assertFalse(any(
            'canvas section' in error or 'viewBox mismatch' in error
            for error in partial_lock['errors']
        ))

        with tempfile.TemporaryDirectory() as tmp_dir:
            project = Path(tmp_dir)
            output_dir = project / 'svg_output'
            output_dir.mkdir()
            (project / 'spec_lock.md').write_text(
                '## canvas\n- format: PPT 16:9\n',
                encoding='utf-8',
            )
            (output_dir / '01.svg').write_text(content, encoding='utf-8')
            checker = SVGQualityChecker()
            with contextlib.redirect_stdout(io.StringIO()):
                missing_project_lock = checker.check_directory(str(project))
            self.assertTrue(any(
                'spec_lock.md canvas section must declare viewBox' in error
                for result in missing_project_lock
                for error in result['errors']
            ))

            (project / 'spec_lock.md').write_text(
                '## colors\n- primary: #000000\n',
                encoding='utf-8',
            )
            checker = SVGQualityChecker()
            with contextlib.redirect_stdout(io.StringIO()):
                missing_canvas_section = checker.check_directory(str(project))
            self.assertTrue(any(
                'spec_lock canvas viewBox is required' in error
                for result in missing_canvas_section
                for error in result['errors']
            ))

        with tempfile.TemporaryDirectory() as tmp_dir:
            output_dir = Path(tmp_dir) / 'svg_output'
            output_dir.mkdir()
            (output_dir / '01.svg').write_text(content, encoding='utf-8')
            (output_dir / '02.svg').write_text(
                content.replace('1280 720', '1280.0 720.0'),
                encoding='utf-8',
            )
            checker = SVGQualityChecker()
            with contextlib.redirect_stdout(io.StringIO()):
                equivalent = checker.check_directory(str(output_dir))
            self.assertFalse(any(
                'viewBox mismatch' in error
                for result in equivalent
                for error in result['errors']
            ))

            (output_dir / '02.svg').write_text(
                content.replace('1280 720', '1024 768'),
                encoding='utf-8',
            )
            checker = SVGQualityChecker()
            with contextlib.redirect_stdout(io.StringIO()):
                mixed = checker.check_directory(str(output_dir))
            self.assertTrue(any(
                "first SVG 01.svg requires '0 0 1280 720'" in error
                for result in mixed
                for error in result['errors']
            ))

        with tempfile.TemporaryDirectory() as tmp_dir:
            template_dir = Path(tmp_dir)
            (template_dir / 'design_spec.md').write_text(
                '---\nkind: layout\n---\n# Test\n',
                encoding='utf-8',
            )
            (template_dir / '01.svg').write_text(content, encoding='utf-8')
            checker = SVGQualityChecker(template_mode=True)
            with contextlib.redirect_stdout(io.StringIO()):
                missing_template_lock = checker.check_directory(str(template_dir))
            self.assertTrue(any(
                'design_spec canvas_viewbox viewBox is required' in error
                for result in missing_template_lock
                for error in result['errors']
            ))

    def test_builder_canvas_guard_covers_public_and_internal_svgs(self):
        self.assertIn('ppt43', PACKAGE_CANVAS_FORMATS)
        svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {dims}"/>'
        with tempfile.TemporaryDirectory() as tmp_dir:
            root = Path(tmp_dir)
            public = root / 'public.svg'
            other = root / 'other.svg'
            public.write_text(svg.format(dims='1280 720'), encoding='utf-8')
            other.write_text(svg.format(dims='1024 768'), encoding='utf-8')

            with self.assertRaisesRegex(CanvasContractError, 'other.svg'):
                create_pptx_with_native_svg(
                    [public, other],
                    root / 'mixed-public.pptx',
                    pptx_structure='flat',
                    verbose=False,
                )
            with self.assertRaisesRegex(CanvasContractError, 'other.svg'):
                create_pptx_with_native_svg(
                    [public],
                    root / 'mixed-internal.pptx',
                    pptx_structure='structured',
                    layout_definition_files=[other],
                    verbose=False,
                )
            with self.assertRaisesRegex(CanvasContractError, 'canvas format'):
                resolve_svg_canvas([other], canvas_format='ppt169')

            too_small = root / 'too-small.svg'
            too_small.write_text(svg.format(dims='100 80'), encoding='utf-8')
            with self.assertRaisesRegex(
                CanvasContractError,
                "PowerPoint's supported slide range",
            ):
                resolve_svg_canvas([too_small])

    def test_template_preview_canvas_lock_only_reads_frontmatter(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            spec_path = Path(tmp_dir) / 'design_spec.md'
            spec_path.write_text(
                '---\nkind: layout\n---\n\n```yaml\n'
                'canvas_viewbox: "0 0 1280 720"\n```\n',
                encoding='utf-8',
            )
            self.assertIsNone(preview_canvas_viewbox(spec_path))

            spec_path.write_text(
                '---\nkind: layout\ncanvas_viewbox: "0 0 1280 720"\n---\n',
                encoding='utf-8',
            )
            self.assertEqual(
                preview_canvas_viewbox(spec_path),
                '0 0 1280 720',
            )

    def test_imported_fractional_root_canvas_round_trips_exact_slide_emu(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            root = Path(tmp_dir)
            source_path = root / 'source.pptx'
            svg_path = root / 'fractional.svg'
            output_path = root / 'fractional.pptx'
            source_width = 12192001
            source_height = 6858001
            presentation = Presentation()
            presentation.slide_width = Emu(source_width)
            presentation.slide_height = Emu(source_height)
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(source_path)

            imported = convert_pptx_to_svg(
                source_path,
                options=ConvertOptions(inheritance_mode='flat'),
            )
            imported_svg = imported.slides[0].svg
            svg_path.write_text(imported_svg, encoding='utf-8')
            imported_root = ET.fromstring(imported_svg)
            imported_viewbox = parse_project_viewbox(
                imported_root.get('viewBox')
            )
            self.assertEqual(
                imported_viewbox.emu_dimensions,
                (source_width, source_height),
            )

            # Keep the historical first four positional arguments covered;
            # new optional contracts must be appended to the public API.
            self.assertTrue(create_pptx_with_native_svg(
                [svg_path],
                output_path,
                None,
                False,
                pptx_structure='flat',
                transition=None,
                enable_notes=False,
            ))
            with zipfile.ZipFile(output_path) as package:
                presentation = ET.fromstring(package.read('ppt/presentation.xml'))
            size = presentation.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}sldSz'
            )
            self.assertIsNotNone(size)
            self.assertEqual(
                (int(size.get('cx')), int(size.get('cy'))),
                (source_width, source_height),
            )

    def test_pptx_import_respects_master_shape_visibility(self):
        for raw, expected in (
            (None, True),
            ('1', True),
            ('true', True),
            ('0', False),
            ('false', False),
        ):
            with self.subTest(boolean=raw):
                self.assertEqual(
                    parse_ooxml_boolean(
                        raw,
                        default=True,
                        context='test showMasterSp',
                    ),
                    expected,
                )
        with self.assertRaisesRegex(RuntimeError, 'invalid boolean value'):
            parse_ooxml_boolean(
                'TRUE',
                default=True,
                context='test showMasterSp',
            )

        with tempfile.TemporaryDirectory() as tmp_dir:
            root = Path(tmp_dir)
            invalid_path = root / 'invalid-layout-visibility.pptx'
            self._write_show_master_fixture(
                invalid_path,
                slide_values=['false'],
                layout_value='TRUE',
            )
            with self.assertRaisesRegex(
                RuntimeError,
                'slideLayout1.xml showMasterSp: invalid boolean value',
            ):
                convert_pptx_to_svg(
                    invalid_path,
                    options=ConvertOptions(inheritance_mode='flat'),
                )

            source_path = root / 'slide-visibility.pptx'
            output_dir = root / 'slide-visibility-output'
            self._write_show_master_fixture(
                source_path,
                slide_values=[None, 'false'],
            )
            converted = convert_pptx_to_svg(
                source_path,
                output_dir=output_dir,
                options=ConvertOptions(inheritance_mode='both'),
            )
            visible, hidden = (artifact.svg for artifact in converted.flat_slides)
            self.assertIn('#FF0000', visible)
            self.assertIn('#0000FF', visible)
            self.assertIn('#00FF00', visible)
            self.assertNotIn('#FF0000', hidden)
            self.assertNotIn('#0000FF', hidden)
            self.assertIn('#00FF00', hidden)
            self.assertIn('SLIDE LOCAL 2', hidden)

            hidden_root = ET.fromstring(hidden)
            hidden_title = next(
                element
                for element in hidden_root.iter()
                if element.tag.rsplit('}', 1)[-1] == 'text'
                and 'SLIDE LOCAL 2' in ''.join(element.itertext())
            )
            self.assertGreater(float(hidden_title.get('x', '0')), 0)
            self.assertGreater(float(hidden_title.get('y', '0')), 0)

            self.assertTrue(any(
                '#FF0000' in artifact.svg for artifact in converted.masters
            ))
            self.assertTrue(any(
                '#0000FF' in artifact.svg for artifact in converted.layouts
            ))
            graph = json.loads(
                (output_dir / 'svg' / 'inheritance.json').read_text(
                    encoding='utf-8'
                )
            )
            self.assertEqual(
                [slide['showInheritedShapes'] for slide in graph['slides']],
                [True, False],
            )
            self.assertTrue(all(
                slide['layout'] and slide['master']
                for slide in graph['slides']
            ))
            source_layout = next(
                layout
                for layout in graph['layouts']
                if layout['partPath'].endswith('slideLayout1.xml')
            )
            self.assertTrue(source_layout['showMasterShapes'])

            manifest = build_manifest(source_path, root / 'manifest-default')
            self.assertEqual(
                [slide['showInheritedShapes'] for slide in manifest['slides']],
                [True, False],
            )
            structure = build_native_structure(source_path, manifest)
            self.assertEqual(
                [slide['showInheritedShapes'] for slide in structure['slides']],
                [True, False],
            )

            layout_hidden_path = root / 'layout-visibility.pptx'
            layout_output_dir = root / 'layout-visibility-output'
            self._write_show_master_fixture(
                layout_hidden_path,
                slide_values=['1', 'true'],
                layout_value='false',
            )
            layout_hidden = convert_pptx_to_svg(
                layout_hidden_path,
                output_dir=layout_output_dir,
                options=ConvertOptions(inheritance_mode='both'),
            )
            for artifact in layout_hidden.flat_slides:
                flat = artifact.svg
                self.assertNotIn('#FF0000', flat)
                self.assertIn('#0000FF', flat)
                self.assertIn('#00FF00', flat)

            layout_graph = json.loads(
                (layout_output_dir / 'svg' / 'inheritance.json').read_text(
                    encoding='utf-8'
                )
            )
            source_layout = next(
                layout
                for layout in layout_graph['layouts']
                if layout['partPath'].endswith('slideLayout1.xml')
            )
            self.assertFalse(source_layout['showMasterShapes'])
            self.assertTrue(all(
                slide['showInheritedShapes']
                for slide in layout_graph['slides']
            ))

            manifest = build_manifest(
                layout_hidden_path,
                root / 'manifest-layout-hidden',
            )
            manifest_layout = next(
                layout
                for layout in manifest['layouts']
                if layout['path'].endswith('slideLayout1.xml')
            )
            self.assertFalse(manifest_layout['showMasterShapes'])
            self.assertTrue(all(
                slide['showInheritedShapes'] for slide in manifest['slides']
            ))
            structure = build_native_structure(layout_hidden_path, manifest)
            structure_layout = next(
                layout
                for layout in structure['layouts']
                if layout['packagePart'].endswith('slideLayout1.xml')
            )
            self.assertFalse(structure_layout['showMasterShapes'])
            self.assertTrue(all(
                slide['showInheritedShapes'] for slide in structure['slides']
            ))

        layout_assets = {
            'backgroundAsset': 'layout-background.png',
            'shapeImageAssets': ['layout-shape.png'],
            'showMasterShapes': False,
        }
        master_assets = {
            # The same media is both the overridden Master background and a
            # visible Master picture; shapeImageAssets must retain that use.
            'backgroundAsset': 'master-shared.png',
            'shapeImageAssets': ['master-shared.png'],
        }
        self.assertEqual(
            _effective_inherited_image_assets(
                show_inherited_shapes=False,
                layout_record=layout_assets,
                master_record=master_assets,
            ),
            set(),
        )
        self.assertEqual(
            _effective_inherited_image_assets(
                show_inherited_shapes=True,
                layout_record=layout_assets,
                master_record=master_assets,
            ),
            {'layout-shape.png'},
        )
        layout_assets['showMasterShapes'] = True
        self.assertEqual(
            _effective_inherited_image_assets(
                show_inherited_shapes=True,
                layout_record=layout_assets,
                master_record=master_assets,
            ),
            {'layout-shape.png', 'master-shared.png'},
        )

    def test_manifest_text_style_uses_explicit_element_precedence(self):
        pml = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        shape = ET.fromstring(f'''<p:sp xmlns:p="{pml}" xmlns:a="{dml}">
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r>
    <a:rPr sz="2400" b="1"/><a:t>Primary</a:t>
  </a:r><a:endParaRPr sz="1200" i="1"/></a:p></p:txBody>
</p:sp>''')
        fallback_shape = ET.fromstring(
            f'''<p:sp xmlns:p="{pml}" xmlns:a="{dml}">
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p>
    <a:endParaRPr sz="1200" i="1"/>
  </a:p></p:txBody>
</p:sp>'''
        )

        with warnings.catch_warnings():
            warnings.simplefilter('error', DeprecationWarning)
            primary = extract_placeholder_text_style(shape)
            fallback = extract_placeholder_text_style(fallback_shape)

        self.assertEqual(primary['fontSizePx'], 32.0)
        self.assertTrue(primary['bold'])
        self.assertNotIn('italic', primary)
        self.assertEqual(fallback['fontSizePx'], 16.0)
        self.assertTrue(fallback['italic'])

    def test_top_level_group_ids_are_unique_animation_anchors(self):
        duplicate = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <g id="card"><rect x="10" y="10" width="20" height="20" fill="#000000"/></g>
  <g id="card"><rect x="40" y="10" width="20" height="20" fill="#000000"/></g>
</svg>'''
        result = self._check(duplicate)
        self.assertTrue(any(
            "Duplicate top-level group id 'card'" in error
            for error in result['errors']
        ))

        nested_anonymous = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <g id="card"><g><rect x="10" y="10" width="20" height="20" fill="#000000"/></g></g>
</svg>'''
        nested_result = self._check(nested_anonymous)
        self.assertFalse(any(
            'Top-level visible <g>' in warning
            for warning in nested_result['warnings']
        ))

        blank = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <g id="   ">
    <rect x="10" y="10" width="20" height="20" fill="#000000"/>
    <rect x="40" y="10" width="20" height="20" fill="#000000"/>
  </g>
</svg>'''
        blank_result = self._check(blank)
        self.assertTrue(any(
            'Top-level visible <g> #1 has no id' in warning
            for warning in blank_result['warnings']
        ))
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'blank-id.svg'
            svg_path.write_text(blank, encoding='utf-8')
            conversion = convert_svg_to_slide_shapes(svg_path)
            self.assertEqual(conversion[3], [])
            self.assertIn('<p:grpSp>', conversion[0])

    def test_svg_output_aggregates_only_ordinary_ungrouped_root_atoms(self):
        content = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"
     data-pptx-page-role="content">
  <rect x="0" y="0" width="100" height="100" fill="#FFFFFF"/>
  <rect id="layout-rule" x="0" y="5" width="100" height="2"
        fill="#CCCCCC" data-pptx-layer="layout"/>
  <rect id="footer-rule" x="0" y="95" width="100" height="2"
        fill="#CCCCCC" data-pptx-role="footer"/>
  <g id="card"><g><rect x="10" y="10" width="30" height="20" fill="#EEEEEE"/></g></g>
  <rect id="loose-card" x="50" y="10" width="30" height="20" fill="#EEEEEE"/>
  <text id="loose-title" x="50" y="45" font-size="20">Loose</text>
</svg>'''
        result = self._check_svg_output(content)
        warnings = [
            warning for warning in result['warnings']
            if 'ungrouped top-level Slide-local element(s)' in warning
        ]
        self.assertEqual(len(warnings), 1)
        self.assertIn('2 ungrouped', warnings[0])
        self.assertIn('<rect id="loose-card">', warnings[0])
        self.assertIn('<text id="loose-title">', warnings[0])

        outlined_canvas = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <rect id="outline" x="0" y="0" width="100" height="100"
        fill="none" stroke="#000000"/>
</svg>'''
        outline_result = self._check_svg_output(outlined_canvas)
        self.assertTrue(any(
            '1 ungrouped top-level Slide-local element(s)' in warning
            for warning in outline_result['warnings']
        ))

        inherited_no_fill = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"
     fill="none">
  <rect id="outline" x="0" y="0" width="100" height="100"/>
</svg>'''
        inherited_result = self._check_svg_output(inherited_no_fill)
        self.assertTrue(any(
            '1 ungrouped top-level Slide-local element(s)' in warning
            for warning in inherited_result['warnings']
        ))

        zero_width_stroke = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <rect x="0" y="0" width="100" height="100" fill="#FFFFFF"
        stroke="#000000" stroke-width="0"/>
</svg>'''
        zero_stroke_result = self._check_svg_output(zero_width_stroke)
        self.assertFalse(any(
            'ungrouped top-level Slide-local element(s)' in warning
            for warning in zero_stroke_result['warnings']
        ))

        ordinary_result = self._check(content)
        self.assertFalse(any(
            'ungrouped top-level Slide-local element(s)' in warning
            for warning in ordinary_result['warnings']
        ))

    def test_animation_tools_reject_duplicate_top_level_group_ids(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            project = Path(tmp_dir)
            output_dir = project / 'svg_output'
            output_dir.mkdir()
            (output_dir / 'fixture.svg').write_text(
                '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <g id="card"><rect x="10" y="10" width="20" height="20" fill="#000000"/></g>
  <g id="card"><rect x="40" y="10" width="20" height="20" fill="#000000"/></g>
</svg>''',
                encoding='utf-8',
            )
            config = {
                'version': 1,
                'slides': {'fixture': {'groups': {'card': {}}}},
            }
            diagnostics = validate_animation_config(project, config)
            self.assertEqual(
                sum(
                    'duplicate top-level group id(s)' in item
                    for item in diagnostics
                ),
                1,
            )
            self.assertFalse(any(
                'references missing group' in item
                for item in diagnostics
            ))
            with self.assertRaisesRegex(ValueError, 'duplicate top-level group id'):
                build_scaffold(project)
            with self.assertRaisesRegex(ValueError, 'duplicate top-level group id'):
                build_group_listing(project)

            (project / 'animations.json').write_text(
                json.dumps(config),
                encoding='utf-8',
            )
            for arguments in (
                ['scaffold', str(project), '--force'],
                ['list-groups', str(project)],
                ['validate', str(project)],
            ):
                with self.subTest(command=arguments[0]):
                    stderr = io.StringIO()
                    with patch('sys.stderr', stderr):
                        self.assertEqual(animation_config_main(arguments), 1)
                    self.assertIn(
                        'duplicate top-level group id',
                        stderr.getvalue(),
                    )

    def test_animation_group_ids_are_scoped_to_each_slide(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            project = Path(tmp_dir)
            output_dir = project / 'svg_output'
            output_dir.mkdir()
            for slide_name in ('01_cover', '02_content'):
                (output_dir / f'{slide_name}.svg').write_text(
                    '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <g id="card"><rect x="10" y="10" width="20" height="20" fill="#000000"/></g>
</svg>''',
                    encoding='utf-8',
                )
            scaffold = build_scaffold(project)
            self.assertEqual(
                set(scaffold['slides']),
                {'01_cover', '02_content'},
            )
            self.assertEqual(
                scaffold['slides']['01_cover']['groups'],
                {'card': {}},
            )
            lines, anonymous = build_group_listing(project)
            self.assertEqual(
                lines,
                ['01_cover: card', '02_content: card'],
            )
            self.assertEqual(anonymous, [])

    def test_animation_scan_ignores_anonymous_nested_groups(self):
        with tempfile.TemporaryDirectory() as tmp_dir:
            project = Path(tmp_dir)
            output_dir = project / 'svg_output'
            output_dir.mkdir()
            (output_dir / 'fixture.svg').write_text(
                '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">
  <g id="card"><g><rect x="10" y="10" width="20" height="20" fill="#000000"/></g></g>
  <g id="   "><rect x="40" y="10" width="20" height="20" fill="#000000"/></g>
</svg>''',
                encoding='utf-8',
            )
            lines, anonymous = build_group_listing(project)
            self.assertEqual(lines, ['fixture: card'])
            self.assertEqual(
                anonymous,
                ['fixture: top-level group #2'],
            )

    def _nested_crop_source(
        self,
        *,
        view_box: str = '0.1 0.2 0.7 0.6',
        outer_extra: str = '',
        inner: str | None = None,
        definitions: str = '',
        ancestor_start: str = '',
        ancestor_end: str = '',
    ) -> str:
        inner_markup = inner or (
            f'<image href="{self._TINY_PNG_URI}" x="0" y="0" '
            'width="1" height="1" preserveAspectRatio="none"/>'
        )
        return (
            '<svg xmlns="http://www.w3.org/2000/svg" '
            'xmlns:xlink="http://www.w3.org/1999/xlink" '
            'viewBox="0 0 640 360">'
            f'{definitions}{ancestor_start}'
            '<svg x="100" y="80" width="200" height="100" '
            f'viewBox="{view_box}" preserveAspectRatio="none"'
            f'{outer_extra}>{inner_markup}</svg>{ancestor_end}</svg>'
        )

    def test_imported_nested_crop_round_trips_to_native_src_rect(self):
        source = self._nested_crop_source(
            inner=(
                f'<image href="{self._TINY_PNG_URI}" x="0" y="0" '
                'width="1" height="1" preserveAspectRatio="none" '
                'opacity="0.8"/>'
            ),
        )
        result = self._check(source)
        self.assertTrue(result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'nested-crop.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertEqual(slide_xml.count('<p:pic>'), 1)
        self.assertIn(
            '<a:srcRect l="10000" t="20000" r="20000" b="20000"/>',
            slide_xml,
        )
        self.assertIn('<a:alphaModFix amt="80000"/>', slide_xml)

    def test_nested_crop_accepts_registered_import_and_structure_metadata(self):
        source = self._nested_crop_source(
            outer_extra=(
                ' id="picture-carrier" data-pptx-frame="100 80 200 100"'
                ' data-pptx-layer="layout" data-pptx-object="picture"'
                ' data-pptx-placeholder-carrier="true"'
                ' data-pptx-shape-id="7" data-pptx-shape-name="Picture 7"'
                ' data-pptx-shape-scope="layout"'
            ),
        )
        root = ET.fromstring(source)
        wrappers = list(root.iter('{http://www.w3.org/2000/svg}svg'))
        crop = parse_project_nested_svg_crop(wrappers[1])
        self.assertEqual((crop.width, crop.height), (200.0, 100.0))

    def test_imported_nested_crop_preserves_signed_src_rect_values(self):
        cases = {
            '-0.1 0.2 1.05 0.7': (
                '<a:srcRect l="-10000" t="20000" r="5000" b="10000"/>'
            ),
            '1.2 0 0.1 1': (
                '<a:srcRect l="120000" t="0" r="-30000" b="0"/>'
            ),
        }
        for view_box, expected_src_rect in cases.items():
            with self.subTest(view_box=view_box):
                source = self._nested_crop_source(view_box=view_box)
                result = self._check(source)
                self.assertTrue(result['passed'])
                with tempfile.TemporaryDirectory() as tmp_dir:
                    svg_path = Path(tmp_dir) / 'signed-crop.svg'
                    svg_path.write_text(source, encoding='utf-8')
                    slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
                self.assertIn(expected_src_rect, slide_xml)

    def test_generic_nested_svg_is_not_silently_dropped(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360">
  <svg x="100" y="80" width="200" height="100" viewBox="0 0 200 100">
    <rect x="0" y="0" width="200" height="100" fill="#FF0000"/>
    <text x="20" y="55" font-size="20">Nested</text>
  </svg>
</svg>'''
        self._assert_checker_and_exporter_reject(
            source,
            'invalid imported crop wrapper',
            'invalid nested SVG crop wrapper',
        )

    def test_nested_crop_cannot_hide_extra_visual_siblings(self):
        valid_inner = (
            f'<image href="{self._TINY_PNG_URI}" x="0" y="0" '
            'width="1" height="1" preserveAspectRatio="none"/>'
        )
        source = self._nested_crop_source(
            inner=valid_inner + '<rect x="0" y="0" width="1" height="1"/>',
        )
        self._assert_checker_and_exporter_reject(
            source,
            'expected exactly one direct SVG-namespace <image> child',
            'expected exactly one direct SVG-namespace <image> child',
        )

    def test_nested_crop_clip_marker_and_reference_are_one_closed_contract(self):
        definitions = (
            '<defs><clipPath id="crop-clip" '
            'clipPathUnits="objectBoundingBox">'
            '<ellipse cx="0.5" cy="0.5" rx="0.5" ry="0.5"/>'
            '</clipPath></defs>'
        )
        source = self._nested_crop_source(
            definitions=definitions,
            outer_extra=(
                ' data-pptx-crop="1" clip-path="url(#crop-clip)"'
            ),
        )
        result = self._check(source)
        self.assertTrue(result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'clipped-crop.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertIn('<a:prstGeom prst="ellipse">', slide_xml)

        invalid_sources = {
            'clip without marker': self._nested_crop_source(
                definitions=definitions,
                outer_extra=' clip-path="url(#crop-clip)"',
            ),
            'marker without clip': self._nested_crop_source(
                definitions=definitions,
                outer_extra=' data-pptx-crop="1"',
            ),
            'missing clip': self._nested_crop_source(
                outer_extra=(
                    ' data-pptx-crop="1" clip-path="url(#missing)"'
                ),
            ),
            'malformed clip': self._nested_crop_source(
                outer_extra=' data-pptx-crop="1" clip-path="crop-clip"',
            ),
        }
        for name, invalid_source in invalid_sources.items():
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    invalid_source,
                    'clip',
                    'clip',
                )

    def test_preset_clip_geometry_must_cover_the_complete_image_frame(self):
        def source_for(shape: str, *, units: str = 'userSpaceOnUse') -> str:
            return (
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 100 100"><defs>'
                f'<clipPath id="clip" clipPathUnits="{units}">{shape}'
                '</clipPath></defs>'
                f'<image href="{self._TINY_PNG_URI}" x="10" y="20" '
                'width="40" height="40" clip-path="url(#clip)"/></svg>'
            )

        valid_source = source_for(
            '<circle cx="30" cy="40" r="20"/>'
        )
        result = self._check(valid_source)
        self.assertTrue(result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'full-circle-clip.svg'
            svg_path.write_text(valid_source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertIn('<a:prstGeom prst="ellipse">', slide_xml)

        invalid_sources = {
            'partial circle': source_for(
                '<circle cx="30" cy="40" r="10"/>'
            ),
            'offset rect': source_for(
                '<rect x="15" y="20" width="35" height="40"/>'
            ),
            'non-uniform rounded rect': source_for(
                '<rect x="10" y="20" width="40" height="40" '
                'rx="4" ry="8"/>'
            ),
        }
        for name, source in invalid_sources.items():
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    source,
                    'must cover the complete frame',
                    'invalid project clip-path',
                )

    def test_clip_definition_and_shape_require_the_svg_namespace(self):
        foreign_definition = (
            '<svg xmlns="http://www.w3.org/2000/svg" xmlns:q="urn:foreign" '
            'viewBox="0 0 100 100"><defs>'
            '<q:clipPath id="clip"><q:circle cx="30" cy="40" r="20"/>'
            '</q:clipPath></defs>'
            f'<image href="{self._TINY_PNG_URI}" x="10" y="20" '
            'width="40" height="40" clip-path="url(#clip)"/></svg>'
        )
        foreign_child = foreign_definition.replace(
            '<q:clipPath id="clip"><q:circle',
            '<clipPath id="clip"><q:circle',
        ).replace('</q:clipPath>', '</clipPath>')
        for name, source in {
            'foreign definition': foreign_definition,
            'foreign child': foreign_child,
        }.items():
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    source,
                    'clip',
                    'invalid project clip-path',
                )

    def test_picture_clip_cannot_depend_on_a_winding_rule(self):
        def source_for(clip_extra: str, shape_extra: str) -> str:
            return (
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 100 100"><defs>'
                f'<clipPath id="clip"{clip_extra}>'
                '<path d="M 10 20 L 50 20 L 50 60 L 10 60 Z"'
                f'{shape_extra}/></clipPath></defs>'
                f'<image href="{self._TINY_PNG_URI}" x="10" y="20" '
                'width="40" height="40" clip-path="url(#clip)"/></svg>'
            )

        invalid_sources = {
            'clipPath attribute': source_for(
                ' clip-rule="evenodd"',
                '',
            ),
            'shape attribute': source_for(
                '',
                ' fill-rule="evenodd"',
            ),
            'inline style': source_for(
                '',
                ' style="clip-rule: evenodd"',
            ),
        }
        for name, source in invalid_sources.items():
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    source,
                    'winding-rule control',
                    'invalid project clip-path',
                )

    def test_nested_crop_requires_svg_namespace_and_visual_ancestry(self):
        foreign_inner = (
            f'<q:image xmlns:q="urn:foreign" href="{self._TINY_PNG_URI}" '
            'x="0" y="0" width="1" height="1" '
            'preserveAspectRatio="none"/>'
        )
        invalid_sources = {
            'foreign image': self._nested_crop_source(inner=foreign_inner),
            'defs ancestor': self._nested_crop_source(
                ancestor_start='<defs>',
                ancestor_end='</defs>',
            ),
            'text ancestor': self._nested_crop_source(
                ancestor_start='<text x="0" y="0">',
                ancestor_end='</text>',
            ),
            'render-only ancestor': self._nested_crop_source(
                ancestor_start='<g data-pptx-part="geometry-detail">',
                ancestor_end='</g>',
            ),
            'render-only self': self._nested_crop_source(
                outer_extra=' data-pptx-part="geometry-detail"',
            ),
        }
        for name, source in invalid_sources.items():
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    source,
                    'invalid imported crop wrapper',
                    'invalid nested SVG crop wrapper',
                )

    def test_plain_and_nested_images_share_strict_source_validation(self):
        truncated_emf = bytearray(44)
        truncated_emf[0:4] = b'\x01\x00\x00\x00'
        truncated_emf[40:44] = b' EMF'
        truncated_emf_uri = (
            'data:image/x-emf;base64,'
            + base64.b64encode(bytes(truncated_emf)).decode('ascii')
        )
        truncated_wmf_uri = (
            'data:image/x-wmf;base64,'
            + base64.b64encode(b'\xd7\xcd\xc6\x9a').decode('ascii')
        )
        fake_emf_eof = bytearray(96)
        fake_emf_eof[0:4] = (1).to_bytes(4, 'little')
        fake_emf_eof[4:8] = (88).to_bytes(4, 'little')
        fake_emf_eof[40:44] = b' EMF'
        fake_emf_eof[44:48] = (0x00010000).to_bytes(4, 'little')
        fake_emf_eof[48:52] = (96).to_bytes(4, 'little')
        fake_emf_eof[52:56] = (2).to_bytes(4, 'little')
        fake_emf_eof[88:92] = (14).to_bytes(4, 'little')
        fake_emf_eof[92:96] = (8).to_bytes(4, 'little')
        fake_emf_eof_uri = (
            'data:image/x-emf;base64,'
            + base64.b64encode(bytes(fake_emf_eof)).decode('ascii')
        )
        duplicate_header_emf = bytearray(196)
        duplicate_header_emf[0:88] = fake_emf_eof[0:88]
        duplicate_header_emf[48:52] = (196).to_bytes(4, 'little')
        duplicate_header_emf[52:56] = (3).to_bytes(4, 'little')
        duplicate_header_emf[88:92] = (1).to_bytes(4, 'little')
        duplicate_header_emf[92:96] = (88).to_bytes(4, 'little')
        duplicate_header_emf[176:180] = (14).to_bytes(4, 'little')
        duplicate_header_emf[180:184] = (20).to_bytes(4, 'little')
        duplicate_header_emf[188:192] = (16).to_bytes(4, 'little')
        duplicate_header_emf[192:196] = (20).to_bytes(4, 'little')
        duplicate_header_emf_uri = (
            'data:image/x-emf;base64,'
            + base64.b64encode(bytes(duplicate_header_emf)).decode('ascii')
        )
        early_eof_emf = bytearray(128)
        early_eof_emf[0:88] = fake_emf_eof[0:88]
        early_eof_emf[48:52] = (128).to_bytes(4, 'little')
        early_eof_emf[52:56] = (3).to_bytes(4, 'little')
        for offset in (88, 108):
            early_eof_emf[offset:offset + 4] = (14).to_bytes(4, 'little')
            early_eof_emf[offset + 4:offset + 8] = (20).to_bytes(
                4,
                'little',
            )
            early_eof_emf[offset + 12:offset + 16] = (16).to_bytes(
                4,
                'little',
            )
            early_eof_emf[offset + 16:offset + 20] = (20).to_bytes(
                4,
                'little',
            )
        early_eof_emf_uri = (
            'data:image/x-emf;base64,'
            + base64.b64encode(bytes(early_eof_emf)).decode('ascii')
        )
        bad_size_last_emf = bytearray(108)
        bad_size_last_emf[0:88] = fake_emf_eof[0:88]
        bad_size_last_emf[48:52] = (108).to_bytes(4, 'little')
        bad_size_last_emf[88:92] = (14).to_bytes(4, 'little')
        bad_size_last_emf[92:96] = (20).to_bytes(4, 'little')
        bad_size_last_emf[100:104] = (16).to_bytes(4, 'little')
        bad_size_last_emf[104:108] = (16).to_bytes(4, 'little')
        bad_size_last_emf_uri = (
            'data:image/x-emf;base64,'
            + base64.b64encode(bytes(bad_size_last_emf)).decode('ascii')
        )
        palette_mismatch_emf = bytearray(112)
        palette_mismatch_emf[0:88] = fake_emf_eof[0:88]
        palette_mismatch_emf[48:52] = (112).to_bytes(4, 'little')
        palette_mismatch_emf[88:92] = (14).to_bytes(4, 'little')
        palette_mismatch_emf[92:96] = (24).to_bytes(4, 'little')
        palette_mismatch_emf[96:100] = (1).to_bytes(4, 'little')
        palette_mismatch_emf[100:104] = (16).to_bytes(4, 'little')
        palette_mismatch_emf[108:112] = (24).to_bytes(4, 'little')
        palette_mismatch_emf_uri = (
            'data:image/x-emf;base64,'
            + base64.b64encode(bytes(palette_mismatch_emf)).decode('ascii')
        )
        valid_wmf = bytes.fromhex(
            '0100 0900 0003 0c000000 0000 03000000 0000 '
            '03000000 0000'
        )
        invalid_version_wmf = bytearray(valid_wmf)
        invalid_version_wmf[4:6] = b'\x00\x00'
        invalid_version_wmf_uri = (
            'data:image/x-wmf;base64,'
            + base64.b64encode(bytes(invalid_version_wmf)).decode('ascii')
        )
        inflated_max_wmf = bytearray(valid_wmf)
        inflated_max_wmf[12:16] = (4).to_bytes(4, 'little')
        inflated_max_wmf_uri = (
            'data:image/x-wmf;base64,'
            + base64.b64encode(bytes(inflated_max_wmf)).decode('ascii')
        )
        oversized_eof_wmf = bytes.fromhex(
            '0100 0900 0003 0d000000 0000 04000000 0000 '
            '04000000 0000 0000'
        )
        oversized_eof_wmf_uri = (
            'data:image/x-wmf;base64,'
            + base64.b64encode(oversized_eof_wmf).decode('ascii')
        )
        early_eof_wmf = bytes.fromhex(
            '0100 0900 0003 0f000000 0000 03000000 0000 '
            '03000000 0000 03000000 0000'
        )
        early_eof_wmf_uri = (
            'data:image/x-wmf;base64,'
            + base64.b64encode(early_eof_wmf).decode('ascii')
        )
        plain_sources = {
            'missing href': (
                '<image x="10" y="10" width="40" height="40"/>'
            ),
            'empty href': (
                '<image href="" x="10" y="10" width="40" height="40"/>'
            ),
            'dual href': (
                f'<image href="{self._TINY_PNG_URI}" '
                f'xlink:href="{self._TINY_PNG_URI}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'invalid base64': (
                '<image href="data:image/png;base64,%%%" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'empty payload': (
                '<image href="data:image/png;base64," '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'wrong mime': (
                '<image href="data:text/plain;base64,QQ==" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'invalid image bytes': (
                '<image href="data:image/png;base64,QQ==" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'truncated emf': (
                f'<image href="{truncated_emf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'truncated wmf': (
                f'<image href="{truncated_wmf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'truncated emf eof record': (
                f'<image href="{fake_emf_eof_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'duplicate emf header': (
                f'<image href="{duplicate_header_emf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'nonfinal emf eof record': (
                f'<image href="{early_eof_emf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'mismatched emf eof size-last': (
                f'<image href="{bad_size_last_emf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'mismatched emf palette count': (
                f'<image href="{palette_mismatch_emf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'invalid wmf version': (
                f'<image href="{invalid_version_wmf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'inflated wmf maximum record': (
                f'<image href="{inflated_max_wmf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'oversized wmf eof record': (
                f'<image href="{oversized_eof_wmf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
            'nonfinal wmf eof record': (
                f'<image href="{early_eof_wmf_uri}" '
                'x="10" y="10" width="40" height="40"/>'
            ),
        }
        for name, image_markup in plain_sources.items():
            with self.subTest(kind='plain', name=name):
                source = (
                    '<svg xmlns="http://www.w3.org/2000/svg" '
                    'xmlns:xlink="http://www.w3.org/1999/xlink" '
                    f'viewBox="0 0 100 100">{image_markup}</svg>'
                )
                self._assert_checker_and_exporter_reject(
                    source,
                    'invalid image source',
                    'invalid project image',
                )

        nested_inner = (
            '<image href="data:image/png;base64,%%%" x="0" y="0" '
            'width="1" height="1" preserveAspectRatio="none"/>'
        )
        self._assert_checker_and_exporter_reject(
            self._nested_crop_source(inner=nested_inner),
            'invalid image source',
            'invalid project image',
        )

    def test_external_image_bytes_must_match_their_extension(self):
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<image href="broken.png" x="10" y="10" '
            'width="40" height="40"/></svg>'
        )
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'fixture.svg'
            svg_path.write_text(source, encoding='utf-8')
            (Path(tmp_dir) / 'broken.png').write_bytes(b'not a png')
            result = SVGQualityChecker().check_file(str(svg_path))
            self.assertFalse(result['passed'])
            self.assertIn('does not match', '\n'.join(result['errors']))
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'invalid project image',
            ):
                convert_svg_to_slide_shapes(svg_path)

    def test_images_require_explicit_positive_dimensions(self):
        image_cases = {
            'missing width': (
                f'<image href="{self._TINY_PNG_URI}" x="10" y="10" '
                'height="40"/>'
            ),
            'missing height': (
                f'<image href="{self._TINY_PNG_URI}" x="10" y="10" '
                'width="40"/>'
            ),
            'zero width': (
                f'<image href="{self._TINY_PNG_URI}" x="10" y="10" '
                'width="0" height="40"/>'
            ),
            'zero height': (
                f'<image href="{self._TINY_PNG_URI}" x="10" y="10" '
                'width="40" height="0"/>'
            ),
            'style overrides width with zero': (
                f'<image href="{self._TINY_PNG_URI}" x="10" y="10" '
                'width="40" height="40" style="width: 0"/>'
            ),
        }
        for name, image_markup in image_cases.items():
            with self.subTest(name=name):
                source = (
                    '<svg xmlns="http://www.w3.org/2000/svg" '
                    f'viewBox="0 0 100 100">{image_markup}</svg>'
                )
                self._assert_checker_and_exporter_reject(
                    source,
                    'positive',
                    'invalid project image',
                )

    def test_compatible_inline_image_dimensions_remain_exportable(self):
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            f'<image href="{self._TINY_PNG_URI}" x="10" y="10" '
            'style="width: 40px; height: 40px"/></svg>'
        )
        result = self._check(source)
        self.assertTrue(result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'inline-image-frame.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertEqual(slide_xml.count('<p:pic>'), 1)

    def test_external_svg_image_is_a_registered_source_format(self):
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<image href="asset.svg" x="10" y="10" '
            'width="40" height="40"/></svg>'
        )
        asset = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 10 10">'
            '<rect width="10" height="10" fill="#2563EB"/></svg>'
        )
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'fixture.svg'
            svg_path.write_text(source, encoding='utf-8')
            (Path(tmp_dir) / 'asset.svg').write_text(asset, encoding='utf-8')
            result = SVGQualityChecker().check_file(str(svg_path))
            self.assertTrue(result['passed'])
            _slide_xml, media_files, *_rest = convert_svg_to_slide_shapes(
                svg_path
            )
        self.assertEqual(tuple(media_files), ('s1_img1.svg',))

    def test_complete_wmf_record_stream_is_a_registered_source_format(self):
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<image href="asset.wmf" x="10" y="10" '
            'width="40" height="40"/></svg>'
        )
        # Standard 18-byte METAHEADER plus one 6-byte META_EOF record.
        asset = bytes.fromhex(
            '0100 0900 0003 0c000000 0000 03000000 0000 '
            '03000000 0000'
        )
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'fixture.svg'
            svg_path.write_text(source, encoding='utf-8')
            (Path(tmp_dir) / 'asset.wmf').write_bytes(asset)
            result = SVGQualityChecker().check_file(str(svg_path))
            self.assertTrue(result['passed'])
            _slide_xml, media_files, *_rest = convert_svg_to_slide_shapes(
                svg_path
            )
        self.assertEqual(tuple(media_files), ('s1_img1.wmf',))

    def test_complete_emf_record_stream_is_a_registered_source_format(self):
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<image href="asset.emf" x="10" y="10" '
            'width="40" height="40"/></svg>'
        )
        # Minimal nondegenerate EMR_HEADER plus EMR_EOF. Accept both the
        # MS-EMF all-record count and LibreOffice's post-header count, plus a
        # legal EOF carrying one palette entry.
        cases = ((2, 0), (1, 0), (2, 1))
        for record_count, palette_entries in cases:
            with self.subTest(
                record_count=record_count,
                palette_entries=palette_entries,
            ):
                eof_size = 20 + palette_entries * 4
                asset = bytearray(88 + eof_size)
                asset[0:4] = (1).to_bytes(4, 'little')
                asset[4:8] = (88).to_bytes(4, 'little')
                asset[8:24] = b''.join(
                    value.to_bytes(4, 'little', signed=True)
                    for value in (0, 0, 100, 100)
                )
                asset[24:40] = b''.join(
                    value.to_bytes(4, 'little', signed=True)
                    for value in (0, 0, 2540, 2540)
                )
                asset[40:44] = b' EMF'
                asset[44:48] = (0x00010000).to_bytes(4, 'little')
                asset[48:52] = len(asset).to_bytes(4, 'little')
                asset[52:56] = record_count.to_bytes(4, 'little')
                asset[68:72] = palette_entries.to_bytes(4, 'little')
                asset[72:80] = b''.join(
                    value.to_bytes(4, 'little') for value in (100, 100)
                )
                asset[80:88] = b''.join(
                    value.to_bytes(4, 'little') for value in (25, 25)
                )
                asset[88:92] = (14).to_bytes(4, 'little')
                asset[92:96] = eof_size.to_bytes(4, 'little')
                asset[96:100] = palette_entries.to_bytes(4, 'little')
                asset[100:104] = (16).to_bytes(4, 'little')
                asset[-4:] = eof_size.to_bytes(4, 'little')
                with tempfile.TemporaryDirectory() as tmp_dir:
                    svg_path = Path(tmp_dir) / 'fixture.svg'
                    svg_path.write_text(source, encoding='utf-8')
                    (Path(tmp_dir) / 'asset.emf').write_bytes(asset)
                    result = SVGQualityChecker().check_file(str(svg_path))
                    self.assertTrue(result['passed'])
                    _slide_xml, media_files, *_rest = (
                        convert_svg_to_slide_shapes(svg_path)
                    )
                self.assertEqual(tuple(media_files), ('s1_img1.emf',))

    def test_template_image_token_does_not_bypass_href_cardinality(self):
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" '
            'xmlns:xlink="http://www.w3.org/1999/xlink" '
            'viewBox="0 0 100 100">'
            '<image href="{{IMAGE}}" xlink:href="other.png" '
            'x="10" y="10" width="40" height="40"/></svg>'
        )
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'template.svg'
            svg_path.write_text(source, encoding='utf-8')
            result = SVGQualityChecker(template_mode=True).check_file(
                str(svg_path)
            )
        self.assertFalse(result['passed'])
        self.assertIn(
            'requires exactly one href or xlink:href',
            '\n'.join(result['errors']),
        )

    def test_foreign_namespace_visuals_fail_checker_and_exporter_preflight(self):
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" '
            'xmlns:q="urn:foreign" viewBox="0 0 100 100">'
            f'<q:image href="{self._TINY_PNG_URI}" x="10" y="10" '
            'width="40" height="40"/></svg>'
        )
        self._assert_checker_and_exporter_reject(
            source,
            'must use the SVG namespace',
            'invalid project image',
        )

    def test_unnamespaced_image_fails_checker_and_exporter_preflight(self):
        source = (
            '<svg viewBox="0 0 100 100">'
            f'<image href="{self._TINY_PNG_URI}" x="10" y="10" '
            'width="40" height="40"/></svg>'
        )
        self._assert_checker_and_exporter_reject(
            source,
            'must use the SVG namespace',
            'invalid project image',
        )

    def test_nested_crop_contract_rejects_malformed_transport_wrappers(self):
        valid_inner = (
            f'<image href="{self._TINY_PNG_URI}" x="0" y="0" '
            'width="1" height="1" preserveAspectRatio="none"/>'
        )
        invalid_sources = {
            'indirect image': self._nested_crop_source(
                inner=f'<g>{valid_inner}</g>',
            ),
            'empty href': self._nested_crop_source(
                inner=(
                    '<image href="" x="0" y="0" width="1" height="1" '
                    'preserveAspectRatio="none"/>'
                ),
            ),
            'dual href': self._nested_crop_source(
                inner=(
                    f'<image href="{self._TINY_PNG_URI}" '
                    f'xlink:href="{self._TINY_PNG_URI}" x="0" y="0" '
                    'width="1" height="1" preserveAspectRatio="none"/>'
                ),
            ),
            'non-unit image': self._nested_crop_source(
                inner=valid_inner.replace('width="1"', 'width="2"'),
            ),
            'inner transform': self._nested_crop_source(
                inner=valid_inner.replace('/>', ' transform="translate(1 0)"/>'),
            ),
            'inner style geometry': self._nested_crop_source(
                inner=valid_inner.replace('/>', ' style="width: 1px"/>'),
            ),
            'outer paint': self._nested_crop_source(outer_extra=' fill="#FF0000"'),
            'outer style geometry': self._nested_crop_source(
                outer_extra=' style="width: 200px; height: 100px"',
            ),
            'unknown metadata': self._nested_crop_source(
                outer_extra=' data-pptx-unknown="value"',
            ),
            'wrong outer fit': self._nested_crop_source().replace(
                'viewBox="0.1 0.2 0.7 0.6" preserveAspectRatio="none"',
                'viewBox="0.1 0.2 0.7 0.6" preserveAspectRatio="xMidYMid meet"',
            ),
            'comma viewBox': self._nested_crop_source(view_box='0.1,0.2,0.7,0.6'),
            'exponent viewBox': self._nested_crop_source(view_box='1e-1 0.2 0.7 0.6'),
            'out-of-range percentage': self._nested_crop_source(
                view_box='21474.83648 0 0.1 1',
            ),
            'collapsed crop': self._nested_crop_source(view_box='0 0 0.000001 1'),
            'redundant crop': self._nested_crop_source(view_box='0 0 1 1'),
        }
        for name, source in invalid_sources.items():
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    source,
                    'invalid imported crop wrapper',
                    'invalid nested SVG crop wrapper',
                )

    def test_canonical_generated_spelling_has_no_compatibility_warning(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <rect x="80" y="80" width="300" height="180"
        fill="#FF0000" fill-opacity="0.5"/>
  <text x="80" y="340" font-family="Arial" font-size="28"
        fill="#000080">Canonical</text>
</svg>'''
        )

        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])
        self.assertEqual(result['warnings'], [])

    def test_supported_aliases_are_non_blocking_warnings(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     xmlns:xlink="http://www.w3.org/1999/xlink"
     viewBox="0 0 1280 720">
  <defs>
    <g id="dot"><circle cx="0" cy="0" r="8" fill="#00AA00"/></g>
    <linearGradient id="legacy-gradient">
      <stop offset="0" stop-color="#0000FF" stop-opacity="50%"/>
      <stop offset="1" stop-color="#0000FF"/>
    </linearGradient>
    <pattern id="legacy-pattern" width="8" height="8">
      <rect width="8" height="8" fill="#FFFFFF"/>
      <path d="M0 8 L8 0" stroke="#999999"/>
    </pattern>
  </defs>
  <g id="faded" opacity="0.6">
    <rect x="80" y="80" width="300" height="180"
          fill="rgba(255, 0, 0, 0.5)" fill-opacity="1.2"/>
  </g>
  <rect x="420" y="80" width="300" height="180"
        fill="url(#legacy-pattern)"/>
  <rect x="760" y="80" width="300" height="180"
        fill="url(#legacy-gradient)"/>
  <text x="80" y="340" font-family="Arial" font-size="21pt"
        fill="navy">Aliases</text>
  <use xlink:href="#dot" x="100" y="420"/>
</svg>'''
        )

        warning_text = '\n'.join(result['warnings'])
        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])
        self.assertIn("fill='rgba(255, 0, 0, 0.5)'", warning_text)
        self.assertIn("fill-opacity='1.2'", warning_text)
        self.assertIn("stop-opacity='50%'", warning_text)
        self.assertIn('group opacity', warning_text)
        self.assertIn('font-size value(s) 21pt', warning_text)
        self.assertIn('legacy xlink:href', warning_text)
        self.assertIn('compatible `ltUpDiag` fallback', warning_text)
        self.assertTrue(all('No change is required' in item or 'does not require' in item
                            for item in result['warnings']))

    def test_unsupported_values_remain_errors(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <rect x="80" y="80" width="300" height="180"
        fill="var(--brand)" opacity="bogus" fill-opacity="50%"/>
  <text x="80" y="340" font-family="Arial" font-size="12%">Broken</text>
</svg>'''
        )

        error_text = '\n'.join(result['errors'])
        self.assertFalse(result['passed'])
        self.assertIn('must be a supported color', error_text)
        self.assertIn('must be one finite numeric opacity', error_text)
        self.assertIn('Unsupported font-size', error_text)

    def test_pattern_transform_stays_blocking_without_explicit_preset(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <defs>
    <pattern id="legacy-pattern" width="8" height="8"
             patternTransform="rotate(45)">
      <rect width="8" height="8" fill="#FFFFFF"/>
      <path d="M0 8 L8 0" stroke="#999999"/>
    </pattern>
  </defs>
  <rect x="80" y="80" width="300" height="180"
        fill="url(#legacy-pattern)"/>
</svg>'''
        )

        self.assertFalse(result['passed'])
        self.assertTrue(any('cannot use patternTransform' in item
                            for item in result['errors']))
        self.assertTrue(any('compatible `ltUpDiag` fallback' in item
                            for item in result['warnings']))

    def test_invalid_gradient_contract_blocks_checker_and_exporter(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <linearGradient id="broken" x1="0" x2="2">
      <stop offset="120%"/>
    </linearGradient>
  </defs>
  <rect x="80" y="80" width="300" height="180"
        fill="url(#broken)"/>
</svg>''',
            'requires an explicit stop-color',
            'invalid project gradient',
        )

    def test_converter_string_path_preserves_inline_geometry_diagnostic(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 160 120">
  <rect x="10" y="10" height="40" fill="#2563EB"
        style="width: 60"/>
</svg>'''
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'string-path.svg'
            svg_path.write_text(source, encoding='utf-8')

            with self.assertRaisesRegex(
                SvgNativeConversionError,
                "string-path.svg: inline geometry materialization failed",
            ):
                convert_svg_to_slide_shapes(str(svg_path))

    def test_degenerate_gradient_stroke_blocks_checker_and_exporter(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <linearGradient id="flow" x1="0" y1="0" x2="1" y2="0">
      <stop offset="0" stop-color="#2563EB"/>
      <stop offset="1" stop-color="#10B981"/>
    </linearGradient>
  </defs>
  <path d="M100 200 C260 200 420 200 620 200" fill="none"
        stroke="url(#flow)" stroke-width="40"/>
</svg>''',
            'objectBoundingBox gradients do not include stroke width',
            'invalid project gradient',
        )

    def test_isolated_move_does_not_expand_gradient_stroke_bounds(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <linearGradient id="flow" x1="0" y1="0" x2="1" y2="0">
      <stop offset="0" stop-color="#2563EB"/>
      <stop offset="1" stop-color="#10B981"/>
    </linearGradient>
  </defs>
  <path d="M100 200 H620 M100 201" fill="none"
        stroke="url(#flow)" stroke-width="40"/>
</svg>''',
            'zero intrinsic height',
            'invalid project gradient',
        )

    def test_expanded_use_gradient_stroke_matches_exporter_preflight(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <linearGradient id="flow" x1="0" y1="0" x2="1" y2="0">
      <stop offset="0" stop-color="#2563EB"/>
      <stop offset="1" stop-color="#10B981"/>
    </linearGradient>
    <symbol id="edge" viewBox="0 0 100 20">
      <path d="M0 10 H100" fill="none"/>
    </symbol>
  </defs>
  <g stroke="url(#flow)" stroke-width="10">
    <use href="#edge" x="100" y="100" width="200" height="40"/>
  </g>
</svg>''',
            'zero intrinsic height',
            'invalid project gradient',
        )

    def test_invalid_filter_contract_blocks_checker_and_exporter(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <filter id="broken">
      <feGaussianBlur stdDeviation="not-a-number"/>
    </filter>
  </defs>
  <rect x="80" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#broken)"/>
</svg>''',
            'stdDeviation must be a finite number',
            'invalid project filter',
        )

    def test_filter_effect_geometry_must_be_explicit(self):
        cases = (
            (
                'Gaussian blur standard deviation',
                '<feGaussianBlur/>',
                '<feGaussianBlur> requires explicit stdDeviation',
            ),
            (
                'drop shadow standard deviation',
                '<feDropShadow dx="0" dy="4"/>',
                '<feDropShadow> requires explicit stdDeviation',
            ),
            (
                'drop shadow horizontal offset',
                '<feDropShadow dy="4" stdDeviation="6"/>',
                '<feDropShadow> requires explicit dx',
            ),
            (
                'drop shadow vertical offset',
                '<feDropShadow dx="0" stdDeviation="6"/>',
                '<feDropShadow> requires explicit dy',
            ),
        )
        for label, primitive, expected in cases:
            with self.subTest(label=label):
                self._assert_checker_and_exporter_reject(
                    f'''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs><filter id="effect">{primitive}</filter></defs>
  <rect x="80" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#effect)"/>
</svg>''',
                    expected,
                    'invalid project filter',
                )

    def test_filter_alpha_transfer_slope_must_be_explicit(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <filter id="effect">
      <feGaussianBlur stdDeviation="6"/>
      <feComponentTransfer>
        <feFuncA type="linear"/>
      </feComponentTransfer>
    </filter>
  </defs>
  <rect x="80" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#effect)"/>
</svg>''',
            '<feFuncA> requires explicit slope',
            'invalid project filter',
        )

    def test_filter_alpha_transfer_rejects_unmapped_intercept(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <filter id="effect">
      <feGaussianBlur stdDeviation="6"/>
      <feComponentTransfer>
        <feFuncA type="linear" slope="0.4" intercept="0.5"/>
      </feComponentTransfer>
    </filter>
  </defs>
  <rect x="80" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#effect)"/>
</svg>''',
            '<feFuncA> intercept is unsupported',
            'invalid project filter',
        )

    def test_filter_blur_rejects_unmapped_edge_mode(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <filter id="effect">
      <feGaussianBlur stdDeviation="6" edgeMode="wrap"/>
      <feFlood flood-color="#000000" flood-opacity="0.4"/>
    </filter>
  </defs>
  <rect x="80" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#effect)"/>
</svg>''',
            '<feGaussianBlur> edgeMode is unsupported',
            'invalid project filter',
        )

    def test_filter_rejects_object_bounding_box_primitive_units(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <filter id="effect" primitiveUnits="objectBoundingBox">
      <feGaussianBlur stdDeviation="0.1"/>
      <feFlood flood-color="#000000" flood-opacity="0.4"/>
    </filter>
  </defs>
  <rect x="80" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#effect)"/>
</svg>''',
            'primitiveUnits must be userSpaceOnUse',
            'invalid project filter',
        )

        user_space = ET.fromstring('''<svg xmlns="http://www.w3.org/2000/svg">
  <defs>
    <filter id="effect" primitiveUnits="userSpaceOnUse">
      <feGaussianBlur stdDeviation="6"/>
      <feFlood flood-color="#000000" flood-opacity="0.4"/>
    </filter>
  </defs>
</svg>''')
        self.assertEqual(project_filter_errors(user_space), [])

    def test_filter_flood_opacity_must_be_explicit(self):
        cases = (
            (
                'drop-shadow shorthand',
                '<feDropShadow dx="0" dy="4" stdDeviation="6" '
                'flood-color="#000000"/>',
                '<feDropShadow> requires explicit flood-opacity',
            ),
            (
                'expanded flood primitive',
                '<feGaussianBlur stdDeviation="6"/>'
                '<feFlood flood-color="#000000"/>',
                '<feFlood> requires explicit flood-opacity',
            ),
        )
        for label, graph, expected in cases:
            with self.subTest(label=label):
                self._assert_checker_and_exporter_reject(
                    f'''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs><filter id="effect">{graph}</filter></defs>
  <rect x="80" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#effect)"/>
</svg>''',
                    expected,
                    'invalid project filter',
                )

        styled = ET.fromstring('''<svg xmlns="http://www.w3.org/2000/svg">
  <defs>
    <filter id="effect">
      <feGaussianBlur stdDeviation="6"/>
      <feFlood flood-color="#000000" style="flood-opacity:0.2"/>
    </filter>
  </defs>
</svg>''')
        self.assertEqual(project_filter_errors(styled), [])

    def test_filter_coordinates_block_drawingml_overflow(self):
        cases = (
            (
                'glow radius',
                '<feGaussianBlur stdDeviation="2863311530"/>',
                'DrawingML rad must map within',
            ),
            (
                'shadow blur radius',
                '<feDropShadow dx="1" dy="0" '
                'stdDeviation="1431655765" flood-opacity="0.2"/>',
                'DrawingML blurRad must map within',
            ),
            (
                'shadow distance',
                '<feDropShadow dx="2863311530" dy="0" '
                'stdDeviation="1" flood-opacity="0.2"/>',
                'DrawingML dist must map within',
            ),
            (
                'mapped infinity',
                '<feGaussianBlur stdDeviation="1e308"/>',
                'DrawingML rad must be finite after EMU mapping',
            ),
        )
        for label, primitive, expected in cases:
            with self.subTest(label=label):
                self._assert_checker_and_exporter_reject(
                    f'''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs><filter id="effect">{primitive}</filter></defs>
  <rect x="80" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#effect)"/>
</svg>''',
                    expected,
                    'invalid project filter',
                )

    def test_filter_coordinates_accept_drawingml_maximum(self):
        glow_std_dev = OOXML_COORDINATE_MAX / EMU_PER_PX
        shadow_std_dev = OOXML_COORDINATE_MAX / (2 * EMU_PER_PX)
        svg = f'''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <filter id="glow">
      <feGaussianBlur stdDeviation="{glow_std_dev!r}"/>
    </filter>
    <filter id="shadow">
      <feDropShadow dx="{glow_std_dev!r}" dy="0"
                    stdDeviation="{shadow_std_dev!r}"
                    flood-opacity="0.2"/>
    </filter>
  </defs>
  <rect x="80" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#glow)"/>
  <rect x="480" y="80" width="300" height="180"
        fill="#FFFFFF" filter="url(#shadow)"/>
</svg>'''
        root = ET.fromstring(svg)
        self.assertEqual(project_filter_errors(root), [])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'effect-boundary.svg'
            svg_path.write_text(svg, encoding='utf-8')
            checker_result = SVGQualityChecker().check_file(str(svg_path))
            self.assertEqual(checker_result['errors'], [])
            slide_xml, *_rest = convert_svg_to_slide_shapes(svg_path)
        self.assertIn(
            f'<a:glow rad="{OOXML_COORDINATE_MAX}">',
            slide_xml,
        )
        self.assertIn(
            f'<a:outerShdw blurRad="{OOXML_COORDINATE_MAX}" '
            f'dist="{OOXML_COORDINATE_MAX}"',
            slide_xml,
        )

    def test_imported_preset_preview_may_mirror_carrier_filter(self):
        rendered = render_preset_geometry(
            'cube',
            Xfrm(x=100, y=120, w=300, h=160),
        )
        semantic_attrs = {
            'data-pptx-object': 'shape',
            'data-pptx-shape-id': '2',
            'data-pptx-shape-scope': 'slide',
            'data-pptx-frame': '100 120 300 160',
            'data-pptx-prst': 'cube',
        }
        markup = serialize_preset_layers(
            rendered.paths,
            semantic_attrs,
            {
                'fill': '#FFFFFF',
                'stroke': 'none',
                'filter': 'url(#shadow)',
            },
        )
        svg = f'''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <defs>
    <filter id="shadow">
      <feDropShadow dx="0" dy="4" stdDeviation="6"
                    flood-color="#000000" flood-opacity="0.2"/>
    </filter>
  </defs>
  <g id="imported-cube" data-pptx-object="shape"
     data-pptx-shape-id="2" data-pptx-shape-scope="slide"
     data-pptx-frame="100 120 300 160" data-pptx-prst="cube"
     data-pptx-preview-sha256="{markup.preview_hash}">
    {markup.markup}
  </g>
</svg>'''
        root = ET.fromstring(svg)
        self.assertEqual(project_filter_errors(root), [])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'imported-effect.svg'
            svg_path.write_text(svg, encoding='utf-8')
            with contextlib.redirect_stdout(io.StringIO()):
                checker_result = SVGQualityChecker().check_file(str(svg_path))
            self.assertEqual(checker_result['errors'], [])
            slide_xml, *_rest = convert_svg_to_slide_shapes(svg_path)
        self.assertEqual(slide_xml.count('<p:sp>'), 1)
        self.assertNotIn('<p:grpSp>', slide_xml)
        self.assertEqual(slide_xml.count('<a:outerShdw'), 1)

        preview = next(
            elem
            for elem in root.iter()
            if elem.get('data-pptx-part') == 'geometry-preview'
        )
        preview.set('filter', 'url(#different)')
        self.assertTrue(any(
            '<g> cannot use filter' in error
            for error in project_filter_errors(root)
        ))

        def imported_parts(
            document: ET.Element,
        ) -> tuple[ET.Element, ET.Element, ET.Element]:
            logical = next(
                elem
                for elem in document.iter()
                if elem.get('id') == 'imported-cube'
            )
            native_carrier = next(
                elem
                for elem in logical
                if elem.get('data-pptx-part') == 'geometry'
            )
            visible_preview = next(
                elem
                for elem in logical
                if elem.get('data-pptx-part') == 'geometry-preview'
            )
            return logical, native_carrier, visible_preview

        def refresh_preview_hash(
            logical: ET.Element,
            native_carrier: ET.Element,
        ) -> None:
            digest = svg_preset_preview_fingerprint(logical)
            logical.set('data-pptx-preview-sha256', digest)
            native_carrier.set('data-pptx-preview-sha256', digest)

        foreign_child_root = ET.fromstring(svg)
        logical, carrier, visible_preview = imported_parts(foreign_child_root)
        ET.SubElement(visible_preview, '{http://www.w3.org/2000/svg}rect', {
            'x': '100',
            'y': '120',
            'width': '10',
            'height': '10',
        })
        refresh_preview_hash(logical, carrier)
        self.assertTrue(any(
            '<g> cannot use filter' in error
            for error in project_filter_errors(foreign_child_root)
        ))

        duplicate_preview_root = ET.fromstring(svg)
        logical, carrier, visible_preview = imported_parts(duplicate_preview_root)
        logical.append(ET.fromstring(ET.tostring(visible_preview)))
        refresh_preview_hash(logical, carrier)
        self.assertTrue(any(
            '<g> cannot use filter' in error
            for error in project_filter_errors(duplicate_preview_root)
        ))

        stale_root = ET.fromstring(svg)
        logical, carrier, _visible_preview = imported_parts(stale_root)
        logical.set('data-pptx-preview-sha256', '0' * 64)
        carrier.set('data-pptx-preview-sha256', '0' * 64)
        self.assertTrue(any(
            '<g> cannot use filter' in error
            for error in project_filter_errors(stale_root)
        ))

        mismatched_carrier_root = ET.fromstring(svg)
        _logical, carrier, _visible_preview = imported_parts(
            mismatched_carrier_root
        )
        carrier.set('data-pptx-prst', 'rect')
        self.assertTrue(any(
            '<g> cannot use filter' in error
            for error in project_filter_errors(mismatched_carrier_root)
        ))

        ordinary = ET.fromstring('''<svg xmlns="http://www.w3.org/2000/svg">
  <defs><filter id="shadow"><feDropShadow dx="0" dy="4"
    stdDeviation="6"/></filter></defs>
  <g filter="url(#shadow)"><rect x="0" y="0" width="10" height="10"/></g>
</svg>''')
        self.assertTrue(any(
            '<g> cannot use filter' in error
            for error in project_filter_errors(ordinary)
        ))

    def test_pptx_effect_import_is_closed_and_fail_closed(self):
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        def effect_result(container: str):
            sp_pr = ET.fromstring(
                f'<p:spPr xmlns:p="http://schemas.openxmlformats.org/'
                f'presentationml/2006/main" xmlns:a="{dml}">'
                f'{container}</p:spPr>'
            )
            return convert_effects(sp_pr, None)

        outer = effect_result(f'''<a:effectLst xmlns:a="{dml}">
  <a:outerShdw blurRad="38100" dist="38100" dir="5400000">
    <a:srgbClr val="000000"><a:alpha val="50000"/></a:srgbClr>
  </a:outerShdw>
</a:effectLst>''')
        self.assertIsNotNone(outer.filter_id)
        self.assertEqual(dict(outer.metadata), {})
        self.assertIn('stdDeviation="2"', ''.join(outer.defs))

        glow = effect_result(f'''<a:effectLst xmlns:a="{dml}">
  <a:glow rad="95250"><a:srgbClr val="2563EB"/></a:glow>
</a:effectLst>''')
        glow_def = ''.join(glow.defs)
        self.assertIn('stdDeviation="10"', glow_def)
        self.assertNotIn('feMorphology', glow_def)

        zero_glow = effect_result(f'''<a:effectLst xmlns:a="{dml}">
  <a:glow rad="0"><a:srgbClr val="2563EB"/></a:glow>
</a:effectLst>''')
        self.assertIsNotNone(zero_glow.filter_id)
        self.assertEqual(dict(zero_glow.metadata), {})
        self.assertIn('stdDeviation="0"', ''.join(zero_glow.defs))

        high_saturation = effect_result(f'''<a:effectLst xmlns:a="{dml}">
  <a:glow rad="95250"><a:srgbClr val="2563EB">
    <a:satMod val="175000"/>
  </a:srgbClr></a:glow>
</a:effectLst>''')
        self.assertIsNotNone(high_saturation.filter_id)
        self.assertEqual(dict(high_saturation.metadata), {})

        system_color = effect_result(f'''<a:effectLst xmlns:a="{dml}">
  <a:glow rad="9525"><a:sysClr val="windowText" lastClr="112233"/>
  </a:glow>
</a:effectLst>''')
        self.assertIsNotNone(system_color.filter_id)
        self.assertEqual(dict(system_color.metadata), {})
        self.assertIn('flood-color="#112233"', ''.join(system_color.defs))

        unsupported_cases = {
            'innerShdw': (
                f'<a:effectLst xmlns:a="{dml}"><a:innerShdw '
                'blurRad="38100" dist="38100" dir="5400000">'
                '<a:srgbClr val="000000"/></a:innerShdw></a:effectLst>'
            ),
            'blur': (
                f'<a:effectLst xmlns:a="{dml}"><a:blur rad="38100"/>'
                '</a:effectLst>'
            ),
            'softEdge': (
                f'<a:effectLst xmlns:a="{dml}"><a:softEdge rad="38100"/>'
                '</a:effectLst>'
            ),
            'reflection': (
                f'<a:effectLst xmlns:a="{dml}"><a:reflection blurRad="0"/>'
                '</a:effectLst>'
            ),
        }
        for effect_name, container in unsupported_cases.items():
            with self.subTest(effect=effect_name):
                result = effect_result(container)
                self.assertIsNone(result.filter_id)
                self.assertEqual(result.defs, ())
                self.assertEqual(
                    dict(result.metadata)['data-pptx-effect-reason'],
                    f'unsupported-effect:{effect_name}',
                )

        zero_offset = effect_result(
            f'<a:effectLst xmlns:a="{dml}"><a:outerShdw '
            'blurRad="38100" dist="0" dir="0">'
            '<a:srgbClr val="000000"/></a:outerShdw></a:effectLst>'
        )
        self.assertIn(
            'offset-is-not-classifiable',
            dict(zero_offset.metadata)['data-pptx-effect-reason'],
        )
        invalid_glow = effect_result(
            f'<a:effectLst xmlns:a="{dml}"><a:glow rad="1.5"/>'
            '</a:effectLst>'
        )
        self.assertIn(
            "rad='1.5'",
            dict(invalid_glow.metadata)['data-pptx-effect-reason'],
        )
        multiple = effect_result(
            f'<a:effectLst xmlns:a="{dml}"><a:outerShdw dist="38100"/>'
            '<a:reflection/></a:effectLst>'
        )
        self.assertEqual(
            dict(multiple.metadata)['data-pptx-effect-reason'],
            'multiple-effects:outerShdw,reflection',
        )
        effect_dag = effect_result(f'<a:effectDag xmlns:a="{dml}"/>')
        self.assertEqual(
            dict(effect_dag.metadata)['data-pptx-effect-reason'],
            'unsupported-effect-container:effectDag',
        )

        invalid_subset_cases = {
            'duplicate containers': ((
                f'<a:effectLst xmlns:a="{dml}"/>'
                f'<a:effectLst xmlns:a="{dml}"/>'
            ), 'multiple-effect-containers'),
            'foreign container namespace': ((
                '<x:effectLst xmlns:x="urn:foreign"/>'
            ), 'invalid-effect-container-namespace'),
            'foreign effect namespace': ((
                f'<a:effectLst xmlns:a="{dml}" xmlns:x="urn:foreign">'
                '<x:outerShdw/></a:effectLst>'
            ), 'invalid-effect-namespace'),
            'missing color': ((
                f'<a:effectLst xmlns:a="{dml}"><a:outerShdw '
                'dist="38100"/></a:effectLst>'
            ), 'missing-color'),
            'missing glow radius': ((
                f'<a:effectLst xmlns:a="{dml}"><a:glow>'
                '<a:srgbClr val="000000"/></a:glow></a:effectLst>'
            ), 'missing-rad'),
            'empty glow radius': ((
                f'<a:effectLst xmlns:a="{dml}"><a:glow rad="">'
                '<a:srgbClr val="000000"/></a:glow></a:effectLst>'
            ), "rad=''"),
            'invalid srgb color': ((
                f'<a:effectLst xmlns:a="{dml}"><a:outerShdw '
                'dist="38100"><a:srgbClr val="bogus"/>'
                '</a:outerShdw></a:effectLst>'
            ), 'invalid-color:srgbClr'),
            'missing system color value': ((
                f'<a:effectLst xmlns:a="{dml}"><a:glow rad="38100">'
                '<a:sysClr lastClr="112233"/></a:glow></a:effectLst>'
            ), 'invalid-color:sysClr'),
            'empty system color value': ((
                f'<a:effectLst xmlns:a="{dml}"><a:glow rad="38100">'
                '<a:sysClr val="" lastClr="112233"/>'
                '</a:glow></a:effectLst>'
            ), 'invalid-color:sysClr'),
            'missing system color fallback': ((
                f'<a:effectLst xmlns:a="{dml}"><a:glow rad="38100">'
                '<a:sysClr val="windowText"/></a:glow></a:effectLst>'
            ), 'unresolvable-color:sysClr'),
            'invalid system color fallback': ((
                f'<a:effectLst xmlns:a="{dml}"><a:glow rad="38100">'
                '<a:sysClr val="windowText" lastClr="bogus"/>'
                '</a:glow></a:effectLst>'
            ), 'unresolvable-color:sysClr'),
            'multiple colors': ((
                f'<a:effectLst xmlns:a="{dml}"><a:outerShdw '
                'dist="38100"><a:srgbClr val="000000"/>'
                '<a:srgbClr val="FFFFFF"/></a:outerShdw></a:effectLst>'
            ), 'multiple-colors'),
            'invalid alpha': ((
                f'<a:effectLst xmlns:a="{dml}"><a:glow rad="38100">'
                '<a:srgbClr val="000000"><a:alpha val="bad"/>'
                '</a:srgbClr></a:glow></a:effectLst>'
            ), 'invalid-alpha-val'),
            'unsupported color modifier': ((
                f'<a:effectLst xmlns:a="{dml}"><a:glow rad="38100">'
                '<a:srgbClr val="000000"><a:gamma/>'
                '</a:srgbClr></a:glow></a:effectLst>'
            ), 'unsupported-color-modifier:gamma'),
            'oversized distance': ((
                f'<a:effectLst xmlns:a="{dml}"><a:outerShdw '
                f'dist="{"9" * 400}"><a:srgbClr val="000000"/>'
                '</a:outerShdw></a:effectLst>'
            ), 'dist='),
            'oversized radius': ((
                f'<a:effectLst xmlns:a="{dml}"><a:glow '
                f'rad="{"9" * 400}"><a:srgbClr val="000000"/>'
                '</a:glow></a:effectLst>'
            ), 'rad='),
            'oversized direction': ((
                f'<a:effectLst xmlns:a="{dml}"><a:outerShdw '
                f'dist="38100" dir="{"9" * 400}">'
                '<a:srgbClr val="000000"/></a:outerShdw></a:effectLst>'
            ), 'dir='),
        }
        for case_name, (container, expected_reason) in invalid_subset_cases.items():
            with self.subTest(case=case_name):
                result = effect_result(container)
                self.assertIsNone(result.filter_id)
                self.assertEqual(result.defs, ())
                self.assertIn(
                    expected_reason,
                    dict(result.metadata)['data-pptx-effect-reason'],
                )

        def round_trip_supported(result):
            filter_id = result.filter_id
            svg = (
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 1280 720"><defs>'
                + ''.join(result.defs)
                + f'</defs><rect x="10" y="10" width="100" height="60" '
                f'fill="#FFFFFF" filter="url(#{filter_id})"/></svg>'
            )
            with tempfile.TemporaryDirectory() as tmp_dir:
                svg_path = Path(tmp_dir) / 'tiny-effect.svg'
                svg_path.write_text(svg, encoding='utf-8')
                output, *_rest = convert_svg_to_slide_shapes(svg_path)
                return output

        tiny_outer = effect_result(
            f'<a:effectLst xmlns:a="{dml}"><a:outerShdw '
            'blurRad="1" dist="100" dir="0">'
            '<a:srgbClr val="000000"/></a:outerShdw></a:effectLst>'
        )
        tiny_outer_xml = round_trip_supported(tiny_outer)
        self.assertIn('<a:outerShdw blurRad="1" dist="100"', tiny_outer_xml)
        self.assertNotIn('<a:glow', tiny_outer_xml)

        low_alpha_outer = effect_result(
            f'<a:effectLst xmlns:a="{dml}"><a:outerShdw '
            'blurRad="9525" dist="9525" dir="0">'
            '<a:srgbClr val="000000"><a:alpha val="1"/>'
            '</a:srgbClr></a:outerShdw></a:effectLst>'
        )
        self.assertIn(
            '<a:alpha val="1"/>',
            round_trip_supported(low_alpha_outer),
        )

        tiny_glow = effect_result(
            f'<a:effectLst xmlns:a="{dml}"><a:glow rad="1">'
            '<a:srgbClr val="000000"/></a:glow></a:effectLst>'
        )
        self.assertIn('<a:glow rad="1">', round_trip_supported(tiny_glow))

        for alpha in (0, 1, 2, 3, 7, 13, 29, 12345, 50001, 99999, 100000):
            with self.subTest(glow_alpha=alpha):
                alpha_glow = effect_result(
                    f'<a:effectLst xmlns:a="{dml}"><a:glow rad="9525">'
                    '<a:srgbClr val="000000">'
                    f'<a:alpha val="{alpha}"/></a:srgbClr>'
                    '</a:glow></a:effectLst>'
                )
                self.assertIn(
                    f'<a:alpha val="{alpha}"/>',
                    round_trip_supported(alpha_glow),
                )

        saturated_alpha = effect_result(
            f'<a:effectLst xmlns:a="{dml}"><a:glow rad="9525">'
            '<a:srgbClr val="000000"><a:alphaMod val="200000"/>'
            '<a:alphaOff val="-50000"/></a:srgbClr>'
            '</a:glow></a:effectLst>'
        )
        self.assertIn(
            'flood-opacity="0.5"',
            ''.join(saturated_alpha.defs),
        )
        self.assertIn(
            '<a:alpha val="50000"/>',
            round_trip_supported(saturated_alpha),
        )

        marked_svg = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <rect x="80" y="80" width="300" height="180" fill="#FFFFFF"
        data-pptx-effect-status="unsupported"
        data-pptx-effect-reason="unsupported-effect:reflection"/>
</svg>'''
        self._assert_checker_and_exporter_reject(
            marked_svg,
            'unsupported source PPTX effect: unsupported-effect:reflection',
            'unsupported imported PPTX effect',
        )

    def test_real_pptx_effect_import_matrix(self):
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        containers = [
            f'''<a:effectLst xmlns:a="{dml}">
  <a:outerShdw blurRad="38100" dist="38100" dir="5400000">
    <a:srgbClr val="000000"><a:alpha val="50000"/></a:srgbClr>
  </a:outerShdw>
</a:effectLst>''',
            f'''<a:effectLst xmlns:a="{dml}">
  <a:glow rad="95250"><a:srgbClr val="2563EB"/></a:glow>
</a:effectLst>''',
            f'''<a:effectLst xmlns:a="{dml}"><a:innerShdw
  blurRad="38100" dist="38100" dir="5400000"><a:srgbClr val="000000"/>
</a:innerShdw></a:effectLst>''',
            f'<a:effectLst xmlns:a="{dml}"><a:blur rad="38100"/></a:effectLst>',
            f'<a:effectLst xmlns:a="{dml}"><a:softEdge rad="38100"/></a:effectLst>',
            f'<a:effectLst xmlns:a="{dml}"><a:reflection/></a:effectLst>',
            f'''<a:effectLst xmlns:a="{dml}"><a:outerShdw
  blurRad="38100" dist="0" dir="0"><a:srgbClr val="000000"/>
</a:outerShdw></a:effectLst>''',
            f'''<a:effectLst xmlns:a="{dml}">
  <a:outerShdw blurRad="38100" dist="38100" dir="5400000">
    <a:srgbClr val="000000"/>
  </a:outerShdw><a:reflection/>
</a:effectLst>''',
            f'<a:effectDag xmlns:a="{dml}"/>',
        ]
        with tempfile.TemporaryDirectory() as tmp_dir:
            root_dir = Path(tmp_dir)
            pptx_path = root_dir / 'effects.pptx'
            self._write_effect_fixture(pptx_path, containers)
            imported = convert_pptx_to_svg(
                pptx_path,
                options=ConvertOptions(inheritance_mode='flat'),
            )
            self.assertEqual(len(imported.slides), len(containers))

            outer_svg = imported.slides[0].svg
            outer_root = ET.fromstring(outer_svg)
            self.assertEqual(project_effect_status_errors(outer_root), [])
            self.assertEqual(project_filter_errors(outer_root), [])
            outer_path = root_dir / 'outer.svg'
            outer_path.write_text(outer_svg, encoding='utf-8')
            outer_xml, *_rest = convert_svg_to_slide_shapes(outer_path)
            self.assertEqual(outer_xml.count('<a:outerShdw'), 1)

            glow_svg = imported.slides[1].svg
            glow_root = ET.fromstring(glow_svg)
            self.assertEqual(project_effect_status_errors(glow_root), [])
            self.assertEqual(project_filter_errors(glow_root), [])
            self.assertNotIn('feMorphology', glow_svg)
            self.assertIn('stdDeviation="10"', glow_svg)
            glow_path = root_dir / 'glow.svg'
            glow_path.write_text(glow_svg, encoding='utf-8')
            glow_xml, *_rest = convert_svg_to_slide_shapes(glow_path)
            self.assertIn('<a:glow rad="95250">', glow_xml)

            expected_reasons = (
                'unsupported-effect:innerShdw',
                'unsupported-effect:blur',
                'unsupported-effect:softEdge',
                'unsupported-effect:reflection',
                'offset-is-not-classifiable',
                'multiple-effects:outerShdw,reflection',
                'unsupported-effect-container:effectDag',
            )
            for artifact, expected_reason in zip(
                imported.slides[2:],
                expected_reasons,
            ):
                with self.subTest(reason=expected_reason):
                    artifact_root = ET.fromstring(artifact.svg)
                    status_errors = project_effect_status_errors(artifact_root)
                    self.assertEqual(len(status_errors), 1)
                    self.assertTrue(any(
                        expected_reason in error
                        for error in status_errors
                    ))
                    marked = [
                        elem
                        for elem in artifact_root.iter()
                        if elem.get('data-pptx-effect-status') == 'unsupported'
                    ]
                    self.assertEqual(len(marked), 2)
                    self.assertEqual(
                        {elem.get('data-pptx-effect-reason') for elem in marked},
                        {next(
                            elem.get('data-pptx-effect-reason')
                            for elem in marked
                        )},
                    )
                    self.assertTrue(any(
                        elem.get('data-pptx-object') == 'shape'
                        and elem.get('data-pptx-prst') == 'cube'
                        for elem in marked
                    ))
                    self.assertTrue(any(
                        elem.get('data-pptx-part') == 'geometry'
                        and elem.get('data-pptx-prst') == 'cube'
                        for elem in marked
                    ))
                    self.assertNotIn('<filter', artifact.svg)
                    artifact_path = root_dir / f'unsupported-{artifact.index}.svg'
                    artifact_path.write_text(artifact.svg, encoding='utf-8')
                    with self.assertRaisesRegex(
                        SvgNativeConversionError,
                        'unsupported imported PPTX effect',
                    ):
                        convert_svg_to_slide_shapes(artifact_path)

    def test_vertical_text_run_effect_import_fails_closed(self):
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        effects = [
            f'''<a:effectLst xmlns:a="{dml}">
  <a:outerShdw blurRad="38100" dist="38100" dir="5400000">
    <a:srgbClr val="000000"/>
  </a:outerShdw>
</a:effectLst>''',
            f'<a:effectLst xmlns:a="{dml}"/>',
            f'<a:effectDag xmlns:a="{dml}"/>',
            None,
        ]
        with tempfile.TemporaryDirectory() as tmp_dir:
            root_dir = Path(tmp_dir)
            pptx_path = root_dir / 'vertical-run-effects.pptx'
            self._write_vertical_run_effect_fixture(pptx_path, effects)
            imported = convert_pptx_to_svg(
                pptx_path,
                options=ConvertOptions(inheritance_mode='flat'),
            )
            self.assertEqual(len(imported.slides), 4)

            effect_svg = imported.slides[0].svg
            effect_root = ET.fromstring(effect_svg)
            self.assertIn('甲', ''.join(effect_root.itertext()))
            self.assertIn('乙', ''.join(effect_root.itertext()))
            self.assertFalse(any(
                elem.get('data-pptx-part') == 'txbody'
                for elem in effect_root.iter()
            ))
            marked = [
                elem
                for elem in effect_root.iter()
                if elem.get('data-pptx-effect-status') == 'unsupported'
            ]
            self.assertEqual(len(marked), 2)
            self.assertEqual(
                {elem.get('data-pptx-effect-reason') for elem in marked},
                {'unsupported-run-effect-route:vertical-text'},
            )
            status_errors = project_effect_status_errors(effect_root)
            self.assertEqual(len(status_errors), 1)
            self.assertIn(
                'unsupported-run-effect-route:vertical-text',
                status_errors[0],
            )
            self._assert_checker_and_exporter_reject(
                effect_svg,
                'unsupported-run-effect-route:vertical-text',
                'unsupported imported PPTX effect',
            )

            for artifact in imported.slides[1:]:
                with self.subTest(slide=artifact.index):
                    root = ET.fromstring(artifact.svg)
                    self.assertEqual(project_effect_status_errors(root), [])
                    self.assertFalse(any(
                        elem.get('data-pptx-effect-status') is not None
                        for elem in root.iter()
                    ))
                    svg_path = root_dir / f'vertical-{artifact.index}.svg'
                    svg_path.write_text(artifact.svg, encoding='utf-8')
                    checker_result = SVGQualityChecker().check_file(
                        str(svg_path)
                    )
                    self.assertTrue(checker_result['passed'])
                    convert_svg_to_slide_shapes(svg_path)

    def test_compound_shape_and_run_effect_reasons_are_preserved(self):
        reason_attr = 'data-pptx-effect-reason'
        normalized = unsupported_effect_metadata(
            'reason-b',
            'reason-a',
            'reason-b',
        )
        self.assertEqual(
            json.loads(normalized[reason_attr]),
            ['reason-a', 'reason-b'],
        )
        self.assertEqual(
            unsupported_effect_metadata(normalized[reason_attr], 'reason-a'),
            normalized,
        )

        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        run_effect = f'''<a:effectLst xmlns:a="{dml}">
  <a:glow rad="38100"><a:srgbClr val="2563EB"/></a:glow>
</a:effectLst>'''
        shape_effect = (
            f'<a:effectLst xmlns:a="{dml}"><a:reflection/>'
            '</a:effectLst>'
        )
        expected_reasons = [
            'unsupported-effect:reflection',
            'unsupported-run-effect-route:vertical-text',
        ]
        with tempfile.TemporaryDirectory() as tmp_dir:
            root_dir = Path(tmp_dir)
            pptx_path = root_dir / 'compound-run-effects.pptx'
            self._write_vertical_run_effect_fixture(
                pptx_path,
                [run_effect],
                shape_effects=[shape_effect],
            )
            imported = convert_pptx_to_svg(
                pptx_path,
                options=ConvertOptions(inheritance_mode='flat'),
            )
            self.assertEqual(len(imported.slides), 1)
            artifact = imported.slides[0]
            root = ET.fromstring(artifact.svg)
            marked = [
                elem
                for elem in root.iter()
                if elem.get('data-pptx-effect-status') == 'unsupported'
            ]
            self.assertEqual(len(marked), 2)
            self.assertEqual(
                {
                    tuple(json.loads(elem.get(reason_attr) or '[]'))
                    for elem in marked
                },
                {tuple(expected_reasons)},
            )
            status_errors = project_effect_status_errors(root)
            self.assertEqual(len(status_errors), 1)
            for reason in expected_reasons:
                self.assertIn(reason, status_errors[0])
            self.assertNotIn('<filter', artifact.svg)
            self._assert_checker_and_exporter_reject(
                artifact.svg,
                expected_reasons[0],
                'unsupported imported PPTX effect',
            )

    def test_relationship_text_run_effect_import_fails_closed(self):
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        run_effect = f'''<a:effectLst xmlns:a="{dml}">
  <a:glow rad="38100"><a:srgbClr val="2563EB"/></a:glow>
</a:effectLst>'''
        cases = [
            (run_effect, True),
            (f'<a:effectLst xmlns:a="{dml}"/>', True),
            (f'<a:effectDag xmlns:a="{dml}"/>', True),
            (None, True),
            (run_effect, False),
        ]
        with tempfile.TemporaryDirectory() as tmp_dir:
            root_dir = Path(tmp_dir)
            pptx_path = root_dir / 'relationship-run-effects.pptx'
            self._write_relationship_run_effect_fixture(pptx_path, cases)
            imported = convert_pptx_to_svg(
                pptx_path,
                options=ConvertOptions(inheritance_mode='flat'),
            )
            self.assertEqual(len(imported.slides), 5)

            effect_svg = imported.slides[0].svg
            effect_root = ET.fromstring(effect_svg)
            self.assertIn('LINK 1', ''.join(effect_root.itertext()))
            self.assertFalse(any(
                elem.get('data-pptx-part') == 'txbody'
                for elem in effect_root.iter()
            ))
            marked = [
                elem
                for elem in effect_root.iter()
                if elem.get('data-pptx-effect-status') == 'unsupported'
            ]
            self.assertEqual(len(marked), 2)
            self.assertEqual(
                {elem.get('data-pptx-effect-reason') for elem in marked},
                {'unsupported-run-effect-route:relationship-bearing-text'},
            )
            status_errors = project_effect_status_errors(effect_root)
            self.assertEqual(len(status_errors), 1)
            self.assertIn(
                'unsupported-run-effect-route:relationship-bearing-text',
                status_errors[0],
            )
            self._assert_checker_and_exporter_reject(
                effect_svg,
                'unsupported-run-effect-route:relationship-bearing-text',
                'unsupported imported PPTX effect',
            )

            for artifact in imported.slides[1:4]:
                with self.subTest(slide=artifact.index):
                    plain_root = ET.fromstring(artifact.svg)
                    self.assertIn(
                        f'LINK {artifact.index}',
                        ''.join(plain_root.itertext()),
                    )
                    self.assertFalse(any(
                        elem.get('data-pptx-part') == 'txbody'
                        for elem in plain_root.iter()
                    ))
                    self.assertEqual(
                        project_effect_status_errors(plain_root),
                        [],
                    )
                    self.assertFalse(any(
                        elem.get('data-pptx-effect-status') is not None
                        for elem in plain_root.iter()
                    ))
                    plain_path = (
                        root_dir / f'relationship-{artifact.index}.svg'
                    )
                    plain_path.write_text(artifact.svg, encoding='utf-8')
                    checker_result = SVGQualityChecker().check_file(
                        str(plain_path)
                    )
                    self.assertTrue(checker_result['passed'])
                    convert_svg_to_slide_shapes(plain_path)

            native_svg = imported.slides[4].svg
            native_root = ET.fromstring(native_svg)
            self.assertTrue(any(
                elem.get('data-pptx-part') == 'txbody'
                for elem in native_root.iter()
            ))
            self.assertEqual(project_effect_status_errors(native_root), [])
            native_path = root_dir / 'relationship-control.svg'
            native_path.write_text(native_svg, encoding='utf-8')
            native_xml = convert_svg_to_slide_shapes(native_path)[0]
            self.assertEqual(native_xml.count('<a:glow'), 1)

    def test_inherited_text_run_effect_import_fails_closed(self):
        expected_reasons = {
            1: 'unsupported-run-effect-route:inherited-text-style',
            2: 'unsupported-run-effect-route:vertical-text',
            3: 'unsupported-run-effect-route:relationship-bearing-text',
        }

        def assert_artifact(artifact) -> None:
            reason = expected_reasons[artifact.index]
            root = ET.fromstring(artifact.svg)
            self.assertIn('INHERITED', ''.join(root.itertext()))
            self.assertEqual(project_filter_errors(root), [])
            marked = [
                elem
                for elem in root.iter()
                if elem.get('data-pptx-effect-status') == 'unsupported'
            ]
            self.assertEqual(len(marked), 2)
            self.assertEqual(
                {
                    elem.get('data-pptx-effect-reason')
                    for elem in marked
                },
                {reason},
            )
            logical_shape = next(
                elem
                for elem in marked
                if elem.tag.rsplit('}', 1)[-1] == 'g'
            )
            self.assertEqual(
                logical_shape.get('data-name'),
                'Content Placeholder 2',
            )
            title_shape = next(
                elem
                for elem in root.iter()
                if elem.get('data-name') == 'Title 1'
            )
            self.assertIsNone(title_shape.get('data-pptx-effect-status'))
            status_errors = project_effect_status_errors(root)
            self.assertEqual(len(status_errors), 1)
            self.assertIn(reason, status_errors[0])

            metadata = [
                elem
                for elem in logical_shape.iter()
                if elem.get('data-pptx-part') == 'txbody'
            ]
            if artifact.index == 1:
                self.assertEqual(len(metadata), 1)
                payload = base64.b64decode(
                    (metadata[0].text or '').strip()
                ).decode('utf-8')
                self.assertNotIn('<a:effectLst', payload)
            else:
                self.assertEqual(metadata, [])

            self._assert_checker_and_exporter_reject(
                artifact.svg,
                reason,
                'unsupported imported PPTX effect',
            )

        for style_source in ('layout', 'master'):
            with self.subTest(style_source=style_source):
                with tempfile.TemporaryDirectory() as tmp_dir:
                    pptx_path = (
                        Path(tmp_dir) / f'{style_source}-run-effects.pptx'
                    )
                    self._write_inherited_run_effect_fixture(
                        pptx_path,
                        style_source=style_source,
                    )
                    imported = convert_pptx_to_svg(
                        pptx_path,
                        options=ConvertOptions(inheritance_mode='flat'),
                    )
                    self.assertEqual(len(imported.slides), 3)
                    for artifact in imported.slides:
                        with self.subTest(slide=artifact.index):
                            assert_artifact(artifact)

    def test_table_cell_run_effect_import_fails_closed(self):
        cases = [
            ('rPr:effectLst', False),
            ('defRPr', False),
            ('endParaRPr', False),
            ('rPr:empty-effectLst', False),
            ('rPr:empty-effectDag', False),
            ('none', False),
            ('rPr:effectLst', True),
        ]
        effect_slides = {1, 2, 3, 7}
        with tempfile.TemporaryDirectory() as tmp_dir:
            root_dir = Path(tmp_dir)
            pptx_path = root_dir / 'table-run-effects.pptx'
            self._write_table_run_effect_fixture(pptx_path, cases)
            imported = convert_pptx_to_svg(
                pptx_path,
                options=ConvertOptions(inheritance_mode='flat'),
            )
            self.assertEqual(len(imported.slides), len(cases))

            for artifact in imported.slides:
                with self.subTest(slide=artifact.index):
                    root = ET.fromstring(artifact.svg)
                    group = next(
                        elem
                        for elem in root.iter()
                        if (
                            elem.tag.rsplit('}', 1)[-1] == 'g'
                            and elem.get('data-name')
                            == f'TABLE EFFECT TEXT {artifact.index}'
                        )
                    )
                    self.assertIn(
                        f'CELL {artifact.index}',
                        ''.join(root.itertext()),
                    )
                    svg_path = root_dir / f'table-{artifact.index}.svg'
                    svg_path.write_text(artifact.svg, encoding='utf-8')

                    if artifact.index in effect_slides:
                        self.assertIsNone(
                            group.get('data-pptx-replace-with')
                        )
                        expected_status = (
                            'unsupported-native-transform'
                            if artifact.index == 7
                            else 'unsupported-table-direct-formatting'
                        )
                        self.assertEqual(
                            group.get('data-pptx-replacement-status'),
                            expected_status,
                        )
                        self.assertEqual(
                            group.get('data-pptx-effect-status'),
                            'unsupported',
                        )
                        self.assertEqual(
                            group.get('data-pptx-effect-reason'),
                            'unsupported-run-effect-route:table-cell-text',
                        )
                        self.assertFalse(any(
                            elem.tag.rsplit('}', 1)[-1] == 'metadata'
                            and elem.get('type') == 'application/json'
                            for elem in group
                        ))
                        status_errors = project_effect_status_errors(root)
                        self.assertEqual(len(status_errors), 1)
                        checker_result = SVGQualityChecker().check_file(
                            str(svg_path)
                        )
                        self.assertFalse(checker_result['passed'])
                        self.assertIn(
                            'unsupported-run-effect-route:table-cell-text',
                            '\n'.join(checker_result['errors']),
                        )
                        for native_objects in (False, True):
                            with self.assertRaisesRegex(
                                SvgNativeConversionError,
                                'unsupported imported PPTX effect',
                            ):
                                convert_svg_to_slide_shapes(
                                    svg_path,
                                    native_objects=native_objects,
                                )
                        continue

                    self.assertEqual(project_effect_status_errors(root), [])
                    self.assertEqual(
                        group.get('data-pptx-replace-with'),
                        'table',
                    )
                    self.assertIsNone(
                        group.get('data-pptx-replacement-status')
                    )
                    self.assertIsNone(group.get('data-pptx-effect-status'))
                    self.assertTrue(any(
                        elem.tag.rsplit('}', 1)[-1] == 'metadata'
                        and elem.get('type') == 'application/json'
                        for elem in group
                    ))
                    checker_result = SVGQualityChecker().check_file(
                        str(svg_path)
                    )
                    self.assertTrue(checker_result['passed'])
                    default_xml = convert_svg_to_slide_shapes(svg_path)[0]
                    native_xml = convert_svg_to_slide_shapes(
                        svg_path,
                        native_objects=True,
                    )[0]
                    self.assertNotIn('<a:tbl>', default_xml)
                    self.assertIn('<a:tbl>', native_xml)

    def test_picture_and_group_effect_import_is_explicit(self):
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        empty_sp_pr = ET.fromstring(
            f'<p:spPr xmlns:p="http://schemas.openxmlformats.org/'
            f'presentationml/2006/main" xmlns:a="{dml}">'
            '<a:effectLst/></p:spPr>'
        )
        self.assertEqual(
            unsupported_target_effect_metadata(empty_sp_pr, 'picture'),
            {},
        )
        nested_groups = ET.fromstring('''<svg>
  <g data-pptx-object="group" data-pptx-shape-id="10"
     data-pptx-shape-scope="slide" data-pptx-effect-status="unsupported"
     data-pptx-effect-reason="unsupported-effect-target:group:reflection">
    <g data-pptx-object="group" data-pptx-shape-id="11"
       data-pptx-shape-scope="slide" data-pptx-effect-status="unsupported"
       data-pptx-effect-reason="unsupported-effect-target:group:reflection"/>
  </g>
</svg>''')
        self.assertEqual(len(project_effect_status_errors(nested_groups)), 2)

        with tempfile.TemporaryDirectory() as tmp_dir:
            root_dir = Path(tmp_dir)
            pptx_path = root_dir / 'target-effects.pptx'
            self._write_picture_group_effect_fixture(pptx_path)
            imported = convert_pptx_to_svg(
                pptx_path,
                options=ConvertOptions(
                    inheritance_mode='flat',
                    embed_images=True,
                ),
            )
            self.assertEqual(len(imported.slides), 1)
            svg = imported.slides[0].svg
            root = ET.fromstring(svg)
            status_errors = project_effect_status_errors(root)
            self.assertEqual(len(status_errors), 2)
            self.assertTrue(any(
                'unsupported-effect-target:picture:outerShdw' in error
                for error in status_errors
            ))
            self.assertTrue(any(
                'unsupported-effect-target:group:reflection' in error
                for error in status_errors
            ))
            marked = [
                elem
                for elem in root.iter()
                if elem.get('data-pptx-effect-status') == 'unsupported'
            ]
            self.assertEqual(len(marked), 3)
            self.assertEqual(
                sum(elem.get('data-pptx-object') == 'picture' for elem in marked),
                2,
            )
            self.assertEqual(
                sum(elem.get('data-pptx-object') == 'group' for elem in marked),
                1,
            )
            self.assertNotIn('<filter', svg)
            self.assertTrue(any(
                elem.tag.endswith('svg')
                and elem.get('data-pptx-object') == 'picture'
                for elem in root.iter()
            ))
            self.assertTrue(any(
                elem.tag.endswith('image')
                for elem in root.iter()
            ))
            self.assertTrue(any(
                elem.get('data-pptx-object') == 'shape'
                for elem in root.iter()
            ))

            svg_path = root_dir / 'target-effects.svg'
            svg_path.write_text(svg, encoding='utf-8')
            checker_result = SVGQualityChecker().check_file(str(svg_path))
            self.assertFalse(checker_result['passed'])
            self.assertEqual(len(checker_result['errors']), 2)
            self.assertEqual(
                sum(
                    'unsupported source PPTX effect' in error
                    for error in checker_result['errors']
                ),
                2,
            )
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'unsupported imported PPTX effect',
            ):
                convert_svg_to_slide_shapes(svg_path)

    def test_missing_paint_reference_blocks_checker_and_exporter(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 1280 720">
  <rect x="80" y="80" width="300" height="180"
        fill="url(#missing)"/>
</svg>''',
            'has no matching direct <defs> definition',
            'invalid project paint reference',
        )

    def test_chart_table_replacement_attributes_resolve_canonical_and_legacy(self):
        import xml.etree.ElementTree as ET

        canonical = ET.fromstring(
            '<g data-pptx-replace-with="table" '
            'data-pptx-replacement-status="unsupported-table-style" '
            'data-pptx-import-source="pptx" '
            'data-pptx-fallback-kind="normalized"/>'
        )
        legacy = ET.fromstring(
            '<g data-pptx-native="table" '
            'data-pptx-native-status="unsupported-table-style" '
            'data-pptx-native-source="pptx" '
            'data-pptx-visual-status="normalized"/>'
        )
        for elem in (canonical, legacy):
            self.assertEqual(native_replacement_kind(elem), 'table')
            self.assertEqual(
                native_replacement_status(elem),
                'unsupported-table-style',
            )
            self.assertEqual(native_import_source(elem), 'pptx')
            self.assertEqual(native_fallback_kind(elem), 'normalized')

    def test_conflicting_chart_table_replacement_aliases_are_errors(self):
        import xml.etree.ElementTree as ET

        elem = ET.fromstring(
            '<g data-pptx-replace-with="chart" data-pptx-native="table"/>'
        )
        with self.assertRaisesRegex(
            NativeMarkerAttributeError,
            'data-pptx-replace-with.*conflicts',
        ):
            native_replacement_kind(elem)

        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="conflict" data-pptx-replace-with="chart" data-pptx-native="table"/>
</svg>'''
        )
        self.assertFalse(result['passed'])
        self.assertIn(
            'data-pptx-replace-with',
            '\n'.join(result['errors']),
        )
        self.assertNotIn(
            'legacy attribute data-pptx-native',
            '\n'.join(result['warnings']),
        )
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'conflict.svg'
            svg_path.write_text(
                '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="conflict" data-pptx-replace-with="chart" data-pptx-native="table"/>
</svg>''',
                encoding='utf-8',
            )
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'invalid chart/table replacement metadata',
            ):
                convert_svg_to_slide_shapes(svg_path)

    def test_all_chart_table_replacement_alias_pairs_reject_conflicts(self):
        import xml.etree.ElementTree as ET

        cases = (
            (
                'data-pptx-replace-with="chart" data-pptx-native="table"',
                native_replacement_kind,
            ),
            (
                'data-pptx-replacement-status="one" data-pptx-native-status="two"',
                native_replacement_status,
            ),
            (
                'data-pptx-import-source="pptx" data-pptx-native-source="other"',
                native_import_source,
            ),
            (
                'data-pptx-fallback-kind="normalized" '
                'data-pptx-visual-status="source-preview"',
                native_fallback_kind,
            ),
        )
        for attributes, getter in cases:
            with self.subTest(attributes=attributes):
                elem = ET.fromstring(f'<g {attributes}/>')
                with self.assertRaises(NativeMarkerAttributeError):
                    getter(elem)

    def test_matching_chart_table_replacement_aliases_are_accepted(self):
        import xml.etree.ElementTree as ET

        cases = (
            (
                '<g data-pptx-replace-with="chart" data-pptx-native="chart"/>',
                native_replacement_kind,
                'chart',
            ),
            (
                '<g data-pptx-replacement-status="reason" '
                'data-pptx-native-status="reason"/>',
                native_replacement_status,
                'reason',
            ),
            (
                '<g data-pptx-import-source="pptx" '
                'data-pptx-native-source="pptx"/>',
                native_import_source,
                'pptx',
            ),
            (
                '<g data-pptx-fallback-kind="normalized" '
                'data-pptx-visual-status="normalized"/>',
                native_fallback_kind,
                'normalized',
            ),
        )
        for markup, getter, expected in cases:
            with self.subTest(markup=markup):
                self.assertEqual(getter(ET.fromstring(markup)), expected)

    def test_invalid_replacement_tokens_block_checker_and_default_export(self):
        cases = (
            ('data-pptx-replace-with="diagram"', 'unsupported data-pptx-replace-with'),
            ('data-pptx-replace-with="Chart"', 'must use lowercase chart or table'),
            ('data-pptx-replace-with=" chart "', 'surrounding whitespace'),
            ('data-pptx-replacement-status=""', 'must not be empty'),
            ('data-pptx-replacement-status=" reason "', 'surrounding whitespace'),
            ('data-pptx-import-source="other"', 'unsupported data-pptx-import-source'),
            ('data-pptx-fallback-kind="Normalized"', 'unsupported data-pptx-fallback-kind'),
        )
        for attributes, expected in cases:
            content = (
                '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
                f'<g id="marker" {attributes}/></svg>'
            )
            with self.subTest(attributes=attributes):
                result = self._check(content)
                self.assertFalse(result['passed'])
                self.assertIn(expected, '\n'.join(result['errors']))
                with tempfile.TemporaryDirectory() as tmp_dir:
                    svg_path = Path(tmp_dir) / 'invalid-token.svg'
                    svg_path.write_text(content, encoding='utf-8')
                    with self.assertRaisesRegex(
                        SvgNativeConversionError,
                        'invalid chart/table replacement metadata',
                    ):
                        convert_svg_to_slide_shapes(svg_path)

    def test_replacement_status_uses_closed_importer_reason_codes(self):
        legal = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="fallback" data-pptx-replacement-status="unsupported-table-style">
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''
        made_up = legal.replace('unsupported-table-style', 'made-up-reason')

        legal_result = self._check(legal)
        self.assertTrue(legal_result['passed'])
        self.assertEqual(legal_result['errors'], [])

        made_up_result = self._check(made_up)
        self.assertFalse(made_up_result['passed'])
        self.assertIn(
            "unsupported data-pptx-replacement-status value: 'made-up-reason'",
            '\n'.join(made_up_result['errors']),
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            legal_path = Path(tmp_dir) / 'legal.svg'
            made_up_path = Path(tmp_dir) / 'made-up.svg'
            legal_path.write_text(legal, encoding='utf-8')
            made_up_path.write_text(made_up, encoding='utf-8')
            convert_svg_to_slide_shapes(legal_path)
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'unsupported data-pptx-replacement-status value',
            ):
                convert_svg_to_slide_shapes(made_up_path)

    def test_legacy_replacement_attributes_remain_non_blocking(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="legacy-table" data-pptx-native="table"
     data-pptx-native-source="pptx" data-pptx-visual-status="normalized"
     data-pptx-x="80" data-pptx-y="80"
     data-pptx-width="320" data-pptx-height="180">
    <metadata data-pptx-kind="table">{"rows":[["A"]]}</metadata>
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''
        )
        warning_text = '\n'.join(result['warnings'])
        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])
        self.assertIn('legacy attribute data-pptx-native', warning_text)
        self.assertIn('legacy metadata attribute data-pptx-kind', warning_text)

    def test_canonical_placeholder_fallback_does_not_require_route_status(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="placeholder" data-pptx-fallback-kind="placeholder">
    <rect x="80" y="80" width="320" height="180" fill="#EEEEEE"/>
  </g>
</svg>'''
        )
        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])
        self.assertIn(
            'reconstruction-only placeholder',
            '\n'.join(result['warnings']),
        )

    def test_legacy_placeholder_fallback_still_requires_legacy_route_status(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="placeholder" data-pptx-visual-status="placeholder">
    <rect x="80" y="80" width="320" height="180" fill="#EEEEEE"/>
  </g>
</svg>'''
        )
        self.assertFalse(result['passed'])
        self.assertIn(
            "data-pptx-route-status='reconstruction-only'",
            '\n'.join(result['errors']),
        )

    def test_metadata_kind_must_match_parent_replacement_kind(self):
        content = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="chart" data-pptx-replace-with="chart"
     data-pptx-x="80" data-pptx-y="80"
     data-pptx-width="320" data-pptx-height="180">
    <metadata type="application/json" data-pptx-kind="table">{"rows":[["A"]]}</metadata>
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''
        result = self._check(content)
        self.assertFalse(result['passed'])
        self.assertIn(
            "metadata kind 'table' conflicts with parent replacement kind 'chart'",
            '\n'.join(result['errors']),
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'metadata-conflict.svg'
            svg_path.write_text(content, encoding='utf-8')
            for native_objects in (False, True):
                with self.subTest(native_objects=native_objects):
                    with self.assertRaisesRegex(
                        SvgNativeConversionError,
                        "metadata kind 'table' conflicts",
                    ):
                        convert_svg_to_slide_shapes(
                            svg_path,
                            native_objects=native_objects,
                        )

    def test_conflicting_legacy_metadata_kind_aliases_are_errors(self):
        content = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="chart" data-pptx-replace-with="chart"
     data-pptx-x="80" data-pptx-y="80"
     data-pptx-width="320" data-pptx-height="180">
    <metadata type="application/json" data-pptx-native="chart"
              data-pptx-kind="table">{"series":[]}</metadata>
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''
        result = self._check(content)
        self.assertFalse(result['passed'])
        self.assertIn(
            'metadata data-pptx-native conflicts with data-pptx-kind',
            '\n'.join(result['errors']),
        )
        self.assertNotIn(
            'legacy metadata attribute',
            '\n'.join(result['warnings']),
        )

    def test_structured_placeholder_wraps_replacement_alias_conflicts(self):
        import xml.etree.ElementTree as ET

        carrier = ET.fromstring(
            '<g data-pptx-replace-with="chart" data-pptx-native="table"/>'
        )
        with self.assertRaisesRegex(
            TemplateStructureError,
            'conflicting chart/table replacement metadata',
        ):
            _validate_placeholder_carrier(
                carrier,
                'chart',
                svg_path=Path('template.svg'),
                element_id='chart-slot',
            )

    def test_local_use_rejects_all_canonical_replacement_metadata(self):
        import xml.etree.ElementTree as ET

        attributes = (
            'data-pptx-replace-with="chart"',
            'data-pptx-replacement-status="reason"',
            'data-pptx-import-source="pptx"',
            'data-pptx-fallback-kind="normalized"',
            'data-pptx-fallback-sha256="' + ('0' * 64) + '"',
        )
        for attribute in attributes:
            root = ET.fromstring(
                '<svg xmlns="http://www.w3.org/2000/svg">'
                f'<defs><g id="source" {attribute}><rect width="10" height="10"/></g></defs>'
                '<use href="#source"/></svg>'
            )
            with self.subTest(attribute=attribute):
                with self.assertRaises(UseExpansionError):
                    expand_local_use_references(root)

    def test_canonical_and_legacy_table_markers_export_identically(self):
        payload = '{"x":80,"y":80,"width":320,"height":180,"rows":[["A"]]}'
        canonical = f'''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="table" data-pptx-replace-with="table">
    <metadata type="application/json">{payload}</metadata>
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''
        legacy = f'''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="table" data-pptx-native="table">
    <metadata data-pptx-native="table">{payload}</metadata>
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''

        with tempfile.TemporaryDirectory() as tmp_dir:
            canonical_path = Path(tmp_dir) / 'canonical.svg'
            legacy_path = Path(tmp_dir) / 'legacy.svg'
            canonical_path.write_text(canonical, encoding='utf-8')
            legacy_path.write_text(legacy, encoding='utf-8')
            canonical_result = convert_svg_to_slide_shapes(
                canonical_path,
                native_objects=True,
            )
            legacy_result = convert_svg_to_slide_shapes(
                legacy_path,
                native_objects=True,
            )

        self.assertEqual(canonical_result, legacy_result)

    def test_authored_hashless_native_marker_has_no_baseline_warning(self):
        content = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="table" data-pptx-replace-with="table">
    <metadata type="application/json">{"x":80,"y":80,"width":320,"height":180,"rows":[["A"]]}</metadata>
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''
        result = self._check(content)
        self.assertTrue(result['passed'])
        self.assertEqual(result['warnings'], [])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'authored.svg'
            svg_path.write_text(content, encoding='utf-8')
            stderr = io.StringIO()
            with patch('sys.stderr', new=stderr):
                convert_svg_to_slide_shapes(svg_path, native_objects=True)

        self.assertNotIn('fallback-sha256', stderr.getvalue())

    def test_imported_hashless_native_marker_warns_but_remains_compatible(self):
        for source_attribute in (
            'data-pptx-import-source="pptx"',
            'data-pptx-native-source="pptx"',
        ):
            content = f'''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="table" data-pptx-replace-with="table" {source_attribute}>
    <metadata type="application/json">{{"x":80,"y":80,"width":320,"height":180,"rows":[["A"]]}}</metadata>
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''
            with self.subTest(source_attribute=source_attribute):
                result = self._check(content)
                self.assertTrue(result['passed'])
                self.assertIn(
                    'imported marker has no data-pptx-fallback-sha256 baseline',
                    '\n'.join(result['warnings']),
                )

                with tempfile.TemporaryDirectory() as tmp_dir:
                    svg_path = Path(tmp_dir) / 'imported.svg'
                    svg_path.write_text(content, encoding='utf-8')
                    stderr = io.StringIO()
                    with patch('sys.stderr', new=stderr):
                        convert_svg_to_slide_shapes(
                            svg_path,
                            native_objects=True,
                        )
                self.assertIn(
                    'imported marker has no data-pptx-fallback-sha256 baseline',
                    stderr.getvalue(),
                )

    def test_explicit_authored_fallback_baseline_still_blocks_stale_native_export(self):
        root = ET.fromstring(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="table" data-pptx-replace-with="table">
    <metadata type="application/json">{"x":80,"y":80,"width":320,"height":180,"rows":[["A"]]}</metadata>
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''
        )
        marker = next(elem for elem in root if elem.tag.endswith('g'))
        stamp_native_fallback_baseline(marker, document_root=root)

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'sealed-authored.svg'
            ET.ElementTree(root).write(svg_path, encoding='unicode')
            fresh_result = SVGQualityChecker().check_file(str(svg_path))
            self.assertTrue(fresh_result['passed'])
            self.assertEqual(fresh_result['warnings'], [])
            convert_svg_to_slide_shapes(svg_path, native_objects=True)

            rect = next(elem for elem in marker if elem.tag.endswith('rect'))
            rect.set('fill', '#EEEEEE')
            ET.ElementTree(root).write(svg_path, encoding='unicode')
            stale_result = SVGQualityChecker().check_file(str(svg_path))
            self.assertTrue(stale_result['passed'])
            self.assertIn(
                'visible SVG fallback differs from its recorded baseline',
                '\n'.join(stale_result['warnings']),
            )
            convert_svg_to_slide_shapes(svg_path)
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'fallback was edited after its baseline',
            ):
                convert_svg_to_slide_shapes(svg_path, native_objects=True)

    def test_explicit_authored_invalid_fallback_baseline_still_fails_native_export(self):
        content = '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">
  <g id="table" data-pptx-replace-with="table"
     data-pptx-fallback-sha256="invalid">
    <metadata type="application/json">{"x":80,"y":80,"width":320,"height":180,"rows":[["A"]]}</metadata>
    <rect x="80" y="80" width="320" height="180" fill="#FFFFFF"/>
  </g>
</svg>'''
        result = self._check(content)
        self.assertTrue(result['passed'])
        self.assertIn(
            'must be a 64-digit SHA-256',
            '\n'.join(result['warnings']),
        )
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'invalid-baseline.svg'
            svg_path.write_text(content, encoding='utf-8')
            convert_svg_to_slide_shapes(svg_path)
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'must be a 64-digit SHA-256',
            ):
                convert_svg_to_slide_shapes(svg_path, native_objects=True)

    def test_authored_template_markers_do_not_preseed_static_fallback_hashes(self):
        template_root = SCRIPT_DIR.parent / 'templates'
        offenders = []
        for svg_path in template_root.rglob('*.svg'):
            root = ET.parse(svg_path).getroot()
            for elem in root.iter():
                if not native_replacement_kind(elem):
                    continue
                if native_import_source(elem) == 'pptx':
                    continue
                if elem.get('data-pptx-fallback-sha256') is not None:
                    offenders.append(str(svg_path.relative_to(template_root)))
        self.assertEqual(offenders, [])

    def test_chart_templates_name_all_top_level_groups(self):
        chart_root = SCRIPT_DIR.parent / 'templates' / 'charts'
        anonymous_groups = []
        duplicate_ids = []
        for svg_path in chart_root.glob('*.svg'):
            root = ET.parse(svg_path).getroot()
            seen_ids = set()
            for elem in root.iter():
                element_id = elem.get('id')
                if not element_id:
                    continue
                if element_id in seen_ids:
                    duplicate_ids.append(
                        f'{svg_path.name}: {element_id}'
                    )
                seen_ids.add(element_id)
            for index, child in enumerate(list(root), start=1):
                if child.tag.rsplit('}', 1)[-1] != 'g':
                    continue
                if (child.get('id') or '').strip():
                    continue
                anonymous_groups.append(
                    f'{svg_path.name}: root child {index}'
                )
        self.assertEqual(anonymous_groups, [])
        self.assertEqual(duplicate_ids, [])

    def test_chart_templates_use_canonical_readable_font_sizes(self):
        chart_root = SCRIPT_DIR.parent / 'templates' / 'charts'
        noncanonical_sizes = []
        undersized_text = []
        for svg_path in chart_root.glob('*.svg'):
            root = ET.parse(svg_path).getroot()
            for elem in root.iter():
                styles = parse_inline_style(elem.get('style'))
                raw_size = styles.get('font-size', elem.get('font-size'))
                if raw_size is None:
                    continue
                try:
                    size = float(raw_size)
                except ValueError:
                    noncanonical_sizes.append(
                        f'{svg_path.name}: {raw_size}'
                    )
                    continue
                if not math.isfinite(size):
                    noncanonical_sizes.append(
                        f'{svg_path.name}: {raw_size}'
                    )
                    continue
                if size < 12:
                    element_id = elem.get('id') or elem.tag.rsplit('}', 1)[-1]
                    undersized_text.append(
                        f'{svg_path.name}: {element_id}={raw_size}'
                    )
        self.assertEqual(noncanonical_sizes, [])
        self.assertEqual(undersized_text, [])

    def test_chart_text_bound_warnings_match_real_overflows(self):
        chart_root = SCRIPT_DIR.parent / 'templates' / 'charts'
        warned = set()
        for svg_path in chart_root.glob('*.svg'):
            result = SVGQualityChecker().check_file(str(svg_path))
            if self._text_bound_warnings(result):
                warned.add(svg_path.name)
        self.assertEqual(
            warned,
            {'bullet_chart.svg', 'pyramid_isometric.svg'},
        )

    def test_chart_catalog_passes_checker_and_export_routes(self):
        chart_root = SCRIPT_DIR.parent / 'templates' / 'charts'
        index = json.loads(
            (chart_root / 'charts_index.json').read_text(encoding='utf-8')
        )
        indexed_keys = set(index['charts'])
        svg_paths = sorted(chart_root.glob('*.svg'))
        svg_keys = {path.stem for path in svg_paths}
        self.assertEqual(index['meta']['total'], len(indexed_keys))
        self.assertEqual(svg_keys, indexed_keys)

        for svg_path in svg_paths:
            with self.subTest(chart=svg_path.stem):
                checker_result = SVGQualityChecker().check_file(str(svg_path))
                self.assertEqual(checker_result['errors'], [])

                default_result = convert_svg_to_slide_shapes(svg_path)
                default_xml = default_result[0]
                self.assertNotIn('<c:chart ', default_xml)
                self.assertNotIn('<cx:chart ', default_xml)
                self.assertNotIn('<a:tbl>', default_xml)

                root = ET.parse(svg_path).getroot()
                replacement_kinds = {
                    kind
                    for elem in root.iter()
                    if (kind := native_replacement_kind(elem))
                }
                stderr = io.StringIO()
                with patch('sys.stderr', new=stderr):
                    native_result = convert_svg_to_slide_shapes(
                        svg_path,
                        native_objects=True,
                    )
                native_xml = native_result[0]

                if not replacement_kinds:
                    self.assertEqual(native_result, default_result)
                    continue

                self.assertEqual(len(replacement_kinds), 1)
                self.assertNotEqual(native_result, default_result)
                if replacement_kinds == {'chart'}:
                    self.assertTrue(
                        '<c:chart ' in native_xml
                        or '<cx:chart ' in native_xml
                    )
                    self.assertNotIn('<a:tbl>', native_xml)
                elif replacement_kinds == {'table'}:
                    self.assertIn('<a:tbl>', native_xml)
                    self.assertNotIn('<c:chart ', native_xml)
                    self.assertNotIn('<cx:chart ', native_xml)
                else:
                    self.fail(
                        f'unsupported replacement kinds: {replacement_kinds}'
                    )

    def test_compact_authored_preset_exports_one_native_shape(self):
        fragment = render_preset_shape_fragment(
            'rightArrow',
            (10, 20, 100, 40),
            adjustments={'adj1': 50000, 'adj2': 50000},
            element_id='next-step',
            style={'fill': '#2563EB', 'stroke': 'none'},
        )
        self.assertNotIn('data-pptx-part', fragment)
        self.assertNotIn('data-pptx-preview-sha256', fragment)
        self.assertNotIn('visibility="hidden"', fragment)
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 120">'
            f'{fragment}</svg>'
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'compact.svg'
            svg_path.write_text(source, encoding='utf-8')
            result = SVGQualityChecker().check_file(str(svg_path))
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])
        self.assertEqual(result['warnings'], [])
        self.assertEqual(slide_xml.count('<p:sp>'), 1)
        self.assertEqual(
            slide_xml.count('<a:prstGeom prst="rightArrow">'),
            1,
        )
        self.assertNotIn('<p:cxnSp>', slide_xml)

    def test_compact_multi_path_preset_exports_one_native_shape(self):
        fragment = render_preset_shape_fragment(
            'cube',
            (10, 20, 100, 80),
            element_id='cube-node',
            style={'fill': '#2563EB', 'stroke': 'none'},
        )
        group = ET.fromstring(fragment)
        self.assertEqual(len(group), 4)
        self.assertTrue(all(child.tag == 'path' for child in group))
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 140">'
            f'{fragment}</svg>'
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'compact-multi.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertEqual(slide_xml.count('<p:sp>'), 1)
        self.assertEqual(slide_xml.count('<a:prstGeom prst="cube">'), 1)

    def test_every_registered_preset_supports_compact_authoring(self):
        registry = get_preset_registry()
        for index, preset in enumerate(registry.names):
            connector = preset in CONNECTOR_PRESET_TYPES
            style = (
                {
                    'fill': 'none',
                    'stroke': '#334155',
                    'stroke-width': '2',
                }
                if connector
                else {'fill': '#2563EB', 'stroke': 'none'}
            )
            with self.subTest(preset=preset):
                fragment = render_preset_shape_fragment(
                    preset,
                    (10, 20, 100, 60),
                    object_kind='connector' if connector else 'shape',
                    element_id=f'preset-{index}',
                    style=style,
                )
                root = ET.fromstring(
                    '<svg xmlns="http://www.w3.org/2000/svg">'
                    f'{fragment}</svg>'
                )
                self.assertEqual(validate_authored_preset_tree(root), [])
                self.assertEqual(
                    materialize_compact_authored_preset_tree(root),
                    1,
                )
                self.assertEqual(validate_authored_preset_tree(root), [])

    def test_compact_authored_connector_exports_one_native_connector(self):
        fragment = render_preset_shape_fragment(
            'bentConnector3',
            (10, 20, 100, 80),
            object_kind='connector',
            element_id='relationship',
            style={
                'fill': 'none',
                'stroke': '#334155',
                'stroke-width': '2',
            },
        )
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 140">'
            f'{fragment}</svg>'
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'compact-connector.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertEqual(slide_xml.count('<p:cxnSp>'), 1)
        self.assertEqual(
            slide_xml.count('<a:prstGeom prst="bentConnector3">'),
            1,
        )
        self.assertNotIn('<p:sp>', slide_xml)

    def test_compact_authored_preset_rejects_stale_visible_path(self):
        fragment = render_preset_shape_fragment(
            'diamond',
            (20, 20, 80, 60),
            element_id='decision',
            style={'fill': '#F59E0B', 'stroke': 'none'},
        )
        stale = fragment.replace('M 20 50', 'M 21 50', 1)
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 160 120">'
            f'{stale}</svg>'
        )
        self._assert_checker_and_exporter_reject(
            source,
            'Compact authored preset path 1 differs from registry output',
            'Invalid authored preset structure',
        )

    def test_compact_authored_preset_rejects_noncanonical_frame_spelling(self):
        fragment = render_preset_shape_fragment(
            'diamond',
            (20, 20, 80, 60),
            element_id='decision',
            style={'fill': '#F59E0B', 'stroke': 'none'},
        ).replace(
            'data-pptx-frame="20 20 80 60"',
            'data-pptx-frame="2e1,20,+80,60.0"',
            1,
        )
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 160 120">'
            f'{fragment}</svg>'
        )

        self._assert_checker_and_exporter_reject(
            source,
            'data-pptx-frame must use the helper',
            'Invalid authored preset structure',
        )

    def test_compact_preset_ancestor_paint_is_compatible_with_warning(self):
        fragment = render_preset_shape_fragment(
            'diamond',
            (20, 20, 80, 60),
            element_id='decision',
            style={'fill': '#F59E0B', 'stroke': 'none'},
        )
        result = self._check(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 160 120">'
            '<g id="faded" fill-opacity="0.25">'
            f'{fragment}</g></svg>'
        )

        self.assertTrue(result['passed'])
        self.assertIn(
            'Compact authored preset(s) use compatible ancestor paint',
            '\n'.join(result['warnings']),
        )

    def test_registry_derived_preset_layers_do_not_drift_from_spec_lock(self):
        fragment = render_preset_shape_fragment(
            'cube',
            (20, 20, 80, 60),
            element_id='cube',
            style={'fill': '#2563EB', 'stroke': 'none'},
        )
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 160 120">'
            f'{fragment}</svg>'
        )

        result = self._check_with_spec_lock(
            source,
            (
                '# Execution Lock\n\n'
                '## colors\n'
                '- primary: #2563EB\n'
            ),
        )

        self.assertTrue(result['passed'])
        self.assertNotIn('spec_lock drift', '\n'.join(result['warnings']))

    def test_spec_lock_color_fallback_normalizes_every_declared_color(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 160 120">
  <rect x="10" y="10" width="60" height="40" fill="#2563EB"/>
  <rect x="90" y="10" width="60" height="40" fill="#10B981"/>
</svg>'''
        lock_text = (
            '# Execution Lock\n\n'
            '## colors\n'
            '- primary: #2563EB\n'
            '- success: #10B981\n'
        )

        with patch('svg_quality_checker._parse_export_color', None):
            result = self._check_with_spec_lock(source, lock_text)

        self.assertTrue(result['passed'])
        self.assertNotIn('spec_lock drift', '\n'.join(result['warnings']))

    def test_spec_lock_color_fallback_accepts_empty_color_section(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 160 120">
  <rect x="10" y="10" width="60" height="40" fill="#FFFFFF"/>
</svg>'''
        lock_text = '# Execution Lock\n\n## colors\n'

        with patch('svg_quality_checker._parse_export_color', None):
            result = self._check_with_spec_lock(source, lock_text)

        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])
        self.assertIn('spec_lock drift', '\n'.join(result['warnings']))

    def test_foreign_namespace_cannot_impersonate_compact_svg_path(self):
        fragment = render_preset_shape_fragment(
            'diamond',
            (20, 20, 80, 60),
            element_id='decision',
            style={'fill': '#F59E0B', 'stroke': 'none'},
        )
        root = ET.fromstring(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 160 120">'
            f'{fragment}</svg>'
        )
        group = list(root)[0]
        list(group)[0].tag = '{urn:foreign}path'
        source = ET.tostring(root, encoding='unicode')

        self._assert_checker_and_exporter_reject(
            source,
            'may contain only direct SVG paths',
            'Invalid authored preset structure',
        )

    def test_expanded_authored_preset_remains_compatible(self):
        fragment = render_preset_shape_fragment(
            'triangle',
            (20, 20, 80, 60),
            element_id='triangle-node',
            style={'fill': '#7C3AED', 'stroke': 'none'},
        )
        root = ET.fromstring(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 160 120">'
            f'{fragment}</svg>'
        )
        self.assertEqual(materialize_compact_authored_preset_tree(root), 1)

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'expanded.svg'
            ET.ElementTree(root).write(
                svg_path,
                encoding='utf-8',
                xml_declaration=True,
            )
            result = SVGQualityChecker().check_file(str(svg_path))
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertTrue(result['passed'])
        self.assertIn(
            'Compatible expanded authored-preset fragment(s) detected',
            '\n'.join(result['warnings']),
        )
        self.assertEqual(slide_xml.count('<a:prstGeom prst="triangle">'), 1)

    def test_malformed_expanded_preset_does_not_receive_migration_warning(self):
        fragment = render_preset_shape_fragment(
            'triangle',
            (20, 20, 80, 60),
            element_id='triangle-node',
            style={'fill': '#7C3AED', 'stroke': 'none'},
        )
        root = ET.fromstring(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 160 120">'
            f'{fragment}</svg>'
        )
        self.assertEqual(materialize_compact_authored_preset_tree(root), 1)
        group = list(root)[0]
        preview = next(
            child
            for child in group
            if child.get('data-pptx-part') == 'geometry-preview'
        )
        group.remove(preview)

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'malformed-expanded.svg'
            ET.ElementTree(root).write(
                svg_path,
                encoding='utf-8',
                xml_declaration=True,
            )
            result = SVGQualityChecker().check_file(str(svg_path))

        self.assertFalse(result['passed'])
        self.assertNotIn(
            'Compatible expanded authored-preset fragment(s) detected',
            '\n'.join(result['warnings']),
        )

    def test_expanded_authored_preset_rejects_compact_only_group_style(self):
        fragment = render_preset_shape_fragment(
            'triangle',
            (20, 20, 80, 60),
            element_id='triangle-node',
            style={'fill': '#7C3AED', 'stroke': 'none'},
        )
        root = ET.fromstring(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 160 120">'
            f'{fragment}</svg>'
        )
        self.assertEqual(materialize_compact_authored_preset_tree(root), 1)
        list(root)[0].set('fill-opacity', '0.5')
        source = ET.tostring(root, encoding='unicode')

        self._assert_checker_and_exporter_reject(
            source,
            'unsupported attributes: fill-opacity',
            'Invalid authored preset structure',
        )

    def test_compact_authored_preset_is_a_structured_fixed_atom(self):
        fragment = render_preset_shape_fragment(
            'rightArrow',
            (20, 20, 80, 40),
            element_id='master-arrow',
            style={'fill': '#2563EB', 'stroke': 'none'},
        ).replace(
            ' fill="#2563EB"',
            ' data-pptx-layer="master" data-pptx-editable="false" '
            'fill="#2563EB"',
            1,
        )
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720" '
            'data-pptx-master="master" data-pptx-master-name="Master" '
            'data-pptx-layout="layout" data-pptx-layout-name="Layout">'
            f'{fragment}</svg>'
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / '01_cover.svg'
            svg_path.write_text(source, encoding='utf-8')
            spec = parse_template_slide(svg_path, 1)
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertEqual(len(spec.elements), 1)
        self.assertEqual(spec.elements[0].element_id, 'master-arrow')
        self.assertEqual(spec.elements[0].layer, 'master')
        self.assertEqual(
            slide_xml.count('<a:prstGeom prst="rightArrow">'),
            1,
        )

    def test_native_placeholder_match_uses_effective_zero_index(self):
        def slide_spec(idx: int) -> TemplateSlideSpec:
            element = TemplateElementSpec(
                element_id='title-slot',
                order=1,
                tag='g',
                placeholder='title',
                placeholder_idx=idx,
            )
            return TemplateSlideSpec(
                slide_num=1,
                svg_path=Path('01_cover.svg'),
                master_key='master',
                master_name='Master',
                layout_key='layout',
                layout_name='Layout',
                elements=(element,),
            )

        source_placeholder = NativePlaceholderSpec(
            semantic_role='title',
            placeholder_type='title',
            idx=None,
        )
        layout = NativeLayoutSpec(
            key='layout',
            name='Layout',
            package_part='ppt/slideLayouts/slideLayout1.xml',
            master_key='master',
            placeholders=(source_placeholder,),
        )

        matches = match_native_placeholders(slide_spec(0), layout)
        self.assertEqual(matches[0][1], source_placeholder)
        self.assertIsNone(source_placeholder.idx)
        self.assertEqual(source_placeholder.effective_idx, 0)
        with self.assertRaisesRegex(
            TemplateStructureError,
            'title idx=1.*has no compatible source placeholder',
        ):
            match_native_placeholders(slide_spec(1), layout)

    def test_stale_compact_preset_is_not_a_structured_atom_or_carrier(self):
        fragment = render_preset_shape_fragment(
            'diamond',
            (20, 20, 80, 60),
            element_id='decision',
            style={'fill': '#F59E0B', 'stroke': 'none'},
        ).replace('M 20 50', 'M 21 50', 1)
        carrier = ET.fromstring(fragment)
        carrier.set('data-pptx-placeholder-carrier', 'true')
        with self.assertRaisesRegex(
            TemplateStructureError,
            'object placeholder carrier must be',
        ):
            _validate_placeholder_carrier(
                carrier,
                'object',
                svg_path=Path('template.svg'),
                element_id='object-slot',
            )

        carrier.attrib.pop('data-pptx-placeholder-carrier')
        carrier.set('data-pptx-layer', 'master')
        carrier.set('data-pptx-editable', 'false')
        root = ET.fromstring(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720" '
            'data-pptx-master="master" data-pptx-master-name="Master" '
            'data-pptx-layout="layout" data-pptx-layout-name="Layout"/>'
        )
        root.append(carrier)
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / '01_cover.svg'
            ET.ElementTree(root).write(
                svg_path,
                encoding='utf-8',
                xml_declaration=True,
            )
            with self.assertRaisesRegex(
                TemplateStructureError,
                'is a <g> on the master layer',
            ):
                parse_template_slide(svg_path, 1)

    def test_expanded_authored_preset_does_not_gain_structured_atom_status(self):
        fragment = render_preset_shape_fragment(
            'rightArrow',
            (20, 20, 80, 40),
            element_id='master-arrow',
            style={'fill': '#2563EB', 'stroke': 'none'},
        )
        root = ET.fromstring(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720" '
            'data-pptx-master="master" data-pptx-master-name="Master" '
            'data-pptx-layout="layout" data-pptx-layout-name="Layout">'
            f'{fragment}</svg>'
        )
        self.assertEqual(materialize_compact_authored_preset_tree(root), 1)
        group = list(root)[0]
        group.set('data-pptx-layer', 'master')
        group.set('data-pptx-editable', 'false')

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / '01_cover.svg'
            ET.ElementTree(root).write(
                svg_path,
                encoding='utf-8',
                xml_declaration=True,
            )
            with self.assertRaisesRegex(
                TemplateStructureError,
                'is a <g> on the master layer',
            ):
                parse_template_slide(svg_path, 1)

    def test_compact_authored_preset_is_an_object_slot_carrier(self):
        carrier = ET.fromstring(render_preset_shape_fragment(
            'diamond',
            (20, 20, 80, 60),
            element_id='object-preset',
            style={'fill': '#F59E0B', 'stroke': 'none'},
        ))
        carrier.set('data-pptx-placeholder-carrier', 'true')

        _validate_placeholder_carrier(
            carrier,
            'object',
            svg_path=Path('template.svg'),
            element_id='object-slot',
        )

        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 160 120">'
            f'{ET.tostring(carrier, encoding="unicode")}</svg>'
        )
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'object-carrier.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertEqual(
            slide_xml.count('<a:prstGeom prst="diamond">'),
            1,
        )

    def test_local_use_rejects_authored_preset_metadata(self):
        attributes = (
            'data-pptx-authoring="preset"',
            'data-pptx-object="shape"',
            'data-pptx-prst="diamond"',
            'data-pptx-frame="20 20 80 60"',
            'data-pptx-av-adj="val 50000"',
        )
        for attribute in attributes:
            root = ET.fromstring(
                '<svg xmlns="http://www.w3.org/2000/svg">'
                f'<defs><g id="source" {attribute}>'
                '<path d="M 0 5 L 5 0 L 10 5 L 5 10 Z"/>'
                '</g></defs><use href="#source"/></svg>'
            )
            with self.subTest(attribute=attribute):
                with self.assertRaises(UseExpansionError):
                    expand_local_use_references(root)

    def test_single_line_text_warning_uses_rendered_bounds(self):
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 100">'
            '<text x="10" y="30" font-size="20">\n'
            f'  {" " * 120}Fit\n'
            '</text>'
            '<text x="150" y="70" font-size="20">WWWW</text>'
            '</svg>'
        )
        result = self._check(source)
        bounds_warnings = self._text_bound_warnings(result)
        self.assertEqual(len(bounds_warnings), 1)
        self.assertIn("'WWWW'", bounds_warnings[0])
        self.assertNotIn("'Fit'", bounds_warnings[0])

    def test_single_line_text_warning_allows_one_pixel_tolerance(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 200 100" font-size="20">
  <text x="141" y="30">WWWW</text>
  <text x="142" y="70">MMMM</text>
</svg>'''
        )
        bounds_warnings = self._text_bound_warnings(result)
        self.assertEqual(len(bounds_warnings), 1)
        self.assertIn("'MMMM'", bounds_warnings[0])
        self.assertNotIn("'WWWW'", bounds_warnings[0])

    def test_single_line_text_warning_respects_text_anchor(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 200 100" font-size="20">
  <text x="100" y="20" text-anchor="middle">WWWWWWWW</text>
  <text x="40" y="20" text-anchor="middle">WWWWWWWW</text>
  <text x="160" y="20" text-anchor="middle">WWWWWWWW</text>
  <text x="200" y="50" text-anchor="end">WWWWWWWW</text>
  <text x="100" y="50" text-anchor="end">WWWWWWWW</text>
  <text x="210" y="50" text-anchor="end">WWWWWWWW</text>
</svg>'''
        result = self._check(source)
        bounds_warnings = self._text_bound_warnings(result)
        self.assertEqual(len(bounds_warnings), 1)
        self.assertIn('Detected 4 single-line text element(s)', bounds_warnings[0])

    def test_single_line_text_warning_resolves_inline_runs_and_transforms(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 200 100" font-size="20">
  <defs><g id="hidden"><text x="190" y="20">WWWW</text></g></defs>
  <g font-weight="semibold">
    <text x="150" y="20">WW<tspan>WW</tspan></text>
  </g>
  <text x="150" y="30" letter-spacing="50">A<tspan
        fill="#FF0000">B</tspan></text>
  <text x="50" y="40"><tspan x="10">WWWWWW</tspan>
    <tspan x="10" dy="20">WWWWWW</tspan></text>
  <g transform="translate(-200 0)"><text x="250" y="60">WWWW</text></g>
  <g transform="translate(100 0)"><text x="50" y="80">WWWW</text></g>
  <g transform="translate(100 0)"><g transform="rotate(180 0 45)">
    <text x="10" y="45" font-size="10">ii</text>
  </g></g>
  <text x="190" y="90">{{LONG_PLACEHOLDER}}</text>
</svg>'''
        result = self._check(source)
        bounds_warnings = self._text_bound_warnings(result)
        self.assertEqual(len(bounds_warnings), 1)
        self.assertIn('Detected 2 single-line text element(s)', bounds_warnings[0])

    def test_single_line_text_warning_includes_inherited_tracking(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 200 100">
  <g font-size="20" font-weight="semibold" letter-spacing="10">
    <text x="100" y="50">iiiiiiii</text>
  </g>
</svg>'''
        )
        bounds_warnings = self._text_bound_warnings(result)
        self.assertEqual(len(bounds_warnings), 1)
        self.assertIn("'iiiiiiii'", bounds_warnings[0])

    def test_single_line_text_warning_measures_negative_tracking_bounds(self):
        result = self._check(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 200 100">
  <text x="5" y="25" font-size="20" letter-spacing="-20">AB</text>
  <text x="195" y="50" font-size="20" letter-spacing="-20">AB</text>
</svg>'''
        )
        bounds_warnings = self._text_bound_warnings(result)
        self.assertEqual(len(bounds_warnings), 1)
        self.assertIn('Detected 2 single-line text element(s)', bounds_warnings[0])

    def test_non_positive_text_frame_extents_block_checker_and_exporter(self):
        for anchor in ('start', 'middle', 'end'):
            source = (
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<text x="320" y="100" font-size="20" '
                f'text-anchor="{anchor}" letter-spacing="-30">AB</text>'
                '</svg>'
            )
            with self.subTest(anchor=anchor):
                self._assert_checker_and_exporter_reject(
                    source,
                    'non-positive DrawingML text-frame extent',
                    'non-positive DrawingML text-frame extent',
                )

    def test_zero_emu_text_frame_extent_is_rejected(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360">
  <text x="40" y="100" font-size="19.9733333333333"
        letter-spacing="-23.4266666666667">AB</text>
</svg>''',
            'non-positive DrawingML text-frame extent (cx=0)',
            r'non-positive DrawingML text-frame extent \(cx=0',
        )

    def test_representable_negative_tracking_keeps_positive_text_frames(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360">
  <text x="40" y="100" font-size="20" letter-spacing="-20">AB</text>
  <text x="40" y="160" font-size="20" letter-spacing="-100">A</text>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'negative-tracking.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertIn(' spc="-1500"', slide_xml)
        self.assertIn(' spc="-7500"', slide_xml)
        extents = [
            int(value)
            for value in re.findall(r'<a:ext cx="(-?\d+)" cy="\d+"', slide_xml)
            if int(value) != 0
        ]
        self.assertEqual(len(extents), 2)
        self.assertTrue(all(value >= 1 for value in extents))

    def test_non_positive_inline_run_advance_blocks_checker_and_exporter(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="20">
  <text x="40" y="100">CONTROL PREFIX <tspan
        letter-spacing="-100">AB</tspan></text>
</svg>''',
            'non-positive DrawingML text-run advance',
            'non-positive DrawingML text-run advance',
        )

    def test_equivalent_inline_runs_cannot_bypass_tracking_validation(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="20">
  <text x="40" y="100">CONTROL PREFIX <tspan
        letter-spacing="-100">A</tspan><tspan
        letter-spacing="-100">B</tspan></text>
</svg>''',
            'non-positive DrawingML text-run advance',
            'non-positive DrawingML text-run advance',
        )

    def test_equivalent_inline_runs_are_emitted_as_one_drawingml_run(self):
        equivalent_sources = (
            'A<tspan>B</tspan>',
            '<tspan font-family="Helvetica">A</tspan>'
            '<tspan font-family="Arial">B</tspan>',
            '<tspan font-size="20">A</tspan>'
            '<tspan font-size="20.01">B</tspan>',
        )
        for inline_text in equivalent_sources:
            with self.subTest(inline_text=inline_text):
                source = (
                    '<svg xmlns="http://www.w3.org/2000/svg" '
                    'viewBox="0 0 640 360">'
                    '<text x="40" y="100" font-size="20" '
                    f'letter-spacing="2">{inline_text}</text></svg>'
                )
                with tempfile.TemporaryDirectory() as tmp_dir:
                    svg_path = Path(tmp_dir) / 'coalesced-runs.svg'
                    svg_path.write_text(source, encoding='utf-8')
                    slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
                self.assertEqual(slide_xml.count('<a:r>'), 1)
                self.assertIn('<a:t>AB</a:t>', slide_xml)

    def test_distinct_output_run_properties_remain_separate(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="20" letter-spacing="-100">
  <text x="40" y="100"><tspan fill="#FF0000">A</tspan><tspan
        fill="#0000FF">B</tspan><tspan fill="#0000FF">中</tspan></text>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'distinct-runs.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertEqual(slide_xml.count('<a:r>'), 3)
        self.assertIn('<a:rPr lang="en-US"', slide_xml)
        self.assertIn('<a:rPr lang="zh-CN"', slide_xml)

    def test_project_text_clusters_match_rendered_tracking_units(self):
        self.assertEqual(split_project_text_clusters('e\u0301'), ['e\u0301'])
        self.assertEqual(split_project_text_clusters('中\ufe00'), ['中\ufe00'])
        self.assertEqual(split_project_text_clusters('✊🏽'), ['✊🏽'])
        self.assertEqual(split_project_text_clusters('👨\u200d👩\u200d👧\u200d👦'), [
            '👨\u200d👩\u200d👧\u200d👦',
        ])
        self.assertEqual(split_project_text_clusters('🇨🇳🇺🇸🇫'), ['🇨🇳', '🇺🇸', '🇫'])
        self.assertEqual(split_project_text_clusters('क्ष'), ['क्ष'])
        self.assertEqual(split_project_text_clusters('A\u200dB'), ['A\u200d', 'B'])
        self.assertEqual(split_project_text_clusters('A्B'), ['A्', 'B'])
        self.assertEqual(split_project_text_clusters('क्A'), ['क्', 'A'])
        self.assertEqual(split_project_text_clusters('क्中'), ['क्', '中'])
        self.assertEqual(split_project_text_clusters('ᨠ᩠ᨡ'), ['ᨠ᩠ᨡ'])
        self.assertEqual(split_project_text_clusters('ᨠ᩠ꪀ'), ['ᨠ᩠', 'ꪀ'])
        self.assertEqual(
            estimate_text_width('e\u0301', 80),
            estimate_text_width('é', 80),
        )

    def test_single_graphemes_do_not_receive_internal_tracking(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="80" letter-spacing="-88">
  <text x="40" y="100">é</text>
  <text x="40" y="220">🇨🇳</text>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'single-graphemes.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertEqual(slide_xml.count(' spc="-6600"'), 2)

    def test_virama_does_not_join_text_across_unicode_scripts(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="80" letter-spacing="-100">
  <text x="40" y="100">A्B</text>
</svg>''',
            'non-positive DrawingML',
            'non-positive DrawingML',
        )

    def test_soft_break_joining_space_does_not_inherit_negative_tracking(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="20" letter-spacing="-21">
  <text x="40" y="80" data-paragraph-line-height="24">
    <tspan x="40">AB</tspan>
    <tspan x="40" dy="24" data-paragraph-soft-break="1">CD</tspan>
  </text>
</svg>'''
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'soft-break.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertEqual(slide_xml.count('<a:p>'), 1)
        self.assertEqual(slide_xml.count('<a:r>'), 3)
        self.assertEqual(slide_xml.count(' spc="-1575"'), 2)
        self.assertIn('<a:t xml:space="preserve"> </a:t>', slide_xml)

    def test_zero_text_run_advance_is_rejected_without_narrowing_safe_tracking(self):
        invalid_metrics = (
            ('20', '-22'),
            ('20', '-21.999999'),
            ('20.01', '-22'),
        )
        for font_size, letter_spacing in invalid_metrics:
            with self.subTest(
                font_size=font_size,
                letter_spacing=letter_spacing,
            ):
                self._assert_checker_and_exporter_reject(
                    '<svg xmlns="http://www.w3.org/2000/svg" '
                    'viewBox="0 0 640 360">'
                    f'<text x="40" y="100" font-size="{font_size}" '
                    f'letter-spacing="{letter_spacing}">AB</text></svg>',
                    "non-positive DrawingML text-run advance for 'AB' "
                    '(advance=0px)',
                    r"non-positive DrawingML text-run advance for 'AB' "
                    r"\(advance=0px\)",
                )

        safe = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360">
  <text x="40" y="100" font-size="20" letter-spacing="-21">AB</text>
</svg>'''
        safe_result = self._check(safe)
        self.assertTrue(safe_result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'positive-run-advance.svg'
            svg_path.write_text(safe, encoding='utf-8')
            safe_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertIn(' spc="-1575"', safe_xml)

        quantized_safe = safe.replace(
            'letter-spacing="-21"',
            'letter-spacing="-21.993"',
        )
        quantized_result = self._check(quantized_safe)
        self.assertTrue(quantized_result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'quantized-positive-run-advance.svg'
            svg_path.write_text(quantized_safe, encoding='utf-8')
            quantized_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertIn(' spc="-1649"', quantized_xml)

    def test_inherited_inline_tracking_cannot_hide_collapsed_text_frame(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="20">
  <g letter-spacing="-30">
    <text x="40" y="100">A<tspan style="letter-spacing:-100">BC</tspan></text>
  </g>
</svg>''',
            'non-positive DrawingML text-frame extent',
            'non-positive DrawingML text-frame extent',
        )

    def test_positioned_paragraph_tracking_is_still_blocked_by_exporter(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="20" letter-spacing="-30">
  <text x="40" y="80" data-paragraph-line-height="24">
    <tspan x="40">AB</tspan>
    <tspan x="40" dy="24">AB</tspan>
  </text>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'paragraph-negative-tracking.svg'
            svg_path.write_text(source, encoding='utf-8')
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'non-positive DrawingML text-frame extent',
            ):
                convert_svg_to_slide_shapes(svg_path)

    def test_bad_paragraph_run_cannot_hide_behind_a_wider_line(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="20">
  <text x="40" y="80" data-paragraph-line-height="24">
    <tspan x="40">HEALTHY PARAGRAPH LINE</tspan>
    <tspan x="40" dy="24" letter-spacing="-100">AB</tspan>
  </text>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'mixed-paragraph-tracking.svg'
            svg_path.write_text(source, encoding='utf-8')
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'non-positive DrawingML text-run advance',
            ):
                convert_svg_to_slide_shapes(svg_path)

    def test_bad_bullet_run_cannot_hide_behind_normal_text(self):
        self._assert_checker_and_exporter_reject(
            '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="20">
  <text x="40" y="100">• HEALTHY PREFIX <tspan
        letter-spacing="-100">AB</tspan></text>
</svg>''',
            'non-positive DrawingML text-run advance',
            'non-positive DrawingML text-run advance',
        )

    def test_bullet_text_uses_the_same_positive_frame_boundary(self):
        safe = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360">
  <text x="40" y="100" font-size="20" letter-spacing="-20">• AB</text>
</svg>'''
        safe_result = self._check(safe)
        self.assertTrue(safe_result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'safe-bullet.svg'
            svg_path.write_text(safe, encoding='utf-8')
            safe_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertIn('<a:buChar char="•"/>', safe_xml)
        self.assertIn(' spc="-1500"', safe_xml)

        self._assert_checker_and_exporter_reject(
            safe.replace('letter-spacing="-20"', 'letter-spacing="-100"'),
            'non-positive DrawingML text-frame extent',
            'non-positive DrawingML text-frame extent',
        )

    def test_whitespace_only_text_does_not_create_a_false_frame_error(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360">
  <text x="40" y="100" font-size="20" letter-spacing="-100"
        xml:space="preserve">  </text>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'whitespace-only.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertNotIn('<p:txBody>', slide_xml)

    def test_unchanged_imported_txbody_bypasses_generated_frame_check(self):
        root = ET.fromstring(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 640 360">
  <g id="imported-shape" data-pptx-object="shape" data-pptx-prst="rect">
    <rect data-pptx-part="geometry" x="20" y="40" width="160" height="80"
          fill="#FFFFFF"/>
    <text x="40" y="100" font-size="20" letter-spacing="-100">AB</text>
  </g>
</svg>'''
        )
        group = next(root.iter('{http://www.w3.org/2000/svg}g'))
        digest = svg_text_fingerprint(group)
        metadata = ET.SubElement(
            group,
            '{http://www.w3.org/2000/svg}metadata',
            {
                'data-pptx-part': 'txbody',
                'data-pptx-encoding': 'base64',
                'data-pptx-text-sha256': digest,
            },
        )
        metadata.text = self._native_txbody_payload()
        result = {'errors': [], 'warnings': []}

        SVGQualityChecker()._check_text_output_geometry(root, result)

        self.assertEqual(result['errors'], [])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'imported-txbody.svg'
            svg_path.write_text(
                ET.tostring(root, encoding='unicode'),
                encoding='utf-8',
            )
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertEqual(slide_xml.count('<p:txBody'), 1)
        self.assertIn(' spc="-7500"', slide_xml)
        self.assertNotIn('cx="-816864"', slide_xml)

    def test_invalid_unchanged_txbody_payload_does_not_bypass_checker(self):
        root = ET.fromstring(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 640 360">
  <g id="imported-shape" data-pptx-object="shape" data-pptx-prst="rect">
    <text x="40" y="100" font-size="20" letter-spacing="-100">AB</text>
  </g>
</svg>'''
        )
        group = next(root.iter('{http://www.w3.org/2000/svg}g'))
        digest = svg_text_fingerprint(group)
        metadata = ET.SubElement(
            group,
            '{http://www.w3.org/2000/svg}metadata',
            {
                'data-pptx-part': 'txbody',
                'data-pptx-encoding': 'base64',
                'data-pptx-text-sha256': digest,
            },
        )
        metadata.text = 'not-base64'
        result = {'errors': [], 'warnings': []}

        SVGQualityChecker()._check_preserved_txbody_contract(root, result)

        self.assertIn('cannot preserve source txBody', '\n'.join(result['errors']))

    def test_edited_imported_text_cannot_drop_run_effects(self):
        def imported_text_root(*, run_effect: bool) -> ET.Element:
            root = ET.fromstring(
                '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 640 360">
  <g id="imported-shape" data-pptx-object="shape" data-pptx-prst="rect">
    <rect data-pptx-part="geometry" x="20" y="40" width="160" height="80"
          fill="#FFFFFF"/>
    <text x="40" y="100" font-size="20">AB</text>
  </g>
</svg>'''
            )
            group = next(root.iter('{http://www.w3.org/2000/svg}g'))
            metadata = ET.SubElement(
                group,
                '{http://www.w3.org/2000/svg}metadata',
                {
                    'data-pptx-part': 'txbody',
                    'data-pptx-encoding': 'base64',
                    'data-pptx-text-sha256': svg_text_fingerprint(group),
                },
            )
            metadata.text = self._native_txbody_payload(
                run_effect=run_effect,
            )
            return root

        unchanged = imported_text_root(run_effect=True)
        with tempfile.TemporaryDirectory() as tmp_dir:
            unchanged_path = Path(tmp_dir) / 'unchanged-run-effect.svg'
            unchanged_path.write_text(
                ET.tostring(unchanged, encoding='unicode'),
                encoding='utf-8',
            )
            unchanged_xml = convert_svg_to_slide_shapes(unchanged_path)[0]
        self.assertEqual(unchanged_xml.count('<a:outerShdw'), 1)

        structurally_changed = imported_text_root(run_effect=True)
        group = next(
            structurally_changed.iter('{http://www.w3.org/2000/svg}g')
        )
        group.append(ET.fromstring(
            '<circle xmlns="http://www.w3.org/2000/svg" cx="220" cy="80" '
            'r="12" fill="#2563EB"/>'
        ))
        with tempfile.TemporaryDirectory() as tmp_dir:
            structural_path = Path(tmp_dir) / 'structural-run-effect.svg'
            structural_path.write_text(
                ET.tostring(structurally_changed, encoding='unicode'),
                encoding='utf-8',
            )
            structural_result = SVGQualityChecker().check_file(
                str(structural_path)
            )
            self.assertFalse(structural_result['passed'])
            self.assertIn(
                'cannot be restored as one native text shape',
                '\n'.join(structural_result['errors']),
            )
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'cannot be restored as one native text shape',
            ):
                convert_svg_to_slide_shapes(structural_path)

        edited = imported_text_root(run_effect=True)
        next(edited.iter('{http://www.w3.org/2000/svg}text')).text = 'AC'
        with tempfile.TemporaryDirectory() as tmp_dir:
            edited_path = Path(tmp_dir) / 'edited-run-effect.svg'
            edited_path.write_text(
                ET.tostring(edited, encoding='unicode'),
                encoding='utf-8',
            )
            checker_result = SVGQualityChecker().check_file(str(edited_path))
            self.assertFalse(checker_result['passed'])
            self.assertIn(
                'source txBody contains run-level effects',
                '\n'.join(checker_result['errors']),
            )
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'source txBody contains run-level effects',
            ):
                convert_svg_to_slide_shapes(edited_path)

        editable = imported_text_root(run_effect=False)
        next(editable.iter('{http://www.w3.org/2000/svg}text')).text = 'AC'
        with tempfile.TemporaryDirectory() as tmp_dir:
            editable_path = Path(tmp_dir) / 'edited-plain-text.svg'
            editable_path.write_text(
                ET.tostring(editable, encoding='unicode'),
                encoding='utf-8',
            )
            editable_result = SVGQualityChecker().check_file(str(editable_path))
            self.assertTrue(editable_result['passed'])
            editable_xml = convert_svg_to_slide_shapes(editable_path)[0]
        self.assertNotIn('<a:outerShdw', editable_xml)
        self.assertIn('<a:t>AC</a:t>', editable_xml)

    def test_txbody_run_effect_detection_covers_inherited_properties(self):
        dml = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        pml = 'http://schemas.openxmlformats.org/presentationml/2006/main'

        def txbody(properties: str) -> ET.Element:
            return ET.fromstring(
                f'<p:txBody xmlns:p="{pml}" xmlns:a="{dml}">'
                f'<a:bodyPr/><a:lstStyle/>{properties}</p:txBody>'
            )

        inherited_list = txbody(
            '<a:p><a:pPr><a:defRPr><a:effectLst>'
            '<a:glow rad="38100"><a:srgbClr val="000000"/></a:glow>'
            '</a:effectLst></a:defRPr></a:pPr></a:p>'
        )
        paragraph_end_dag = txbody(
            '<a:p><a:endParaRPr><a:effectDag><a:blur rad="38100"/>'
            '</a:effectDag></a:endParaRPr></a:p>'
        )
        empty_containers = txbody(
            '<a:p><a:r><a:rPr><a:effectLst/><a:effectDag/></a:rPr>'
            '<a:t>AB</a:t></a:r></a:p>'
        )
        local_text = txbody('<a:p><a:r><a:t>AB</a:t></a:r></a:p>')

        self.assertTrue(txbody_has_run_effects(inherited_list))
        self.assertTrue(txbody_has_run_effects(paragraph_end_dag))
        self.assertFalse(txbody_has_run_effects(empty_containers))
        self.assertTrue(txbody_has_run_effects(local_text, inherited_list))
        self.assertFalse(txbody_has_run_effects(local_text, empty_containers))

    def test_checker_does_not_trust_authored_runtime_txbody_snapshot(self):
        root = ET.fromstring(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 640 360">
  <g id="imported-shape" data-pptx-object="shape" data-pptx-prst="rect"
     data-pptx-runtime-txbody-unchanged="1">
    <rect data-pptx-part="geometry" x="20" y="40" width="160" height="80"
          fill="#FFFFFF"/>
    <text x="40" y="100" font-size="20" letter-spacing="-100">AB</text>
    <metadata data-pptx-part="txbody" data-pptx-encoding="base64"
              data-pptx-text-sha256="0000000000000000000000000000000000000000000000000000000000000000"/>
  </g>
</svg>'''
        )
        metadata = next(root.iter('{http://www.w3.org/2000/svg}metadata'))
        metadata.text = self._native_txbody_payload()
        result = {'errors': [], 'warnings': []}

        SVGQualityChecker()._check_text_output_geometry(root, result)

        self.assertIn(
            'non-positive DrawingML text-frame extent',
            '\n'.join(result['errors']),
        )
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'forged-runtime-snapshot.svg'
            svg_path.write_text(
                ET.tostring(root, encoding='unicode'),
                encoding='utf-8',
            )
            with self.assertRaisesRegex(
                SvgNativeConversionError,
                'non-positive DrawingML text-frame extent',
            ):
                convert_svg_to_slide_shapes(svg_path)

    def test_text_property_contract_preserves_registered_drawingml_semantics(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-style="italic">
  <g id="text-group" style="font-weight:600; text-anchor:end;
     letter-spacing:2; text-decoration:underline line-through">
    <text x="500" y="160" font-family="Arial" font-size="24">
      <tspan font-weight="500">Regular</tspan><tspan font-weight="600">Bold</tspan>
    </text>
  </g>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])
        self.assertEqual(result['warnings'], [])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'text-contract.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertEqual(slide_xml.count('<a:r>'), 2)
        self.assertEqual(slide_xml.count(' b="1"'), 1)
        self.assertEqual(slide_xml.count(' i="1"'), 2)
        self.assertEqual(slide_xml.count(' u="sng"'), 2)
        self.assertEqual(slide_xml.count(' strike="sngStrike"'), 2)
        self.assertEqual(slide_xml.count(' spc="150"'), 2)
        self.assertIn('algn="r"', slide_xml)

    def test_text_whitespace_contract_matches_rendered_run_boundaries(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-family="Arial" font-size="20">
  <text x="20" y="50">A <tspan fill="#FF0000"> B </tspan> C</text>
  <text x="20" y="100"><tspan xml:space="preserve">  D  </tspan></text>
  <text x="20" y="150" xml:space="preserve"><tspan
        xml:space="default">  E  </tspan></text>
  <text x="20" y="200">F  G　H</text>
  <text x="20" y="250" xml:space="preserve">H
	I</text>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'text-whitespace.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        texts = re.findall(r'<a:t(?: [^>]*)?>(.*?)</a:t>', slide_xml, re.S)
        self.assertEqual(
            texts,
            ['A ', 'B ', 'C', '  D  ', 'E', 'F  G　H', 'H  I'],
        )
        self.assertNotIn('<a:t>H\n', slide_xml)
        self.assertEqual(slide_xml.count(' xml:space="preserve"'), 4)

    def test_paragraph_runs_use_the_same_xml_space_inheritance(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-family="Arial" font-size="20">
  <text x="20" y="50" data-paragraph-line-height="24">
    <tspan x="20" xml:space="preserve">  First  </tspan>
    <tspan x="20" dy="24" xml:space="default">  Second  </tspan>
  </text>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'paragraph-whitespace.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertEqual(slide_xml.count('<a:p>'), 1)
        self.assertIn(
            '<a:t xml:space="preserve">  First  Second</a:t>',
            slide_xml,
        )

    def test_paragraph_merge_preserves_inline_sibling_space_and_inheritance(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360">
  <text x="40" y="80" font-family="Arial" font-size="24">
    <tspan x="40" dy="0" font-family="Arial Black"
           font-weight="bold">&#183; Stable:</tspan>
    <tspan fill="#FF0000">distributed ownership</tspan>
    <tspan x="40" dy="40" font-family="Arial Black"
           font-weight="bold">&#183; Future:</tspan>
    <tspan fill="#FF0000">long-term incentives</tspan>
  </text>
</svg>'''
        result = self._check(source)
        self.assertTrue(result['passed'])
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'paragraph-inline-siblings.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        texts = re.findall(r'<a:t(?: [^>]*)?>(.*?)</a:t>', slide_xml, re.S)
        self.assertEqual(
            ''.join(texts),
            'Stable: distributed ownership'
            'Future: long-term incentives',
        )
        self.assertEqual(slide_xml.count('<a:buChar char="•"/>'), 2)
        visible_runs = [
            run
            for run in re.findall(r'<a:r>(.*?)</a:r>', slide_xml, re.S)
            if re.search(r'<a:t(?: [^>]*)?>[^ <]', run)
        ]
        self.assertEqual(len(visible_runs), 4)
        self.assertIn('typeface="Arial Black"', visible_runs[0])
        self.assertIn(' b="1"', visible_runs[0])
        self.assertIn('typeface="Arial"', visible_runs[1])
        self.assertNotIn(' b="1"', visible_runs[1])
        self.assertIn('typeface="Arial Black"', visible_runs[2])
        self.assertIn(' b="1"', visible_runs[2])
        self.assertIn('typeface="Arial"', visible_runs[3])
        self.assertNotIn(' b="1"', visible_runs[3])

    def test_flattened_lines_keep_tail_and_xml_space_overrides(self):
        cases = {
            'tail': (
                '<text x="10" y="20" font-size="20">'
                '<tspan x="10" dy="20">A</tspan> B</text>',
                ['A B'],
            ),
            'child_preserve': (
                '<text x="10" y="20" font-size="20">'
                '<tspan x="10" dy="20" xml:space="preserve">'
                '  A  </tspan></text>',
                ['  A  '],
            ),
            'child_default_reset': (
                '<text x="10" y="20" font-size="20" '
                'xml:space="preserve"><tspan x="10" dy="20" '
                'xml:space="default">  A  </tspan></text>',
                ['A'],
            ),
            'preserved_leading_text': (
                '<text x="10" y="20" font-size="20" '
                'xml:space="preserve">  A  '
                '<tspan x="10" dy="20">B</tspan></text>',
                ['  A  ', 'B'],
            ),
        }
        for name, (body, expected) in cases.items():
            source = (
                '<svg xmlns="http://www.w3.org/2000/svg" '
                f'viewBox="0 0 640 360">{body}</svg>'
            )
            with self.subTest(name=name), tempfile.TemporaryDirectory() as tmp_dir:
                svg_path = Path(tmp_dir) / f'{name}.svg'
                svg_path.write_text(source, encoding='utf-8')
                slide_xml = convert_svg_to_slide_shapes(
                    svg_path,
                    merge_paragraphs=False,
                )[0]
                texts = re.findall(
                    r'<a:t(?: [^>]*)?>(.*?)</a:t>',
                    slide_xml,
                    re.S,
                )
                self.assertEqual(texts, expected)

    def test_explicit_default_xml_space_matches_the_omitted_default(self):
        source = (
            '<svg xmlns="http://www.w3.org/2000/svg" '
            'viewBox="0 0 640 360"><text x="40" y="100" '
            'font-size="24">  A  B  </text></svg>'
        )
        explicit = source.replace(
            'font-size="24"',
            'font-size="24" xml:space="default"',
        )
        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'default-space.svg'
            svg_path.write_text(source, encoding='utf-8')
            omitted_xml = convert_svg_to_slide_shapes(svg_path)[0]
            svg_path.write_text(explicit, encoding='utf-8')
            explicit_xml = convert_svg_to_slide_shapes(svg_path)[0]
        self.assertEqual(omitted_xml, explicit_xml)
        self.assertIn('<a:t>A B</a:t>', explicit_xml)

    def test_vertical_text_import_keeps_literal_space_advances(self):
        tx_body = ET.fromstring('''
<p:txBody xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:bodyPr vert="eaVert"/><a:lstStyle/><a:p><a:r>
    <a:rPr sz="1800"/><a:t xml:space="preserve">甲  乙</a:t>
  </a:r></a:p>
</p:txBody>''')
        result = convert_vertical_txbody(
            tx_body,
            Xfrm(x=0, y=0, w=100, h=400),
            None,
        )
        self.assertEqual(result.svg.count('<tspan'), 4)
        root = ET.fromstring(
            '<svg xmlns="http://www.w3.org/2000/svg">'
            f'{result.svg}</svg>'
        )
        text = next(root.iter('{http://www.w3.org/2000/svg}text'))
        self.assertEqual(''.join(text.itertext()), '甲  乙')

    def test_txbody_preview_leading_break_advances_once_at_its_size(self):
        cases = {'1800': 69.2, '3600': 98.0}
        for break_size, expected_baseline in cases.items():
            with self.subTest(break_size=break_size):
                svg = self._txbody_preview(f'''
  <a:p>
    <a:pPr><a:defRPr sz="1800"/></a:pPr>
    <a:br><a:rPr sz="{break_size}"/></a:br>
    <a:r><a:rPr sz="1800"/><a:t>B</a:t></a:r>
  </a:p>''')
                baselines = self._visible_text_baselines(svg)
                self.assertAlmostEqual(baselines['B'][0], expected_baseline)

    def test_txbody_preview_empty_break_line_uses_its_own_break_size(self):
        cases = {
            ('1800', '1800'): 98.0,
            ('3600', '1800'): 98.0,
            ('1800', '3600'): 126.8,
        }
        for (first_break, empty_break), expected_baseline in cases.items():
            with self.subTest(
                first_break=first_break,
                empty_break=empty_break,
            ):
                svg = self._txbody_preview(f'''
  <a:p>
    <a:r><a:rPr sz="1800"/><a:t>A</a:t></a:r>
    <a:br><a:rPr sz="{first_break}"/></a:br>
    <a:br><a:rPr sz="{empty_break}"/></a:br>
    <a:r><a:rPr sz="1800"/><a:t>B</a:t></a:r>
  </a:p>''')
                baselines = self._visible_text_baselines(svg)
                self.assertAlmostEqual(baselines['A'][0], 40.4)
                self.assertAlmostEqual(baselines['B'][0], expected_baseline)

    def test_txbody_preview_empty_paragraph_uses_end_para_size(self):
        cases = {'1800': 98.0, '3600': 126.8}
        for empty_size, expected_baseline in cases.items():
            with self.subTest(empty_size=empty_size):
                svg = self._txbody_preview(f'''
  <a:p><a:r><a:rPr sz="1800"/><a:t>A</a:t></a:r></a:p>
  <a:p>
    <a:pPr><a:defRPr sz="1800"/></a:pPr>
    <a:endParaRPr sz="{empty_size}"/>
  </a:p>
  <a:p><a:r><a:rPr sz="1800"/><a:t>B</a:t></a:r></a:p>''')
                baselines = self._visible_text_baselines(svg)
                self.assertAlmostEqual(baselines['A'][0], 40.4)
                self.assertAlmostEqual(baselines['B'][0], expected_baseline)

    def test_text_property_compatibility_aliases_warn_and_normalize(self):
        compatible = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360">
  <text x="40" y="100" font-family="Arial" font-size="24"
        font-weight="medium" letter-spacing="2px"
        text-decoration="line-through underline">Medium</text>
  <text x="40" y="180" font-family="Arial" font-size="24"
        font-weight="semibold">Semibold</text>
</svg>'''
        canonical = compatible.replace(
            'font-weight="medium"', 'font-weight="500"'
        ).replace(
            'letter-spacing="2px"', 'letter-spacing="2"'
        ).replace(
            'text-decoration="line-through underline"',
            'text-decoration="underline line-through"',
        ).replace('font-weight="semibold"', 'font-weight="600"')

        result = self._check(compatible)
        warning_text = '\n'.join(result['warnings'])
        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])
        self.assertIn("font-weight='medium'", warning_text)
        self.assertIn('font-weight="500"', warning_text)
        self.assertIn("font-weight='semibold'", warning_text)
        self.assertIn('font-weight="600"', warning_text)
        self.assertIn("letter-spacing='2px'", warning_text)
        self.assertIn('letter-spacing="2"', warning_text)
        self.assertIn(
            "text-decoration='line-through underline'",
            warning_text,
        )

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'text-compatible.svg'
            svg_path.write_text(compatible, encoding='utf-8')
            compatible_xml = convert_svg_to_slide_shapes(svg_path)[0]
            svg_path.write_text(canonical, encoding='utf-8')
            canonical_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertEqual(compatible_xml, canonical_xml)
        self.assertEqual(compatible_xml.count(' b="1"'), 1)
        self.assertIn(' spc="150"', compatible_xml)
        self.assertIn(' u="sng"', compatible_xml)
        self.assertIn(' strike="sngStrike"', compatible_xml)

    def test_invalid_text_properties_block_checker_and_exporter(self):
        cases = {
            'anchor': ('text-anchor="banana"', 'text-anchor'),
            'weight': ('font-weight="heavy"', 'font-weight'),
            'style': ('font-style="oblique"', 'font-style'),
            'decoration': ('text-decoration="notunderline"', 'text-decoration'),
            'spacing': ('letter-spacing="banana"', 'letter-spacing'),
            'spacing_keyword': ('letter-spacing="normal"', 'letter-spacing'),
            'word_spacing': ('word-spacing="2"', 'word-spacing'),
            'baseline': ('style="dominant-baseline:middle"', 'dominant-baseline'),
            'unknown_attribute': ('textLength="120"', 'textLength'),
            'inline_filter': ('style="filter:url(#shadow)"', 'filter'),
            'xml_space': ('xml:space="Preserve"', 'xml:space'),
        }
        for name, (declaration, expected) in cases.items():
            content = (
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                f'<text x="40" y="100" font-size="24" {declaration}>'
                'Broken</text></svg>'
            )
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    content,
                    expected,
                    'invalid project text property',
                )

    def test_xml_space_is_scoped_to_text_and_tspan(self):
        invalid_sources = {
            'root': (
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360" xml:space="preserve">'
                '<text x="40" y="100" font-size="24">Text</text></svg>'
            ),
            'group': (
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360"><g xml:space="preserve">'
                '<text x="40" y="100" font-size="24">Text</text>'
                '</g></svg>'
            ),
        }
        for name, source in invalid_sources.items():
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    source,
                    'xml:space only as a direct attribute on <text> or <tspan>',
                    'invalid project text property',
                )

    def test_compatible_letter_spacing_units_match_unitless_px(self):
        cases = (
            ('2px', '2', '150'),
            ('3pt', '4', '300'),
            ('1em', '24', '1800'),
        )
        for compatible, canonical, spacing in cases:
            source = (
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<text x="40" y="100" font-family="Arial" font-size="24" '
                f'letter-spacing="{compatible}">Tracking</text></svg>'
            )
            normalized = source.replace(
                f'letter-spacing="{compatible}"',
                f'letter-spacing="{canonical}"',
            )
            with self.subTest(value=compatible):
                result = self._check(source)
                self.assertTrue(result['passed'])
                self.assertIn(
                    f'letter-spacing="{canonical}"',
                    '\n'.join(result['warnings']),
                )
                with tempfile.TemporaryDirectory() as tmp_dir:
                    svg_path = Path(tmp_dir) / 'tracking.svg'
                    svg_path.write_text(source, encoding='utf-8')
                    compatible_xml = convert_svg_to_slide_shapes(svg_path)[0]
                    svg_path.write_text(normalized, encoding='utf-8')
                    canonical_xml = convert_svg_to_slide_shapes(svg_path)[0]
                self.assertEqual(compatible_xml, canonical_xml)
                self.assertIn(f' spc="{spacing}"', compatible_xml)

    def test_relative_font_sizes_share_one_inheritance_model(self):
        source = '''<svg xmlns="http://www.w3.org/2000/svg"
     viewBox="0 0 640 360" font-size="20">
  <g id="relative-text" font-size="1.5em" letter-spacing="1em">
    <text x="40" y="100" font-family="Arial" font-size="2em">Parent<tspan
          font-size=".5em">Child</tspan></text>
    <text x="40" y="180" font-family="Arial" font-size="2em">
      <tspan x="40" dy="60" font-size=".5em">Positioned child</tspan>
    </text>
  </g>
</svg>'''
        result = self._check(source)
        warning_text = '\n'.join(result['warnings'])
        self.assertTrue(result['passed'])
        self.assertEqual(result['errors'], [])
        self.assertIn('letter-spacing="30"', warning_text)
        self.assertNotIn('letter-spacing="60"', warning_text)

        with tempfile.TemporaryDirectory() as tmp_dir:
            svg_path = Path(tmp_dir) / 'relative-text.svg'
            svg_path.write_text(source, encoding='utf-8')
            slide_xml = convert_svg_to_slide_shapes(svg_path)[0]

        self.assertIn(' sz="4500"', slide_xml)
        self.assertIn(' sz="2250"', slide_xml)
        self.assertIn(' spc="2250"', slide_xml)
        self.assertNotIn(' spc="4500"', slide_xml)

    def test_invalid_inherited_and_overridden_text_values_cannot_hide(self):
        cases = {
            'root_inheritance': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360" font-weight="heavy">'
                '<text x="40" y="100" font-size="24">Broken</text></svg>'
            ), 'font-weight'),
            'overridden_attribute': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<text x="40" y="100" font-size="24" font-weight="heavy" '
                'style="font-weight:700">Broken</text></svg>'
            ), 'font-weight'),
            'tspan_anchor': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<text x="40" y="100" font-size="24">'
                '<tspan text-anchor="end">Broken</tspan></text></svg>'
            ), 'text-anchor'),
            'spacing_range': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<text x="40" y="100" font-size="24" '
                'letter-spacing="5333.35">Broken</text></svg>'
            ), 'letter-spacing'),
            'malformed_inheritance': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<g style="font-weight"><text x="40" y="100" '
                'font-size="24">Broken</text></g></svg>'
            ), 'malformed inline style'),
            'hidden_invalid_font_size': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<text x="40" y="100" '
                'style="font-size:banana; font-size:24">Broken</text></svg>'
            ), 'font-size'),
            'wrong_element': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<rect x="20" y="20" width="100" height="60" '
                'font-weight="bold"/></svg>'
            ), 'cannot carry text property'),
            'unsupported_group_property': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<g line-height="2"><text x="40" y="100" '
                'font-size="24">Broken</text></g></svg>'
            ), 'line-height'),
            'unsupported_group_style': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<g style="font-feature-settings:smcp"><text x="40" '
                'y="100" font-size="24">Broken</text></g></svg>'
            ), 'font-feature-settings'),
            'unregistered_prefixed_property': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<g font-optical-sizing="auto"><text x="40" y="100" '
                'font-size="24">Broken</text></g></svg>'
            ), 'font-optical-sizing'),
            'malformed_prefixed_property': ((
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'viewBox="0 0 640 360">'
                '<g style="font-optical-sizing"><text x="40" y="100" '
                'font-size="24">Broken</text></g></svg>'
            ), 'malformed inline style'),
        }
        for name, (content, expected) in cases.items():
            with self.subTest(name=name):
                self._assert_checker_and_exporter_reject(
                    content,
                    expected,
                    'invalid project text property',
                )


if __name__ == '__main__':
    unittest.main()
