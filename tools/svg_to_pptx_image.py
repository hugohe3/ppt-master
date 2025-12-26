#!/usr/bin/env python3
"""
PPT Master - SVG 转 PPTX 工具 (图片嵌入版)

将项目中的 SVG 文件批量转换为 PowerPoint 演示文稿。
每个 SVG 文件先转换为 PNG 图片，然后嵌入到幻灯片中。

用法:
    python svg_to_pptx_image.py <项目路径>
    python svg_to_pptx_image.py <项目路径> -o output.pptx
"""

import sys
import os
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image
import cairosvg

# EMU 转换常量
EMU_PER_INCH = 914400
EMU_PER_PIXEL = EMU_PER_INCH / 96

def find_svg_files(project_path):
    """查找项目中的 SVG 文件"""
    svg_dir = project_path / 'svg_output'
    if not svg_dir.exists():
        svg_dir = project_path
    
    svg_files = list(svg_dir.glob('*.svg'))
    return sorted(svg_files)

def svg_to_png(svg_path, png_path, dpi=300):
    """将 SVG 转换为 PNG 图片"""
    cairosvg.svg2png(url=str(svg_path), write_to=str(png_path), dpi=dpi)

def main():
    parser = argparse.ArgumentParser(description='SVG 转 PPTX 工具 (图片嵌入版)')
    parser.add_argument('project_path', type=str, help='项目目录路径')
    parser.add_argument('-o', '--output', type=str, default=None, help='输出文件路径')
    args = parser.parse_args()
    
    project_path = Path(args.project_path)
    if not project_path.exists():
        print(f"错误: 路径不存在: {project_path}")
        sys.exit(1)
    
    svg_files = find_svg_files(project_path)
    if not svg_files:
        print("错误: 未找到 SVG 文件")
        sys.exit(1)
    
    # 创建 PPTX 演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 遍历 SVG 文件，转换并添加到幻灯片
    for i, svg_path in enumerate(svg_files):
        print(f"处理 {svg_path.name}...")
        
        # 创建空白幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 将 SVG 转换为临时 PNG 文件
        png_path = svg_path.with_suffix('.png')
        try:
            svg_to_png(svg_path, png_path, dpi=300)
            
            # 获取图片尺寸
            with Image.open(png_path) as img:
                width, height = img.size
            
            # 计算图片在幻灯片中的位置（居中）
            slide_width_px = int(prs.slide_width / EMU_PER_PIXEL)
            slide_height_px = int(prs.slide_height / EMU_PER_PIXEL)
            
            # 保持图片比例，确保完全适应幻灯片
            ratio = min(slide_width_px / width, slide_height_px / height)
            new_width = int(width * ratio)
            new_height = int(height * ratio)
            
            left = Emu((slide_width_px - new_width) / 2 * EMU_PER_PIXEL)
            top = Emu((slide_height_px - new_height) / 2 * EMU_PER_PIXEL)
            width_emu = Emu(new_width * EMU_PER_PIXEL)
            height_emu = Emu(new_height * EMU_PER_PIXEL)
            
            # 添加图片到幻灯片
            slide.shapes.add_picture(str(png_path), left, top, width_emu, height_emu)
            
            # 删除临时 PNG 文件
            png_path.unlink()
            
        except Exception as e:
            print(f"  错误处理 {svg_path.name}: {e}")
            # 清理临时文件
            if png_path.exists():
                png_path.unlink()
    
    # 保存 PPTX 文件
    if args.output:
        output_path = Path(args.output)
    else:
        output_path = project_path / f"{project_path.name}_image.pptx"
    
    prs.save(str(output_path))
    print(f"\n转换完成！PPTX 文件已保存到: {output_path}")

if __name__ == '__main__':
    main()
