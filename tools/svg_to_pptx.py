#!/usr/bin/env python3
"""
PPT Master - SVG 转 PPTX 工具

将项目中的 SVG 文件批量转换为 PowerPoint 演示文稿。
每个 SVG 文件对应一张幻灯片，SVG 以原生矢量格式嵌入。

用法:
    python3 tools/svg_to_pptx.py <项目路径>
    python3 tools/svg_to_pptx.py <项目路径> -o output.pptx
    python3 tools/svg_to_pptx.py <项目路径> --use-final

示例:
    python3 tools/svg_to_pptx.py examples/ppt169_demo
    python3 tools/svg_to_pptx.py examples/ppt169_demo -o presentation.pptx
    python3 tools/svg_to_pptx.py examples/ppt169_demo --use-final

依赖:
    pip install python-pptx

注意:
    - SVG 以原生矢量格式嵌入，保持可编辑性
    - 需要 PowerPoint 2016+ 才能正确显示
"""

import sys
import os
import argparse
import re
import zipfile
import shutil
import tempfile
from pathlib import Path
from typing import Optional, Tuple, List
from xml.etree import ElementTree as ET

# 检查 python-pptx 是否已安装
try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("错误: 缺少 python-pptx 库")
    print("请运行: pip install python-pptx")
    sys.exit(1)

# 导入项目工具模块
sys.path.insert(0, str(Path(__file__).parent))
try:
    from project_utils import get_project_info, CANVAS_FORMATS
except ImportError:
    CANVAS_FORMATS = {
        'ppt169': {'name': 'PPT 16:9', 'dimensions': '1280x720', 'viewbox': '0 0 1280 720'},
        'ppt43': {'name': 'PPT 4:3', 'dimensions': '1024x768', 'viewbox': '0 0 1024 768'},
        'wechat': {'name': '微信公众号头图', 'dimensions': '900x383', 'viewbox': '0 0 900 383'},
        'xiaohongshu': {'name': '小红书', 'dimensions': '1242x1660', 'viewbox': '0 0 1242 1660'},
        'moments': {'name': '朋友圈/Instagram', 'dimensions': '1080x1080', 'viewbox': '0 0 1080 1080'},
        'story': {'name': 'Story/竖版', 'dimensions': '1080x1920', 'viewbox': '0 0 1080 1920'},
        'banner': {'name': '横版 Banner', 'dimensions': '1920x1080', 'viewbox': '0 0 1920 1080'},
        'a4': {'name': 'A4 打印', 'dimensions': '1240x1754', 'viewbox': '0 0 1240 1754'},
    }
    
    def get_project_info(path):
        return {'format': 'unknown', 'name': Path(path).name}


# EMU 转换常量
EMU_PER_INCH = 914400
EMU_PER_PIXEL = EMU_PER_INCH / 96

# XML 命名空间
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'asvg': 'http://schemas.microsoft.com/office/drawing/2016/SVG/main',
}

# 注册命名空间
for prefix, uri in NAMESPACES.items():
    ET.register_namespace(prefix, uri)


def get_slide_dimensions(canvas_format: str) -> Tuple[int, int]:
    """获取幻灯片尺寸（EMU 单位）"""
    if canvas_format not in CANVAS_FORMATS:
        canvas_format = 'ppt169'
    
    dimensions = CANVAS_FORMATS[canvas_format]['dimensions']
    match = re.match(r'(\d+)[×x](\d+)', dimensions)
    if match:
        width_px = int(match.group(1))
        height_px = int(match.group(2))
    else:
        width_px, height_px = 1280, 720
    
    return int(width_px * EMU_PER_PIXEL), int(height_px * EMU_PER_PIXEL)


def get_pixel_dimensions(canvas_format: str) -> Tuple[int, int]:
    """获取画布像素尺寸"""
    if canvas_format not in CANVAS_FORMATS:
        canvas_format = 'ppt169'
    
    dimensions = CANVAS_FORMATS[canvas_format]['dimensions']
    match = re.match(r'(\d+)[×x](\d+)', dimensions)
    if match:
        return int(match.group(1)), int(match.group(2))
    return 1280, 720


def detect_format_from_svg(svg_path: Path) -> Optional[str]:
    """从 SVG 文件的 viewBox 检测画布格式"""
    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            content = f.read(2000)
        
        match = re.search(r'viewBox="([^"]+)"', content)
        if match:
            viewbox = match.group(1)
            for fmt_key, fmt_info in CANVAS_FORMATS.items():
                if fmt_info['viewbox'] == viewbox:
                    return fmt_key
    except Exception:
        pass
    return None


def find_svg_files(project_path: Path, source: str = 'output') -> Tuple[List[Path], str]:
    """
    查找项目中的 SVG 文件
    
    Args:
        project_path: 项目目录路径
        source: SVG 来源目录
            - 'output': svg_output（默认，原始版本）
            - 'final': svg_final（带嵌入图标/图片）
            - 'flat': svg_output_flattext（扁平化文本）
            - 'final_flat': svg_final_flattext（最终版+扁平化）
            - 或任意子目录名称
    
    Returns:
        (SVG 文件列表, 实际使用的目录名)
    """
    # 预定义目录映射
    dir_map = {
        'output': 'svg_output',
        'final': 'svg_final',
        'flat': 'svg_output_flattext',
        'final_flat': 'svg_final_flattext',
    }
    
    # 获取目录名（支持预定义别名或直接指定目录名）
    dir_name = dir_map.get(source, source)
    svg_dir = project_path / dir_name
    
    if not svg_dir.exists():
        print(f"  警告: {dir_name} 目录不存在，尝试 svg_output")
        dir_name = 'svg_output'
        svg_dir = project_path / dir_name
    
    if not svg_dir.exists():
        # 直接在指定目录查找
        if project_path.is_dir():
            svg_dir = project_path
            dir_name = project_path.name
        else:
            return [], ''
    
    return sorted(svg_dir.glob('*.svg')), dir_name


def create_slide_xml_with_svg(slide_num: int, svg_rid: str, width_emu: int, height_emu: int) -> str:
    """创建包含 SVG 图片的幻灯片 XML"""
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:pic>
        <p:nvPicPr>
          <p:cNvPr id="2" name="SVG Image {slide_num}"/>
          <p:cNvPicPr>
            <a:picLocks noChangeAspect="1"/>
          </p:cNvPicPr>
          <p:nvPr/>
        </p:nvPicPr>
        <p:blipFill>
          <a:blip r:embed="{svg_rid}"/>
          <a:stretch>
            <a:fillRect/>
          </a:stretch>
        </p:blipFill>
        <p:spPr>
          <a:xfrm>
            <a:off x="0" y="0"/>
            <a:ext cx="{width_emu}" cy="{height_emu}"/>
          </a:xfrm>
          <a:prstGeom prst="rect">
            <a:avLst/>
          </a:prstGeom>
        </p:spPr>
      </p:pic>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr>
    <a:masterClrMapping/>
  </p:clrMapOvr>
</p:sld>'''


def create_slide_rels_xml(svg_rid: str, svg_filename: str) -> str:
    """创建幻灯片关系文件 XML"""
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="{svg_rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/{svg_filename}"/>
</Relationships>'''


def create_pptx_with_native_svg(
    svg_files: List[Path],
    output_path: Path,
    canvas_format: Optional[str] = None,
    verbose: bool = True
) -> bool:
    """创建包含原生 SVG 的 PPTX 文件"""
    if not svg_files:
        print("错误: 没有找到 SVG 文件")
        return False
    
    # 自动检测画布格式
    if canvas_format is None:
        canvas_format = detect_format_from_svg(svg_files[0])
        if canvas_format and verbose:
            format_name = CANVAS_FORMATS.get(canvas_format, {}).get('name', canvas_format)
            print(f"  检测到画布格式: {format_name}")
    
    if canvas_format is None:
        canvas_format = 'ppt169'
        if verbose:
            print(f"  使用默认格式: PPT 16:9")
    
    width_emu, height_emu = get_slide_dimensions(canvas_format)
    pixel_width, pixel_height = get_pixel_dimensions(canvas_format)
    
    if verbose:
        print(f"  幻灯片尺寸: {pixel_width} x {pixel_height} px")
        print(f"  SVG 文件数: {len(svg_files)}")
        print()
    
    # 创建临时目录
    temp_dir = Path(tempfile.mkdtemp())
    
    try:
        # 首先用 python-pptx 创建基础 PPTX
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu
        
        # 添加空白幻灯片作为占位
        blank_layout = prs.slide_layouts[6]
        for _ in svg_files:
            prs.slides.add_slide(blank_layout)
        
        # 保存基础 PPTX
        base_pptx = temp_dir / 'base.pptx'
        prs.save(str(base_pptx))
        
        # 解压 PPTX
        extract_dir = temp_dir / 'pptx_content'
        with zipfile.ZipFile(base_pptx, 'r') as zf:
            zf.extractall(extract_dir)
        
        # 创建 media 目录
        media_dir = extract_dir / 'ppt' / 'media'
        media_dir.mkdir(exist_ok=True)
        
        # 处理每个 SVG 文件
        success_count = 0
        for i, svg_path in enumerate(svg_files, 1):
            slide_num = i
            svg_filename = f'image{i}.svg'
            svg_rid = 'rId2'
            
            try:
                # 复制 SVG 到 media 目录
                shutil.copy(svg_path, media_dir / svg_filename)
                
                # 更新幻灯片 XML
                slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                slide_xml = create_slide_xml_with_svg(slide_num, svg_rid, width_emu, height_emu)
                with open(slide_xml_path, 'w', encoding='utf-8') as f:
                    f.write(slide_xml)
                
                # 创建/更新关系文件
                rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                rels_dir.mkdir(exist_ok=True)
                rels_path = rels_dir / f'slide{slide_num}.xml.rels'
                rels_xml = create_slide_rels_xml(svg_rid, svg_filename)
                with open(rels_path, 'w', encoding='utf-8') as f:
                    f.write(rels_xml)
                
                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name}")
                
                success_count += 1
                
            except Exception as e:
                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name} - 错误: {e}")
        
        # 更新 [Content_Types].xml 添加 SVG 类型
        content_types_path = extract_dir / '[Content_Types].xml'
        with open(content_types_path, 'r', encoding='utf-8') as f:
            content_types = f.read()
        
        # 添加 SVG 扩展类型（如果不存在）
        if 'Extension="svg"' not in content_types:
            content_types = content_types.replace(
                '</Types>',
                '  <Default Extension="svg" ContentType="image/svg+xml"/>\n</Types>'
            )
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)
        
        # 重新打包 PPTX
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(extract_dir)
                    zf.write(file_path, arcname)
        
        if verbose:
            print()
            print(f"[完成] 已保存: {output_path}")
            print(f"  成功: {success_count}, 失败: {len(svg_files) - success_count}")
        
        return success_count == len(svg_files)
        
    finally:
        # 清理临时目录
        shutil.rmtree(temp_dir, ignore_errors=True)


def main():
    parser = argparse.ArgumentParser(
        description='PPT Master - SVG 转 PPTX 工具（原生 SVG 嵌入）',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
示例:
    %(prog)s examples/ppt169_demo
    %(prog)s examples/ppt169_demo -s final
    %(prog)s examples/ppt169_demo -s my_custom_folder
    %(prog)s examples/ppt169_demo -o presentation.pptx

SVG 来源目录 (-s):
    output     - svg_output（默认，原始版本）
    final      - svg_final（带嵌入图标/图片）
    flat       - svg_output_flattext（扁平化文本）
    final_flat - svg_final_flattext（最终版+扁平化）
    <任意名>   - 直接指定项目下的子目录名

特点:
    - SVG 以原生矢量格式嵌入，保持可编辑性
    - 需要 PowerPoint 2016+ 查看
'''
    )
    
    parser.add_argument('project_path', type=str, help='项目目录路径')
    parser.add_argument('-o', '--output', type=str, default=None, help='输出文件路径')
    parser.add_argument('-s', '--source', type=str, default='output', 
                        help='SVG 来源: output/final/flat/final_flat 或任意子目录名')
    parser.add_argument('-f', '--format', type=str, choices=list(CANVAS_FORMATS.keys()), default=None, help='指定画布格式')
    parser.add_argument('-q', '--quiet', action='store_true', help='静默模式')
    
    args = parser.parse_args()
    
    project_path = Path(args.project_path)
    if not project_path.exists():
        print(f"错误: 路径不存在: {project_path}")
        sys.exit(1)
    
    try:
        project_info = get_project_info(str(project_path))
        project_name = project_info.get('name', project_path.name)
        detected_format = project_info.get('format')
    except Exception:
        project_name = project_path.name
        detected_format = None
    
    canvas_format = args.format
    if canvas_format is None and detected_format and detected_format != 'unknown':
        canvas_format = detected_format
    
    svg_files, source_dir_name = find_svg_files(project_path, args.source)
    
    if not svg_files:
        print("错误: 未找到 SVG 文件")
        sys.exit(1)
    
    if args.output:
        output_path = Path(args.output)
    else:
        output_path = project_path / f"{project_name}.pptx"
    
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    verbose = not args.quiet
    if verbose:
        print("PPT Master - SVG 转 PPTX 工具（原生 SVG）")
        print("=" * 50)
        print(f"  项目路径: {project_path}")
        print(f"  SVG 目录: {source_dir_name}")
        print(f"  输出文件: {output_path}")
        print()
    
    success = create_pptx_with_native_svg(
        svg_files,
        output_path,
        canvas_format=canvas_format,
        verbose=verbose
    )
    
    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
